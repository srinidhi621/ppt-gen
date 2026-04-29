"""Extract shape data from designer reference slides.

Reads assets/ground_truth/internal_inbox/designer_reference_slides.pptx
and produces JSON files in examples/source/ for the targeted slides.

Run:
    PYTHONPATH=. .venv/bin/python examples/extract_designer_slides.py
"""

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = PROJECT_ROOT / "assets" / "ground_truth" / "internal_inbox" / "designer_reference_slides.pptx"
OUTPUT_DIR = PROJECT_ROOT / "examples" / "source"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# Slide definitions: (index, file_suffix, name, proposed_archetype)
# Targets reflect the audit at assets/ground_truth/annotations/designer_reference_slides_audit.md.
TARGETS = [
    (0,  "S01", "S01 - Hero title with background visual", "hero_title"),
    (1,  "S02", "S02 - Numbered infographic overlay",      "process_flow"),
    (2,  "S03", "S03 - Image-only section divider",        "section_break"),
    (3,  "S04", "S04 - Image-only section divider (alt)",  "section_break"),
    (5,  "S06", "S06 - Solution flow with connectors",     "content_with_visual"),
    (6,  "S07", "S07 - Five concept cards (prompt eng.)",  "three_cards"),
    (13, "S14", "S14 - Four-quadrant proposition matrix",  "matrix_grid"),
    (14, "S15", "S15 - KPI band + value-prop grid",        "kpi_grid"),
    (17, "S18", "S18 - Priority-stratified scope cards",   "three_cards"),
]


def extract_shape(shape):
    """Extract relevant data from a single shape."""
    data = {
        "name": shape.name,
        "shape_type": str(shape.shape_type) if shape.shape_type else "UNKNOWN",
        "left_emu": shape.left,
        "top_emu": shape.top,
        "width_emu": shape.width,
        "height_emu": shape.height,
        "fill_hex": None,
        "has_text": False,
        "text": None,
        "font_name": None,
        "font_size_pt": None,
        "font_bold": None,
        "font_color_hex": None,
        "is_image": False,
    }

    # Check for image
    if shape.shape_type is not None and shape.shape_type == 13:  # PICTURE
        data["is_image"] = True

    # Check for fill
    try:
        fill = shape.fill
        if fill.type is not None:
            try:
                rgb = fill.fore_color.rgb
                data["fill_hex"] = f"#{rgb}"
            except (AttributeError, TypeError):
                pass
    except (AttributeError, TypeError):
        pass

    # Check for text
    try:
        if shape.has_text_frame:
            tf = shape.text_frame
            full_text = tf.text.strip()
            if full_text:
                data["has_text"] = True
                data["text"] = full_text

            # Get font from first run of first paragraph
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        data["font_name"] = run.font.name
                    if run.font.size:
                        data["font_size_pt"] = run.font.size.pt
                    if run.font.bold is not None:
                        data["font_bold"] = run.font.bold
                    try:
                        if run.font.color and run.font.color.rgb:
                            data["font_color_hex"] = f"#{run.font.color.rgb}"
                    except (AttributeError, TypeError):
                        pass
                    # Only need first run's font info
                    break
                if data["font_name"]:
                    break
    except (AttributeError, TypeError):
        pass

    return data


def extract_slide(prs, slide_index, slide_name, proposed_archetype):
    """Extract all shape data from a slide."""
    slide = prs.slides[slide_index]
    shapes = []
    for shape in slide.shapes:
        shapes.append(extract_shape(shape))

    return {
        "slide_index": slide_index,
        "slide_name": slide_name,
        "proposed_archetype": proposed_archetype,
        "shape_count": len(shapes),
        "shapes": shapes,
    }


def main():
    prs = Presentation(str(PPTX_PATH))
    total_slides = len(prs.slides)
    print(f"Loaded {PPTX_PATH.name} with {total_slides} slides")

    for slide_index, suffix, name, archetype in TARGETS:
        if slide_index >= total_slides:
            print(f"SKIP  slide {slide_index} ({suffix}) - only {total_slides} slides in file")
            continue

        data = extract_slide(prs, slide_index, name, archetype)
        out_path = OUTPUT_DIR / f"designer_{suffix}.json"
        with open(out_path, "w") as f:
            json.dump(data, f, indent=2)
        print(f"OK    {suffix} -> {out_path} ({data['shape_count']} shapes)")


if __name__ == "__main__":
    main()
