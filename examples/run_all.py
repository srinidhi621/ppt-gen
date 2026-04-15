"""Example regression suite runner.

Discovers all example build.py files under examples/ and runs each one.
After each build, runs the objective scanner and checks density bounds
from metadata.json. Reports success/failure for each example.

Run:
    PYTHONPATH=. .venv/bin/python examples/run_all.py
"""

import json
import subprocess
import sys
from pathlib import Path

EXAMPLES_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = EXAMPLES_DIR.parent

# Ensure the project root is on sys.path for scanner/runtime imports
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))
DESIGN_SYSTEM = PROJECT_ROOT / "assets" / "template" / "design_system.json"


def discover_examples():
    """Find all build.py files under examples/."""
    build_files = sorted(EXAMPLES_DIR.glob("*/*/build.py"))
    # Exclude any build.py in __pycache__
    return [
        f for f in build_files
        if "__pycache__" not in str(f.relative_to(EXAMPLES_DIR))
    ]


def run_example(build_path, python_exe):
    """Run a single example build.py and return (success, output)."""
    try:
        result = subprocess.run(
            [python_exe, str(build_path)],
            cwd=str(PROJECT_ROOT),
            env={**__import__("os").environ, "PYTHONPATH": str(PROJECT_ROOT)},
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode == 0:
            return True, result.stdout.strip()
        else:
            return False, result.stderr.strip() or result.stdout.strip()
    except subprocess.TimeoutExpired:
        return False, "TIMEOUT (60s)"
    except Exception as e:
        return False, str(e)


def check_scanner(output_pptx):
    """Run the objective scanner on the output PPTX. Returns (pass, findings)."""
    from src.scan.scanner import scan_pptx

    report = scan_pptx(str(output_pptx), str(DESIGN_SYSTEM))
    return report["pass"], report["findings"]


def check_density(output_pptx, metadata):
    """Check that the built slide's shape count is within density bounds."""
    from pptx import Presentation

    prs = Presentation(str(output_pptx))
    last_slide = prs.slides[-1]
    shape_count = len(last_slide.shapes)

    density = metadata.get("density", {})
    min_shapes = density.get("min_shapes", 0)
    max_shapes = density.get("max_shapes", 999)

    errors = []
    if shape_count < min_shapes:
        errors.append(
            f"Shape count {shape_count} < min_shapes {min_shapes}"
        )
    if shape_count > max_shapes:
        errors.append(
            f"Shape count {shape_count} > max_shapes {max_shapes}"
        )
    return errors


def main():
    python_exe = sys.executable
    examples = discover_examples()

    if not examples:
        print("No examples found!")
        sys.exit(1)

    print(f"Found {len(examples)} examples")
    print(f"Python: {python_exe}")
    print(f"Project root: {PROJECT_ROOT}")
    print("-" * 60)

    passed = 0
    failed = 0

    for build_path in examples:
        rel_path = build_path.relative_to(EXAMPLES_DIR)
        example_dir = build_path.parent
        output_pptx = example_dir / "output.pptx"
        meta_path = example_dir / "metadata.json"

        # Step 1: Build
        success, output = run_example(build_path, python_exe)
        if not success:
            failed += 1
            print(f"  FAIL  {rel_path}  [build failed]")
            for line in output.split("\n")[:3]:
                print(f"        {line}")
            continue

        if not output_pptx.exists():
            failed += 1
            print(f"  FAIL  {rel_path}  [no output.pptx]")
            continue

        # Step 2: Scanner
        scan_ok, findings = check_scanner(output_pptx)
        if not scan_ok:
            blocking = [f for f in findings if f["severity"] == "BLOCKING"]
            failed += 1
            print(f"  FAIL  {rel_path}  [scanner: {len(blocking)} BLOCKING]")
            for f in blocking[:3]:
                print(f"        {f['check_id']}: {f['details'][:80]}")
            continue

        # Step 3: Density bounds
        if meta_path.exists():
            meta = json.loads(meta_path.read_text())
            density_errors = check_density(output_pptx, meta)
            if density_errors:
                failed += 1
                print(f"  FAIL  {rel_path}  [density]")
                for e in density_errors:
                    print(f"        {e}")
                continue

        passed += 1
        print(f"  PASS  {rel_path}")

    print("-" * 60)
    print(f"Total: {len(examples)}  Passed: {passed}  Failed: {failed}")

    if failed > 0:
        sys.exit(1)


if __name__ == "__main__":
    main()
