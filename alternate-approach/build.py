"""Build the 10x Accelerated Talent Program deck on the Ascendion template."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x1E, 0x27, 0x61)
CORAL = RGBColor(0xF9, 0x61, 0x67)
CHARCOAL = RGBColor(0x36, 0x45, 0x4F)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MED_GRAY = RGBColor(0x8A, 0x94, 0xA6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x21, 0x21, 0x21)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
BODY_FONT = "Calibri"

prs = Presentation("template.pptx")
SW, SH = prs.slide_width, prs.slide_height

layout_blank = None
for l in prs.slide_layouts:
    if l.name == "Blank":
        layout_blank = l

# Edit slide 1 title
slide1 = prs.slides[0]
for shp in slide1.shapes:
    if shp.has_text_frame and "Slide heading" in shp.text_frame.text:
        paras = shp.text_frame.paragraphs
        if paras[0].runs:
            paras[0].runs[0].text = "10x Accelerated Talent Program"
        if len(paras) > 1 and paras[1].runs:
            paras[1].runs[0].text = "Operating Plan  |  Leadership Review"

# Delete slide 2 (placeholder)
xml_slides = prs.slides._sldIdLst
slides_list = list(xml_slides)
rId = slides_list[1].get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
prs.part.drop_rel(rId)
xml_slides.remove(slides_list[1])

def add_blank():
    return prs.slides.add_slide(layout_blank)

def add_rect(slide, left, top, width, height, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp

def add_text(slide, left, top, width, height, text, *, size=14, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = BODY_FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb

def slide_header(slide, kicker, title):
    add_rect(slide, Emu(0), Emu(0), SW, Inches(0.08), NAVY)
    add_rect(slide, Inches(0.5), Inches(0.45), Inches(0.18), Inches(0.18), CORAL)
    add_text(slide, Inches(0.78), Inches(0.4), Inches(10), Inches(0.3),
             kicker.upper(), size=10, bold=True, color=CORAL)
    add_text(slide, Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.6),
             title, size=26, bold=True, color=NAVY)


# ==== SLIDE 2: Operating Principle ====
s = add_blank()
slide_header(s, "01  |  The thesis", "What this program is for")

add_rect(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(2.2), NAVY)
add_rect(s, Inches(0.5), Inches(1.6), Inches(0.12), Inches(2.2), CORAL)
add_text(s, Inches(0.95), Inches(1.85), Inches(11.5), Inches(0.4),
         "OPERATING PRINCIPLE", size=11, bold=True, color=CORAL)
add_text(s, Inches(0.95), Inches(2.25), Inches(11.5), Inches(1.5),
         "Find engineers who use judgment and AI leverage to remove work \u2014\nthen put them on a track that doesn\u2019t destroy that capability.",
         size=22, bold=True, color=WHITE)

col_w = Inches(4.0)
col_h = Inches(2.5)
col_top = Inches(4.15)
col_lefts = [Inches(0.5), Inches(4.65), Inches(8.8)]
cols = [
    ("THE PROBLEM", "Conventional hiring rewards pedigree, narrow correctness, and confidence. None of these find force multipliers."),
    ("WHAT WE WON\u2019T DO", "No HackerRank rounds. No essay prompts. No personality interviews. No steering committees. No \u201810x\u2019 marketing language externally."),
    ("WHAT WE WILL DO", "Verifiable artifacts. Scaffolded take-homes with AI allowed and scored. Anchored defense panels. Protected runway after admit."),
]
for left, (h, body) in zip(col_lefts, cols):
    add_rect(s, left, col_top, col_w, col_h, LIGHT_GRAY)
    add_rect(s, left, col_top, Inches(0.5), Inches(0.08), CORAL)
    add_text(s, left + Inches(0.25), col_top + Inches(0.25), col_w - Inches(0.5), Inches(0.35),
             h, size=11, bold=True, color=NAVY)
    add_text(s, left + Inches(0.25), col_top + Inches(0.65), col_w - Inches(0.5), col_h - Inches(0.8),
             body, size=12, color=CHARCOAL)

# ==== SLIDE 3: Two tracks ====
s = add_blank()
slide_header(s, "02  |  Program shape", "Two tracks. Small cohorts. On purpose.")
add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4),
         "A 10x program with 200 admits is not a 10x program. We optimize for precision, not recall.",
         size=13, color=CHARCOAL)

card_w = Inches(6.0)
card_h = Inches(4.4)
card_top = Inches(2.0)

def track_card(left, name, profile, size_, cycle, focus_label, focus_items):
    add_rect(s, left, card_top, card_w, card_h, WHITE, line=MED_GRAY)
    add_rect(s, left, card_top, card_w, Inches(0.85), NAVY)
    add_text(s, left + Inches(0.35), card_top + Inches(0.15), card_w - Inches(0.7), Inches(0.4),
             name, size=20, bold=True, color=WHITE)
    add_text(s, left + Inches(0.35), card_top + Inches(0.5), card_w - Inches(0.7), Inches(0.3),
             profile, size=11, color=ICE)
    rows = [
        ("Cohort size", size_),
        ("Cadence", cycle),
        (focus_label, focus_items),
    ]
    y = card_top + Inches(1.15)
    for label, val in rows:
        add_text(s, left + Inches(0.35), y, Inches(1.8), Inches(0.3),
                 label.upper(), size=9, bold=True, color=CORAL)
        add_text(s, left + Inches(2.15), y - Inches(0.05), card_w - Inches(2.5), Inches(1.3),
                 val, size=12, color=DARK)
        y += Inches(1.05)

track_card(Inches(0.5), "ASCENDER",
           "0\u20132 yrs experience  |  external + internal",
           "8\u201312 per cohort", "Twice yearly",
           "Test for",
           "Decomposition under ambiguity. Taste in deciding what NOT to fix. Healthy AI use the candidate can defend line by line.")

track_card(Inches(6.85), "PRINCIPAL TRACK",
           "2\u20136 yrs experience  |  lateral + skip-level nomination",
           "8\u201312 per cohort", "Twice yearly",
           "Test for",
           "Sequencing and trade-off articulation. Resisting the urge to rewrite. Systems thinking, governance awareness, ability to teach.")

add_text(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.35),
         "Internal nominations run from day one and enter at Stage 2. Highest-yield source. The current draft ignores it.",
         size=11, bold=True, color=NAVY)

# ==== SLIDE 4: The Funnel ====
s = add_blank()
slide_header(s, "03  |  The funnel", "Four stages. Five weeks. ~150 in, 8\u201312 out.")

stages = [
    ("01", "EVIDENCE GATE", "Week 1",
     "GitHub link + 3-min Loom + one shipped artifact. No essays.",
     "~150 \u2192 ~40"),
    ("02", "SCAFFOLDED TAKE-HOME", "Week 2\u20133",
     "48 hrs. Real repo. AI allowed and expected. PR + 10-min walkthrough.",
     "~40 \u2192 ~16"),
    ("03", "DEFENSE PANEL", "Week 4",
     "60 min. Two panelists. Five anchored questions. Independent scoring.",
     "~16 \u2192 ~10"),
    ("04", "ADMIT DECISION", "Week 5",
     "30 min. Director + Principal. Confirms appetite. No new tech eval.",
     "8\u201312 admits"),
]

card_w = Inches(2.95)
card_h = Inches(4.7)
gap = Inches(0.18)
total = card_w * 4 + gap * 3
start_left = (SW - total) // 2
top = Inches(1.55)

for i, (num, name, when, what, ratio) in enumerate(stages):
    left = start_left + (card_w + gap) * i
    add_rect(s, left, top, card_w, card_h, WHITE, line=MED_GRAY)
    add_rect(s, left, top, card_w, Inches(1.05), NAVY)
    add_text(s, left + Inches(0.3), top + Inches(0.15), card_w - Inches(0.6), Inches(0.55),
             num, size=28, bold=True, color=CORAL)
    add_text(s, left + Inches(0.3), top + Inches(0.68), card_w - Inches(0.6), Inches(0.3),
             when.upper(), size=10, bold=True, color=WHITE)
    add_text(s, left + Inches(0.3), top + Inches(1.25), card_w - Inches(0.6), Inches(0.4),
             name, size=13, bold=True, color=NAVY)
    add_text(s, left + Inches(0.3), top + Inches(1.75), card_w - Inches(0.6), Inches(2.2),
             what, size=12, color=CHARCOAL)
    add_rect(s, left, top + card_h - Inches(0.85), card_w, Inches(0.85), LIGHT_GRAY)
    add_text(s, left + Inches(0.3), top + card_h - Inches(0.7), Inches(2), Inches(0.25),
             "FUNNEL", size=8, bold=True, color=MED_GRAY)
    add_text(s, left + Inches(0.3), top + card_h - Inches(0.45), card_w - Inches(0.6), Inches(0.4),
             ratio, size=14, bold=True, color=NAVY)

# ==== SLIDE 5: Discriminators ====
s = add_blank()
slide_header(s, "04  |  Where this differs", "The two design decisions that determine whether it works")

left_x = Inches(0.5)
left_w = Inches(6.7)
top_y = Inches(1.55)

add_text(s, left_x, top_y, left_w, Inches(0.4),
         "AI-USE RUBRIC", size=12, bold=True, color=CORAL)
add_text(s, left_x, top_y + Inches(0.32), left_w, Inches(0.4),
         "Scored explicitly. Without this, calibration collapses.",
         size=11, color=CHARCOAL)

rub_top = top_y + Inches(0.85)
rub_h = Inches(4.4)
col_w_r = (left_w - Inches(0.2)) / 3

rubric_cols = [
    ("HEALTHY USE", NAVY,
     "Used AI for boilerplate, scaffolding, tests. Verified outputs. Can explain every line."),
    ("OVER-RELIANCE", CORAL,
     "Cannot explain choices. Inconsistent style. Hallucinated comments. Tests don\u2019t match code."),
    ("AVOIDANCE", MED_GRAY,
     "Avoided AI in an AI-appropriate task. Slower with no quality gain."),
]
for i, (h, c, body) in enumerate(rubric_cols):
    cl = left_x + (col_w_r + Inches(0.1)) * i
    add_rect(s, cl, rub_top, col_w_r, rub_h, LIGHT_GRAY)
    add_rect(s, cl, rub_top, col_w_r, Inches(0.45), c)
    add_text(s, cl + Inches(0.15), rub_top + Inches(0.1), col_w_r - Inches(0.3), Inches(0.3),
             h, size=10, bold=True, color=WHITE)
    add_text(s, cl + Inches(0.15), rub_top + Inches(0.65), col_w_r - Inches(0.3), rub_h - Inches(0.8),
             body, size=11, color=CHARCOAL)

right_x = Inches(7.45)
right_w = Inches(5.35)
add_text(s, right_x, top_y, right_w, Inches(0.4),
         "BASE RATES PER TRACK PER CYCLE", size=12, bold=True, color=CORAL)
add_text(s, right_x, top_y + Inches(0.32), right_w, Inches(0.4),
         "If they want recall, this is the wrong program.",
         size=11, color=CHARCOAL)

stats_top = top_y + Inches(0.85)
stats = [
    ("~150", "applicants"),
    ("~40", "Stage 2 invitees"),
    ("~16", "Stage 3 invitees"),
    ("8\u201312", "admits"),
    ("~6 hrs", "reviewer cost / admit"),
    ("25%", "false-positive tolerance"),
]
sw_ = (right_w - Inches(0.15)) / 2
sh_ = Inches(1.35)
for idx, (big, label) in enumerate(stats):
    col = idx % 2
    row = idx // 2
    sl = right_x + (sw_ + Inches(0.15)) * col
    st = stats_top + (sh_ + Inches(0.1)) * row
    add_rect(s, sl, st, sw_, sh_, WHITE, line=MED_GRAY)
    add_rect(s, sl, st, Inches(0.08), sh_, CORAL)
    add_text(s, sl + Inches(0.25), st + Inches(0.18), sw_ - Inches(0.4), Inches(0.7),
             big, size=24, bold=True, color=NAVY)
    add_text(s, sl + Inches(0.25), st + Inches(0.78), sw_ - Inches(0.4), Inches(0.5),
             label, size=10, color=CHARCOAL)

# ==== SLIDE 6: Development Arc ====
s = add_blank()
slide_header(s, "05  |  After admit", "Selection is 20%. The development arc is 80%.")

add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.45),
         "If we cannot protect the runway, we should not run the program. This is the failure mode that quietly kills these initiatives.",
         size=13, color=CHARCOAL)

phases = [
    ("MONTHS 1\u20133", "PROTECTED RUNWAY",
     "60% protected time. A real problem. A named sponsor. NOT on standard delivery utilization.",
     "Requires written executive air cover. A Delivery VP will try to break this in week 2."),
    ("MONTHS 3\u20139", "SHADOW + LEAD",
     "Pair with a Principal Engineer on a real client engagement. Lead one workstream end-to-end.",
     "Quarterly review on five anchored dimensions: framing, leverage, taste, ownership, teaching."),
    ("MONTHS 9\u201318", "FORCE MULTIPLIER",
     "Own a small team or a reusable practice asset \u2014 playbook, internal tool, eval harness.",
     "Promotion or off-ramp at month 18. Off-ramp to senior IC track is graduation, not failure."),
]

ph_w = Inches(4.05)
ph_h = Inches(4.5)
gap = Inches(0.2)
ph_top = Inches(2.0)
total = ph_w * 3 + gap * 2
ph_start = (SW - total) // 2

for i, (when, name, what, risk) in enumerate(phases):
    left = ph_start + (ph_w + gap) * i
    add_rect(s, left, ph_top, ph_w, ph_h, WHITE, line=MED_GRAY)
    add_rect(s, left, ph_top, ph_w, Inches(0.12), CORAL)
    add_text(s, left + Inches(0.3), ph_top + Inches(0.3), ph_w - Inches(0.6), Inches(0.3),
             when, size=10, bold=True, color=CORAL)
    add_text(s, left + Inches(0.3), ph_top + Inches(0.6), ph_w - Inches(0.6), Inches(0.5),
             name, size=16, bold=True, color=NAVY)
    add_rect(s, left + Inches(0.3), ph_top + Inches(1.2), ph_w - Inches(0.6), Emu(9525), MED_GRAY)
    add_text(s, left + Inches(0.3), ph_top + Inches(1.35), ph_w - Inches(0.6), Inches(1.7),
             what, size=12, color=DARK)
    rb_top = ph_top + Inches(3.1)
    add_rect(s, left + Inches(0.2), rb_top, ph_w - Inches(0.4), Inches(1.3), LIGHT_GRAY)
    add_text(s, left + Inches(0.35), rb_top + Inches(0.12), ph_w - Inches(0.7), Inches(0.3),
             "WATCH FOR", size=8, bold=True, color=CORAL)
    add_text(s, left + Inches(0.35), rb_top + Inches(0.38), ph_w - Inches(0.7), Inches(0.85),
             risk, size=10, color=CHARCOAL)

# ==== SLIDE 7: Governance, Risks & First Moves ====
s = add_blank()
slide_header(s, "06  |  Make it real", "Governance, the three risks that kill this, and the first moves")

lcol_x = Inches(0.5)
lcol_w = Inches(5.8)
lcol_top = Inches(1.55)

add_text(s, lcol_x, lcol_top, lcol_w, Inches(0.4),
         "GOVERNANCE \u2014 FOUR NAMED OWNERS, NO COMMITTEES", size=11, bold=True, color=CORAL)

owners = [
    ("Executive Sponsor", "One named VP. Air cover for protected runway. Authority to kill the program if quality slips."),
    ("Program Lead", "One named Director. End-to-end accountability. Rubric integrity. Calibration."),
    ("Tech Panel", "4\u20136 named Principals. Task design. Scoring. Defense panels."),
    ("TA Partner", "One named recruiter. Funnel mechanics. Candidate experience. Data."),
]
y = lcol_top + Inches(0.5)
for role, desc in owners:
    add_rect(s, lcol_x, y, Inches(0.08), Inches(0.95), CORAL)
    add_text(s, lcol_x + Inches(0.2), y, lcol_w - Inches(0.3), Inches(0.3),
             role, size=12, bold=True, color=NAVY)
    add_text(s, lcol_x + Inches(0.2), y + Inches(0.3), lcol_w - Inches(0.3), Inches(0.7),
             desc, size=11, color=CHARCOAL)
    y += Inches(1.05)

rcol_x = Inches(6.55)
rcol_w = Inches(6.28)
add_text(s, rcol_x, lcol_top, rcol_w, Inches(0.4),
         "THREE RISKS THAT WILL ACTUALLY KILL THIS", size=11, bold=True, color=CORAL)

risks = [
    "Delivery pressure breaks the protected runway in month 2.",
    "AI-use rubric drifts across panelists across cycles.",
    "Excellent ICs admitted with no senior IC ladder to graduate them onto.",
]
y = lcol_top + Inches(0.5)
for i, r in enumerate(risks, 1):
    add_rect(s, rcol_x, y, rcol_w, Inches(0.55), LIGHT_GRAY)
    add_text(s, rcol_x + Inches(0.2), y + Inches(0.13), Inches(0.4), Inches(0.3),
             str(i), size=14, bold=True, color=CORAL)
    add_text(s, rcol_x + Inches(0.55), y + Inches(0.16), rcol_w - Inches(0.7), Inches(0.35),
             r, size=11, color=DARK)
    y += Inches(0.65)

mv_top = Inches(4.5)
mv_h = Inches(2.4)
add_rect(s, rcol_x, mv_top, rcol_w, mv_h, NAVY)
add_rect(s, rcol_x, mv_top, Inches(0.12), mv_h, CORAL)
add_text(s, rcol_x + Inches(0.3), mv_top + Inches(0.18), rcol_w - Inches(0.5), Inches(0.4),
         "FIRST THREE MOVES", size=11, bold=True, color=CORAL)

moves = [
    ("This week", "Get one VP to sign the protected-runway commitment in writing. No signature, no program."),
    ("Next 2 weeks", "Tech panel writes one Ascender task and one Principal task. Three Principals submit benchmarks."),
    ("Weeks 3\u20134", "Run internal pilot with 20 nominees. Decide on external Cohort 1 from data, not enthusiasm."),
]
y = mv_top + Inches(0.6)
for when, what in moves:
    add_text(s, rcol_x + Inches(0.3), y, Inches(1.5), Inches(0.3),
             when.upper(), size=9, bold=True, color=CORAL)
    add_text(s, rcol_x + Inches(1.7), y - Inches(0.02), rcol_w - Inches(2.0), Inches(0.55),
             what, size=11, color=WHITE)
    y += Inches(0.55)

prs.save("10x_program_plan.pptx")
print("OK saved")
