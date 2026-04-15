"""Rewrite of build.py using ppt_runtime.

Same content, same structural patterns — cards, headers, stat grids,
columns — but composed via the runtime's named-anchor API with
Ascendion design-system tokens instead of inline hex/Inches constants.

This is the runtime validation gate (SLICE-006b): if ppt_runtime can
reproduce this deck, it can reproduce what the LLM builder will generate.
"""

from pathlib import Path
import sys

from pptx.util import Emu

from src.ppt_runtime import (
    Grid,
    Rect,
    Tokens,
    add_line,
    add_rect,
    add_text,
    draw_header_bar,
    load_template,
)

# ---------------------------------------------------------------------------
# Setup
# ---------------------------------------------------------------------------

BASE_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = BASE_DIR.parent
TEMPLATE_PATH = PROJECT_ROOT / "assets" / "template" / "template.pptx"
DS_PATH = PROJECT_ROOT / "assets" / "template" / "design_system.json"
DEFAULT_OUTPUT = BASE_DIR / "10x_program_plan.v3_runtime.pptx"
OUTPUT_PATH = Path(sys.argv[1]).expanduser().resolve() if len(sys.argv) > 1 else DEFAULT_OUTPUT
OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)

canvas = load_template(TEMPLATE_PATH, DS_PATH)
tokens = Tokens.from_design_system(DS_PATH)

# Shorthand colours
C1 = tokens.color("accent_1")       # primary brand (green)
C2 = tokens.color("accent_2")       # highlight / callout (magenta)
BG = tokens.color("bg_primary")     # white
BG_DARK = tokens.color("bg_dark")   # black
TXT = tokens.color("text_primary")  # body text
TXT2 = tokens.color("text_secondary")  # subtle text
TXT_W = tokens.color("text_on_dark")   # text on dark
MUTED = tokens.color("accent_5")    # muted fill
SUBTLE = tokens.color("accent_6")   # borders / deemphasized


# ---------------------------------------------------------------------------
# Helpers (thin wrappers the builder would write at the top of build_deck.py)
# ---------------------------------------------------------------------------

def hero_band(slide, grid, text, kicker_text=None):
    """Dark hero band spanning full width with optional kicker."""
    band = grid.span(col=1, col_span=12,
                     top=canvas.body_top, height_emu=tokens.spacing("xl") * 4)
    add_rect(slide, band, fill=C1)
    # Accent bar on left edge
    add_rect(slide, Rect(band.left, band.top, tokens.spacing("xs"), band.height), fill=C2)

    inner_left = band.left + tokens.spacing("md") * 2
    inner_w = band.width - tokens.spacing("md") * 4

    if kicker_text:
        k_rect = Rect(inner_left, band.top + tokens.spacing("md"), inner_w, tokens.spacing("lg"))
        add_text(slide, k_rect, kicker_text,
                 type_style=tokens.type("kicker"), color=C2)

    txt_top = band.top + (tokens.spacing("lg") if kicker_text else tokens.spacing("md"))
    txt_rect = Rect(inner_left, txt_top + tokens.spacing("sm"),
                    inner_w, band.height - txt_top + band.top - tokens.spacing("lg"))
    add_text(slide, txt_rect, text,
             font_name="Space Grotesk", font_size_pt=22, bold=True, color=TXT_W)
    return band


def support_card(slide, rect, heading, body):
    """Light card with accent bar, heading, and body."""
    add_rect(slide, rect, fill=MUTED)
    add_rect(slide, Rect(rect.left, rect.top, tokens.spacing("lg"), tokens.spacing("xs")), fill=C2)
    pad = tokens.spacing("md")
    add_text(slide, Rect(rect.left + pad, rect.top + pad, rect.width - 2 * pad, tokens.spacing("lg")),
             heading, type_style=tokens.type("kicker"), color=C1)
    add_text(slide, Rect(rect.left + pad, rect.top + pad + tokens.spacing("lg") + tokens.spacing("sm"),
                         rect.width - 2 * pad, rect.height - 3 * pad - tokens.spacing("lg")),
             body, type_style=tokens.type("body"), color=TXT)


def track_card(slide, rect, name, profile, rows):
    """Large card with dark header and label/value rows."""
    header_h = tokens.spacing("xl") + tokens.spacing("md")
    add_rect(slide, rect, fill=BG, line=SUBTLE)
    add_rect(slide, Rect(rect.left, rect.top, rect.width, header_h), fill=C1)
    pad = tokens.spacing("md")
    add_text(slide, Rect(rect.left + pad, rect.top + tokens.spacing("sm"),
                         rect.width - 2 * pad, tokens.spacing("lg")),
             name, font_name="Space Grotesk", font_size_pt=20, bold=True, color=TXT_W)
    add_text(slide, Rect(rect.left + pad, rect.top + tokens.spacing("lg") + tokens.spacing("sm"),
                         rect.width - 2 * pad, tokens.spacing("md")),
             profile, type_style=tokens.type("caption"), color=TXT_W)

    y = rect.top + header_h + tokens.spacing("md")
    for label, val in rows:
        add_text(slide, Rect(rect.left + pad, y, tokens.spacing("xl") * 3, tokens.spacing("lg")),
                 label, type_style=tokens.type("kicker"), color=C2)
        add_text(slide, Rect(rect.left + pad + tokens.spacing("xl") * 3, y,
                             rect.width - 3 * pad - tokens.spacing("xl") * 3,
                             tokens.spacing("xl") * 2),
                 val, type_style=tokens.type("body"), color=TXT)
        y += tokens.spacing("xl") * 2


def stage_card(slide, rect, num, when, name, what, funnel_label, funnel_val):
    """Funnel stage card with number, header, body, and bottom stats."""
    header_h = tokens.spacing("xl") + tokens.spacing("lg")
    pad = tokens.spacing("md")
    add_rect(slide, rect, fill=BG, line=SUBTLE)
    add_rect(slide, Rect(rect.left, rect.top, rect.width, header_h), fill=C1)
    add_text(slide, Rect(rect.left + pad, rect.top + tokens.spacing("sm"),
                         rect.width - 2 * pad, tokens.spacing("xl")),
             num, font_name="Space Grotesk", font_size_pt=28, bold=True, color=C2)
    add_text(slide, Rect(rect.left + pad, rect.top + tokens.spacing("xl"),
                         rect.width - 2 * pad, tokens.spacing("md")),
             when, type_style=tokens.type("kicker"), color=TXT_W)

    body_top = rect.top + header_h + tokens.spacing("sm")
    add_text(slide, Rect(rect.left + pad, body_top, rect.width - 2 * pad, tokens.spacing("lg")),
             name, type_style=tokens.type("subtitle"), color=C1, bold=True)
    add_text(slide, Rect(rect.left + pad, body_top + tokens.spacing("lg") + tokens.spacing("sm"),
                         rect.width - 2 * pad, tokens.spacing("xl") * 4),
             what, type_style=tokens.type("body"), color=TXT)

    # Bottom stat
    bot_h = tokens.spacing("xl") + tokens.spacing("lg")
    bot_top = rect.bottom - bot_h
    add_rect(slide, Rect(rect.left, bot_top, rect.width, bot_h), fill=MUTED)
    add_text(slide, Rect(rect.left + pad, bot_top + tokens.spacing("xs"),
                         rect.width - 2 * pad, tokens.spacing("md")),
             funnel_label, type_style=tokens.type("caption"), color=SUBTLE, bold=True)
    add_text(slide, Rect(rect.left + pad, bot_top + tokens.spacing("md") + tokens.spacing("xs"),
                         rect.width - 2 * pad, tokens.spacing("lg")),
             funnel_val, type_style=tokens.type("subtitle"), color=C1, bold=True)


def phase_card(slide, rect, when, name, what, risk):
    """Development phase card with accent bar, body, and risk callout."""
    pad = tokens.spacing("md")
    add_rect(slide, rect, fill=BG, line=SUBTLE)
    add_rect(slide, Rect(rect.left, rect.top, rect.width, tokens.spacing("xs")), fill=C2)

    add_text(slide, Rect(rect.left + pad, rect.top + tokens.spacing("md"),
                         rect.width - 2 * pad, tokens.spacing("md")),
             when, type_style=tokens.type("kicker"), color=C2)
    add_text(slide, Rect(rect.left + pad, rect.top + tokens.spacing("lg") + tokens.spacing("sm"),
                         rect.width - 2 * pad, tokens.spacing("xl")),
             name, type_style=tokens.type("subtitle"), color=C1, bold=True)

    # Divider line
    div_y = rect.top + tokens.spacing("xl") * 2 + tokens.spacing("md")
    add_line(slide, rect.left + pad, div_y, rect.right - pad, div_y, color=SUBTLE, width_pt=0.5)

    add_text(slide, Rect(rect.left + pad, div_y + tokens.spacing("sm"),
                         rect.width - 2 * pad, tokens.spacing("xl") * 3),
             what, type_style=tokens.type("body"), color=TXT)

    # Risk callout
    risk_h = tokens.spacing("xl") * 2 + tokens.spacing("md")
    risk_top = rect.bottom - risk_h - tokens.spacing("sm")
    risk_rect = Rect(rect.left + tokens.spacing("sm"), risk_top,
                     rect.width - tokens.spacing("md"), risk_h)
    add_rect(slide, risk_rect, fill=MUTED)
    add_text(slide, Rect(risk_rect.left + pad, risk_rect.top + tokens.spacing("xs"),
                         risk_rect.width - 2 * pad, tokens.spacing("md")),
             "WATCH FOR", type_style=tokens.type("caption"), color=C2, bold=True)
    add_text(slide, Rect(risk_rect.left + pad, risk_rect.top + tokens.spacing("md") + tokens.spacing("xs"),
                         risk_rect.width - 2 * pad, risk_h - tokens.spacing("lg")),
             risk, type_style=tokens.type("caption"), color=TXT2)


# ===========================================================================
# SLIDE 1: Operating Principle
# ===========================================================================
slide = canvas.add_slide("header_light")
g = Grid(canvas, cols=12, gutter="md")

draw_header_bar(slide, kicker="01  |  The thesis",
                title="What this program is for", canvas=canvas, tokens=tokens)

band = hero_band(slide, g,
    "Find engineers who use judgment and AI leverage to remove work —\n"
    "then put them on a track that doesn't destroy that capability.",
    kicker_text="OPERATING PRINCIPLE")

cols_data = [
    ("THE PROBLEM",
     "Conventional hiring rewards pedigree, narrow correctness, and confidence. "
     "None of these find force multipliers."),
    ("WHAT WE WON'T DO",
     "No essay prompts. No personality interviews. No steering committees. "
     "No '10x' marketing language externally."),
    ("WHAT WE WILL DO",
     "Verifiable artifacts. Scaffolded take-homes with AI allowed and scored. "
     "Anchored defense panels. Clear runway after admission to ensure continuity"),
]
card_top = band.bottom + tokens.spacing("md")
card_h = canvas.body_top + canvas.body_height - card_top
regions = g.row(top=card_top, height_emu=card_h, items=[(4, "c1"), (4, "c2"), (4, "c3")])
for (heading, body), name in zip(cols_data, ["c1", "c2", "c3"]):
    support_card(slide, regions[name], heading, body)


# ===========================================================================
# SLIDE 2: Two tracks
# ===========================================================================
slide = canvas.add_slide("header_light")
g = Grid(canvas, cols=12, gutter="md")

draw_header_bar(slide, kicker="02  |  Program shape",
                title="Two tracks. Small cohorts. On purpose.", canvas=canvas, tokens=tokens)

sub_rect = g.span(col=1, col_span=12, top=canvas.body_top, height_emu=tokens.spacing("lg"))
add_text(slide, sub_rect,
         "A 10x program with 200 admits is not a 10x program. We optimize for precision, not recall.",
         type_style=tokens.type("body"), color=TXT2)

card_top = canvas.body_top + tokens.spacing("xl")
card_h = canvas.body_top + canvas.body_height - card_top - tokens.spacing("lg")
lr = g.row(top=card_top, height_emu=card_h, items=[(6, "left"), (6, "right")])

track_card(slide, lr["left"], "ASCENDER TRACK",
           "0–2 yrs experience  |  external + internal",
           [("Cohort size", "10–15 per cohort"),
            ("Cadence", "Twice yearly"),
            ("Test for", "Decomposition under ambiguity. Taste in deciding what NOT to fix. "
                         "Healthy AI use the candidate can defend line by line.")])

track_card(slide, lr["right"], "PRINCIPAL TRACK",
           "2–6 yrs experience  |  lateral + skip-level nomination",
           [("Cohort size", "10–15 per cohort"),
            ("Cadence", "Twice yearly"),
            ("Test for", "Sequencing and trade-off articulation. Resisting the urge to rewrite. "
                         "Systems thinking, governance awareness, ability to teach.")])

# Footer note
footer = g.span(col=1, col_span=12,
                top=canvas.body_top + canvas.body_height - tokens.spacing("md"),
                height_emu=tokens.spacing("md"))
add_text(slide, footer,
         "Internal nominations run from day one and enter at Stage 2. Highest-yield source.",
         type_style=tokens.type("kicker"), color=C1)


# ===========================================================================
# SLIDE 3: The Funnel
# ===========================================================================
slide = canvas.add_slide("header_light")
g = Grid(canvas, cols=12, gutter="sm")

draw_header_bar(slide, kicker="03  |  The funnel",
                title="Four stages. Five weeks. ~150 in, 8–12 out.",
                canvas=canvas, tokens=tokens)

stages = [
    ("01", "WEEK 1", "EVIDENCE GATE",
     "GitHub link + 3-min Loom + one shipped artifact. No essays.",
     "FUNNEL", "1,500–15,000 → ~400"),
    ("02", "WEEK 2–3", "SCAFFOLDED TAKE-HOME",
     "48 hrs. Real repo. AI allowed and expected. PR + 10-min walkthrough.",
     "FUNNEL", "~400 → ~50"),
    ("03", "WEEK 4", "DEFENSE PANEL",
     "60 min. Two panelists. Five anchored questions. Independent scoring.",
     "FUNNEL", "~50 → ~15"),
    ("04", "WEEK 5", "ADMIT DECISION",
     "30 min. Director + Principal. Confirms appetite. No new tech eval.",
     "FUNNEL", "10–15 admits"),
]

card_top = canvas.body_top + tokens.spacing("sm")
card_h = canvas.body_top + canvas.body_height - card_top
regions = g.row(top=card_top, height_emu=card_h,
                items=[(3, "s1"), (3, "s2"), (3, "s3"), (3, "s4")])
for (num, when, name, what, fl, fv), key in zip(stages, ["s1", "s2", "s3", "s4"]):
    stage_card(slide, regions[key], num, when, name, what, fl, fv)


# ===========================================================================
# SLIDE 4: Discriminators
# ===========================================================================
slide = canvas.add_slide("header_light")
g = Grid(canvas, cols=12, gutter="md")

draw_header_bar(slide, kicker="04  |  Where this differs",
                title="The two design decisions that determine whether it works",
                canvas=canvas, tokens=tokens)

lr = g.row(top=canvas.body_top, height_emu=canvas.body_height,
           items=[(6, "left"), (6, "right")])

# Left: AI-use rubric
left_r = lr["left"]
pad = tokens.spacing("md")
add_text(slide, Rect(left_r.left, left_r.top, left_r.width, tokens.spacing("md")),
         "AI-USE RUBRIC", type_style=tokens.type("kicker"), color=C2)
add_text(slide, Rect(left_r.left, left_r.top + tokens.spacing("md"),
                     left_r.width, tokens.spacing("md")),
         "Scored explicitly. Without this, calibration collapses.",
         type_style=tokens.type("caption"), color=TXT2)

rubric_top = left_r.top + tokens.spacing("xl")
rubric_h = left_r.height - tokens.spacing("xl")
rubric_cols = [
    ("HEALTHY USE", C1,
     "Used AI for boilerplate, scaffolding, tests. Verified outputs. Can explain every line."),
    ("OVER-RELIANCE", C2,
     "Cannot explain choices. Inconsistent style. Hallucinated comments. Tests don't match code."),
    ("AVOIDANCE", SUBTLE,
     "Avoided AI in an AI-appropriate task. Slower with no quality gain."),
]
rub_gutter = tokens.spacing("xs")
rub_w = (left_r.width - 2 * rub_gutter) // 3
for i, (h, c, body) in enumerate(rubric_cols):
    rl = left_r.left + i * (rub_w + rub_gutter)
    r = Rect(rl, rubric_top, rub_w, rubric_h)
    add_rect(slide, r, fill=MUTED)
    add_rect(slide, Rect(rl, rubric_top, rub_w, tokens.spacing("lg")), fill=c)
    add_text(slide, Rect(rl + pad, rubric_top + tokens.spacing("xs"),
                         rub_w - 2 * pad, tokens.spacing("md")),
             h, type_style=tokens.type("caption"), color=TXT_W, bold=True)
    add_text(slide, Rect(rl + pad, rubric_top + tokens.spacing("lg") + tokens.spacing("sm"),
                         rub_w - 2 * pad, rubric_h - tokens.spacing("xl")),
             body, type_style=tokens.type("caption"), color=TXT)

# Right: Base rates
right_r = lr["right"]
add_text(slide, Rect(right_r.left, right_r.top, right_r.width, tokens.spacing("md")),
         "BASE RATES PER TRACK PER CYCLE", type_style=tokens.type("kicker"), color=C2)
add_text(slide, Rect(right_r.left, right_r.top + tokens.spacing("md"),
                     right_r.width, tokens.spacing("md")),
         "If they want recall, this is the wrong program.",
         type_style=tokens.type("caption"), color=TXT2)

stats = [
    ("1,500–15,000", "applicants"),
    ("~400", "Stage 2 invitees"),
    ("~50", "Stage 3 invitees"),
    ("10–15", "admits"),
    ("~6 hrs", "reviewer cost / admit"),
    ("25%", "false-positive tolerance"),
]
stat_top = right_r.top + tokens.spacing("xl")
stat_gutter = tokens.spacing("xs")
sw_ = (right_r.width - stat_gutter) // 2
sh_ = (right_r.height - tokens.spacing("xl") - 2 * stat_gutter) // 3
for idx, (big, label) in enumerate(stats):
    col, row = idx % 2, idx // 2
    sl = right_r.left + col * (sw_ + stat_gutter)
    st = stat_top + row * (sh_ + stat_gutter)
    cell = Rect(sl, st, sw_, sh_)
    add_rect(slide, cell, fill=BG, line=SUBTLE)
    add_rect(slide, Rect(sl, st, tokens.spacing("xs"), sh_), fill=C2)
    add_text(slide, Rect(sl + pad, st + tokens.spacing("sm"), sw_ - 2 * pad, tokens.spacing("xl")),
             big, font_name="Space Grotesk", font_size_pt=24, bold=True, color=C1)
    add_text(slide, Rect(sl + pad, st + tokens.spacing("xl") + tokens.spacing("sm"),
                         sw_ - 2 * pad, tokens.spacing("lg")),
             label, type_style=tokens.type("caption"), color=TXT2)


# ===========================================================================
# SLIDE 5: Development Arc
# ===========================================================================
slide = canvas.add_slide("header_light")
g = Grid(canvas, cols=12, gutter="md")

draw_header_bar(slide, kicker="05  |  After admit",
                title="Selection is 20%. The development arc is 80%.",
                canvas=canvas, tokens=tokens)

sub = g.span(col=1, col_span=12, top=canvas.body_top, height_emu=tokens.spacing("lg"))
add_text(slide, sub,
         "If we cannot protect the runway, we should not run the program. "
         "This is the failure mode that quietly kills these initiatives.",
         type_style=tokens.type("body"), color=TXT2)

phases = [
    ("MONTHS 1–3", "PROTECTED RUNWAY",
     "60% protected time. A real problem. A named sponsor. "
     "NOT on standard delivery utilization.",
     "Requires written executive air cover. A Delivery VP will try to break this in week 2."),
    ("MONTHS 3–9", "SHADOW + LEAD",
     "Pair with a Principal Engineer on a real client engagement. "
     "Lead one workstream end-to-end.",
     "Quarterly review on five anchored dimensions: framing, leverage, taste, ownership, teaching."),
    ("MONTHS 9–18", "FORCE MULTIPLIER",
     "Own a small team or a reusable practice asset — playbook, internal tool, eval harness.",
     "Promotion or off-ramp at month 18. Off-ramp to senior IC track is graduation, not failure."),
]

ph_top = canvas.body_top + tokens.spacing("xl")
ph_h = canvas.body_top + canvas.body_height - ph_top
ph_regions = g.row(top=ph_top, height_emu=ph_h, items=[(4, "p1"), (4, "p2"), (4, "p3")])
for (when, name, what, risk), key in zip(phases, ["p1", "p2", "p3"]):
    phase_card(slide, ph_regions[key], when, name, what, risk)


# ===========================================================================
# SLIDE 6: Governance, Risks & First Moves
# ===========================================================================
slide = canvas.add_slide("header_light")
g = Grid(canvas, cols=12, gutter="md")

draw_header_bar(slide, kicker="06  |  Make it real",
                title="Governance, the three risks that kill this, and the first moves",
                canvas=canvas, tokens=tokens)

lr = g.row(top=canvas.body_top, height_emu=canvas.body_height,
           items=[(5, "left"), (7, "right")])

# Left: Governance
left_r = lr["left"]
add_text(slide, Rect(left_r.left, left_r.top, left_r.width, tokens.spacing("md")),
         "GOVERNANCE — FOUR NAMED OWNERS, NO COMMITTEES",
         type_style=tokens.type("kicker"), color=C2)

owners = [
    ("Executive Sponsor",
     "One named VP. Air cover for protected runway. Authority to kill the program if quality slips."),
    ("Program Lead",
     "One named Director. End-to-end accountability. Rubric integrity. Calibration."),
    ("Tech Panel",
     "4–6 named Principals. Task design. Scoring. Defense panels."),
    ("TA Partner",
     "One named recruiter. Funnel mechanics. Candidate experience. Data."),
]
y = left_r.top + tokens.spacing("lg")
row_h = tokens.spacing("xl") * 2
for role, desc in owners:
    add_rect(slide, Rect(left_r.left, y, tokens.spacing("xs"), row_h), fill=C2)
    add_text(slide, Rect(left_r.left + tokens.spacing("sm"), y, left_r.width - tokens.spacing("md"),
                         tokens.spacing("md")),
             role, type_style=tokens.type("body"), color=C1, bold=True)
    add_text(slide, Rect(left_r.left + tokens.spacing("sm"), y + tokens.spacing("md"),
                         left_r.width - tokens.spacing("md"), row_h - tokens.spacing("md")),
             desc, type_style=tokens.type("caption"), color=TXT)
    y += row_h + tokens.spacing("xs")

# Right: Risks
right_r = lr["right"]
add_text(slide, Rect(right_r.left, right_r.top, right_r.width, tokens.spacing("md")),
         "THREE RISKS THAT WILL ACTUALLY KILL THIS",
         type_style=tokens.type("kicker"), color=C2)

risks = [
    "Delivery pressure breaks the protected runway in month 2.",
    "AI-use rubric drifts across panelists across cycles.",
    "Excellent ICs admitted with no senior IC ladder to graduate them onto.",
]
y = right_r.top + tokens.spacing("lg")
risk_h = tokens.spacing("xl")
for i, r in enumerate(risks, 1):
    add_rect(slide, Rect(right_r.left, y, right_r.width, risk_h), fill=MUTED)
    add_text(slide, Rect(right_r.left + tokens.spacing("sm"), y + tokens.spacing("xs"),
                         tokens.spacing("md"), tokens.spacing("md")),
             str(i), type_style=tokens.type("subtitle"), color=C2, bold=True)
    add_text(slide, Rect(right_r.left + tokens.spacing("lg"), y + tokens.spacing("xs"),
                         right_r.width - tokens.spacing("xl"), risk_h - tokens.spacing("sm")),
             r, type_style=tokens.type("caption"), color=TXT)
    y += risk_h + tokens.spacing("xs")

# First moves box
mv_top = y + tokens.spacing("sm")
mv_h = right_r.bottom - mv_top
add_rect(slide, Rect(right_r.left, mv_top, right_r.width, mv_h), fill=C1)
add_rect(slide, Rect(right_r.left, mv_top, tokens.spacing("xs"), mv_h), fill=C2)
add_text(slide, Rect(right_r.left + tokens.spacing("md"), mv_top + tokens.spacing("sm"),
                     right_r.width - tokens.spacing("lg"), tokens.spacing("md")),
         "FIRST THREE MOVES", type_style=tokens.type("kicker"), color=C2)

moves = [
    ("This week",
     "Get one VP to sign the protected-runway commitment in writing. No signature, no program."),
    ("Next 2 weeks",
     "Tech panel writes one Ascender task and one Principal task. Three Principals submit benchmarks."),
    ("Weeks 3–4",
     "Run internal pilot with 20 nominees. Decide on external Cohort 1 from data, not enthusiasm."),
]
y = mv_top + tokens.spacing("lg") + tokens.spacing("sm")
for when, what in moves:
    add_text(slide, Rect(right_r.left + tokens.spacing("md"), y,
                         tokens.spacing("xl") * 2 + tokens.spacing("lg"), tokens.spacing("md")),
             when, type_style=tokens.type("kicker"), color=C2)
    add_text(slide, Rect(right_r.left + tokens.spacing("xl") * 3, y,
                         right_r.width - tokens.spacing("xl") * 3 - tokens.spacing("md"),
                         tokens.spacing("xl")),
             what, type_style=tokens.type("caption"), color=TXT_W)
    y += tokens.spacing("xl")


# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
canvas.save(OUTPUT_PATH)
print(f"OK  saved {OUTPUT_PATH}")
