"""DeckIR to PPTX renderer."""

from __future__ import annotations

import json
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation

from ..models.deck_ir import DeckIR
from ..models.render_map import RenderMap, RenderMapEntry
from ..validate.drift import _read_alt_text
from pptx.oxml.ns import qn

RENDERABLE_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp"}

# Regex for parsing Markdown inline formatting into (text, bold, italic) spans.
_MD_INLINE_RE = re.compile(
    r"\*\*\*(.+?)\*\*\*"   # ***bold+italic***
    r"|\*\*(.+?)\*\*"      # **bold**
    r"|\*(.+?)\*"          # *italic*
    r"|([^*]+)"            # plain text
    r"|(\*)"               # lone asterisk
)


def parse_markdown_runs(text: str) -> List[Tuple[str, bool, bool]]:
    """Parse Markdown inline formatting into a list of (text, bold, italic) tuples."""
    if "**" not in text and "*" not in text:
        return [(text, False, False)]
    runs: List[Tuple[str, bool, bool]] = []
    for m in _MD_INLINE_RE.finditer(text):
        if m.group(1) is not None:
            runs.append((m.group(1), True, True))
        elif m.group(2) is not None:
            runs.append((m.group(2), True, False))
        elif m.group(3) is not None:
            runs.append((m.group(3), False, True))
        elif m.group(4) is not None:
            runs.append((m.group(4), False, False))
        elif m.group(5) is not None:
            runs.append((m.group(5), False, False))
    return runs


class Renderer:
    def __init__(
        self, template_path: Path, layout_catalog_path: Path, icons_json_path: Path
    ) -> None:
        self.template_path = template_path
        self.layout_catalog_path = layout_catalog_path
        self.icons_json_path = icons_json_path

    def render(self, deck: DeckIR, output_path: Path) -> RenderMap:
        """Render a DeckIR to PPTX and return a RenderMap."""
        layout_catalog = self._load_layout_catalog()
        icon_index = self._load_icon_index()

        prs = Presentation(str(self.template_path))
        self._remove_existing_slides(prs)
        render_map = RenderMap()

        for slide_spec in deck.slides:
            layout_entry = layout_catalog.get(slide_spec.layout_id)
            if not layout_entry:
                raise ValueError(f"Unknown layout_id: {slide_spec.layout_id}")

            master_index = layout_entry["master_index"]
            layout_index = layout_entry["layout_index"]
            layout = prs.slide_masters[master_index].slide_layouts[layout_index]
            slide = prs.slides.add_slide(layout)

            field_keys: List[str] = []
            field_key_by_idx = self._layout_field_key_by_idx(layout)
            alt_text_to_shape = {}
            for shape in slide.shapes:
                alt_text = _read_alt_text(shape)
                if not alt_text and shape.is_placeholder:
                    alt_text = field_key_by_idx.get(shape.placeholder_format.idx)
                    if alt_text:
                        self._set_alt_text(shape, alt_text)
                if not alt_text:
                    continue
                alt_text_to_shape[alt_text] = shape
                field_keys.append(alt_text)
                if alt_text in slide_spec.fields:
                    self._apply_text(shape, slide_spec.fields[alt_text])

            for asset_ref in slide_spec.asset_refs:
                if not asset_ref.target_field_key:
                    raise ValueError(
                        f"Asset ref missing target_field_key: {asset_ref.asset_id}"
                    )
                target_shape = alt_text_to_shape.get(asset_ref.target_field_key)
                if not target_shape:
                    raise ValueError(
                        f"Missing placeholder for asset target: {asset_ref.target_field_key}"
                    )
                asset_path = self._resolve_asset_path(asset_ref, icon_index)
                self._apply_image(slide, target_shape, asset_path)

            notes_text = ""
            if slide_spec.speaker_notes is not None:
                if isinstance(slide_spec.speaker_notes, dict):
                    notes_text = json.dumps(
                        slide_spec.speaker_notes, sort_keys=True, ensure_ascii=True
                    )
                else:
                    notes_text = str(slide_spec.speaker_notes)
            slide.notes_slide.notes_text_frame.text = notes_text

            render_map.entries[slide_spec.slide_id] = RenderMapEntry(
                slide_id=slide_spec.slide_id,
                slide_index=len(prs.slides) - 1,
                field_keys=sorted(set(field_keys)),
            )

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        return render_map

    def _apply_text(self, shape, value: Any) -> None:
        if not shape.has_text_frame:
            return
        text_frame = shape.text_frame
        text_frame.clear()
        if isinstance(value, list):
            if not value:
                return
            for idx, item in enumerate(value):
                paragraph = (
                    text_frame.paragraphs[0]
                    if idx == 0
                    else text_frame.add_paragraph()
                )
                paragraph.level = 0
                self._set_paragraph_runs(paragraph, str(item))
        else:
            paragraph = text_frame.paragraphs[0]
            self._set_paragraph_runs(paragraph, str(value))

    @staticmethod
    def _set_paragraph_runs(paragraph, text: str) -> None:
        """Replace paragraph content with Markdown-formatted runs."""
        # Clear existing runs
        for run in list(paragraph.runs):
            run._r.getparent().remove(run._r)
        runs = parse_markdown_runs(text)
        for run_text, bold, italic in runs:
            r = paragraph.add_run()
            r.text = run_text
            if bold:
                r.font.bold = True
            if italic:
                r.font.italic = True

    def _apply_image(self, slide, shape, asset_path: Path) -> None:
        try:
            shape.insert_picture(str(asset_path))
        except Exception:
            slide.shapes.add_picture(
                str(asset_path),
                shape.left,
                shape.top,
                width=shape.width,
                height=shape.height,
            )

    def _remove_existing_slides(self, prs: Presentation) -> None:
        while len(prs.slides) > 0:
            slide_id = prs.slides._sldIdLst[0]
            r_id = slide_id.rId
            prs.part.drop_rel(r_id)
            del prs.slides._sldIdLst[0]

    def _set_alt_text(self, shape, field_key: str) -> None:
        elem = shape.element
        nvSpPr = elem.find(qn("p:nvSpPr"))
        if nvSpPr is not None:
            cNvPr = nvSpPr.find(qn("p:cNvPr"))
            if cNvPr is not None:
                cNvPr.set("descr", field_key)
                return
        nvPicPr = elem.find(qn("p:nvPicPr"))
        if nvPicPr is not None:
            cNvPr = nvPicPr.find(qn("p:cNvPr"))
            if cNvPr is not None:
                cNvPr.set("descr", field_key)

    def _layout_field_key_by_idx(self, layout) -> Dict[int, str]:
        mapping: Dict[int, str] = {}
        for shape in layout.shapes:
            if not shape.is_placeholder:
                continue
            alt_text = _read_alt_text(shape)
            if not alt_text:
                continue
            mapping[shape.placeholder_format.idx] = alt_text
        return mapping

    def _load_layout_catalog(self) -> Dict[str, Dict[str, Any]]:
        with open(self.layout_catalog_path, "r", encoding="utf-8") as handle:
            catalog = json.load(handle)
        layouts = catalog.get("layouts", [])
        return {entry["layout_id"]: entry for entry in layouts}

    def _load_icon_index(self) -> Dict[str, str]:
        icon_index: Dict[str, str] = {}
        with open(self.icons_json_path, "r", encoding="utf-8") as handle:
            icons = json.load(handle)
        icon_entries = icons.get("icons", [])
        for entry in icon_entries:
            icon_id = entry["icon_id"]
            icon_index[icon_id] = f"assets/icons/png/{entry['filename']}"

        external_registry_path = (
            self.icons_json_path.parent.parent / "external_assets" / "registry.manifest.json"
        )
        if external_registry_path.exists():
            with open(external_registry_path, "r", encoding="utf-8") as handle:
                registry = json.load(handle)
            for icon in registry.get("icons", []):
                icon_id = icon.get("id")
                pack = icon.get("pack")
                svg_path = icon.get("svg_path")
                if not icon_id or not pack or not svg_path:
                    continue
                # Prefer local PNG mapping when present; fall back to external SVG only if missing.
                icon_index.setdefault(
                    str(icon_id), f"assets/external_assets/{pack}/{svg_path}"
                )

        return icon_index

    def _resolve_asset_path(self, asset_ref, icon_index: Dict[str, str]) -> Path:
        if asset_ref.asset_type == "icon":
            source_path = icon_index.get(asset_ref.asset_id)
            if not source_path:
                raise FileNotFoundError(
                    f"Unknown icon_id: {asset_ref.asset_id}"
                )
            project_root = self.template_path.parents[2]
            candidate = project_root / source_path
            if candidate.exists():
                if candidate.suffix.lower() not in RENDERABLE_IMAGE_SUFFIXES:
                    raise ValueError(
                        f"Unsupported icon format for PowerPoint render: {candidate.suffix} "
                        f"(icon_id={asset_ref.asset_id}). Use PNG/JPG/WebP in MVP."
                    )
                return candidate
            raise FileNotFoundError(f"Missing icon asset file: {candidate}")
        asset_path = Path(asset_ref.asset_id)
        if asset_path.is_absolute():
            return asset_path
        project_root = self.template_path.parents[2]
        candidate = project_root / asset_ref.asset_id
        if candidate.exists():
            return candidate
        assets_candidate = project_root / "assets" / asset_ref.asset_id
        if assets_candidate.exists():
            return assets_candidate
        raise FileNotFoundError(f"Missing asset: {asset_ref.asset_id}")
