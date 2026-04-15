"""Content fidelity checker: compares user input facts against PPTX output.

Usage::

    report = check_content_fidelity("Our Q3 revenue grew 14%.", "output.pptx")
    assert report["visible_coverage_score"] > 0.8
"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation


# ---------------------------------------------------------------------------
# Fact extraction
# ---------------------------------------------------------------------------

# Patterns to extract fact-like tokens from user input (order matters for
# priority — earlier patterns take precedence over later overlapping ones).
_FACT_PATTERNS = [
    re.compile(r"\$[\d,]+(?:\.\d+)?[BMKbmk]?"),               # dollar amounts
    re.compile(r"\d+(?:\.\d+)?%"),                              # percentages
    re.compile(r"\b(?:Q[1-4])\s*\d{4}\b"),                      # quarter+year (Q3 2024)
    re.compile(r"\b(?:Q[1-4])\b"),                              # quarters
    re.compile(r"\b\d{4}\b"),                                    # years
    re.compile(r"\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b"),          # dates (01/15/2024)
    re.compile(r"\b\d+(?:,\d{3})*(?:\.\d+)?\b"),               # numbers
    re.compile(r'"[^"]{3,}"'),                                   # double-quoted phrases
    re.compile(r"'[^']{3,}'"),                                   # single-quoted phrases
]

# Month names for date extraction
_MONTHS = (
    "January|February|March|April|May|June|July|August|September|"
    "October|November|December|Jan|Feb|Mar|Apr|Jun|Jul|Aug|Sep|Oct|Nov|Dec"
)
_DATE_PATTERNS = [
    re.compile(rf"\b(?:{_MONTHS})\s+\d{{1,2}},?\s+\d{{4}}\b"),  # January 15, 2024
    re.compile(rf"\b\d{{1,2}}\s+(?:{_MONTHS})\s+\d{{4}}\b"),     # 15 January 2024
    re.compile(rf"\b(?:{_MONTHS})\s+\d{{4}}\b"),                  # March 2024
]

# Proper nouns: capitalized words not at sentence start, at least 2 chars,
# excluding common title words that are not proper nouns.
_COMMON_TITLE_WORDS = {
    "The", "This", "That", "These", "Those", "Our", "Your", "Their",
    "Its", "His", "Her", "Each", "Every", "Both", "All", "Any", "Some",
    "Most", "Many", "Few", "Several", "Other", "Such", "What", "Which",
    "Who", "How", "When", "Where", "Why", "Here", "There", "Now",
    "Also", "However", "Therefore", "Furthermore", "Moreover",
    "Additionally", "Currently", "Previously", "Next", "First",
    "Second", "Third", "Finally", "Meanwhile", "Instead", "Otherwise",
    "Yes", "No", "Not", "But", "And", "For", "With",
}


def extract_facts(text: str) -> list[str]:
    """Extract factual tokens from user input text.

    Returns a deduplicated list of fact strings: numbers, percentages,
    dates, proper nouns, and quoted phrases.
    """
    facts = []
    seen = set()

    def _add(fact: str):
        if fact and fact not in seen and len(fact) > 1:
            seen.add(fact)
            facts.append(fact)

    # Extract pattern-based facts
    for pattern in _FACT_PATTERNS:
        for match in pattern.finditer(text):
            _add(match.group(0).strip("'\""))

    # Extract dates with month names
    for pattern in _DATE_PATTERNS:
        for match in pattern.finditer(text):
            _add(match.group(0))

    # Extract proper nouns (capitalized words, filtered by exclusion list)
    for word in text.split():
        clean = re.sub(r'[^\w]', '', word)
        if (clean
                and clean[0].isupper()
                and len(clean) >= 2
                and clean not in seen
                and clean not in _COMMON_TITLE_WORDS):
            _add(clean)

    return facts


# ---------------------------------------------------------------------------
# Text extraction from PPTX
# ---------------------------------------------------------------------------

def _extract_visible_text(prs: Presentation) -> str:
    """Extract all visible text from slide shapes."""
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n".join(texts)


def _extract_notes_text(prs: Presentation) -> str:
    """Extract text from slide notes."""
    texts = []
    for slide in prs.slides:
        if slide.has_notes_slide:
            notes = slide.notes_slide
            if notes.notes_text_frame:
                texts.append(notes.notes_text_frame.text)
    return "\n".join(texts)


# ---------------------------------------------------------------------------
# Detection patterns (shared with scanner)
# ---------------------------------------------------------------------------

_PLACEHOLDER_PATTERNS = [
    re.compile(r"\{title\}", re.IGNORECASE),
    re.compile(r"\{body\}", re.IGNORECASE),
    re.compile(r"Lorem ipsum", re.IGNORECASE),
    re.compile(r"\bTODO\b"),
    re.compile(r"\[insert\]", re.IGNORECASE),
    re.compile(r"\bTBD\b"),
    re.compile(r"\{\{.*?\}\}"),
]

_MARKDOWN_PATTERNS = [
    re.compile(r"\*\*"),
    re.compile(r"__"),
    re.compile(r"##"),
    re.compile(r"```"),
    re.compile(r"\[.*?\]\(.*?\)"),
    re.compile(r"(?m)^- "),
]

# Common numbers/ordinals excluded from hallucination detection
_COMMON_NUMBERS = {
    "0", "1", "2", "3", "4", "5", "6", "7", "8", "9", "10",
    "100", "1st", "2nd", "3rd",
}

_NUMBER_RE = re.compile(r"\b\d+(?:,\d{3})*(?:\.\d+)?%?\b")
_DOLLAR_RE = re.compile(r"\$[\d,]+(?:\.\d+)?[BMKbmk]?")
_QUOTED_RE = re.compile(r'"([^"]{3,})"')


def _extract_output_proper_nouns(text: str) -> set[str]:
    """Extract proper nouns from output text for hallucination checking."""
    nouns = set()
    for word in text.split():
        clean = re.sub(r'[^\w]', '', word)
        if (clean
                and clean[0].isupper()
                and len(clean) >= 3
                and clean not in _COMMON_TITLE_WORDS):
            nouns.add(clean)
    return nouns


def _extract_dates(text: str) -> set[str]:
    """Extract date-like strings from text."""
    dates = set()
    for pattern in _DATE_PATTERNS:
        for match in pattern.finditer(text):
            dates.add(match.group(0))
    # Also numeric dates
    for match in re.finditer(r"\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b", text):
        dates.add(match.group(0))
    return dates


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def check_content_fidelity(user_input: str, pptx_path: str | Path) -> dict:
    """Check how faithfully a PPTX represents the user's input content.

    Args:
        user_input: The original user text/prompt.
        pptx_path: Path to the generated PPTX file.

    Returns:
        A content_fidelity_report dict.
    """
    pptx_path = Path(pptx_path)
    prs = Presentation(str(pptx_path))

    # Extract text
    visible_text = _extract_visible_text(prs)
    notes_text = _extract_notes_text(prs)

    # Extract facts from user input
    facts = extract_facts(user_input)
    total_facts = len(facts)

    # Match facts against visible text (case-insensitive substring)
    matched_visible = []
    matched_notes_only = []
    dropped = []

    for fact in facts:
        fact_lower = fact.lower()
        if fact_lower in visible_text.lower():
            matched_visible.append(fact)
        elif fact_lower in notes_text.lower():
            matched_notes_only.append(fact)
        else:
            dropped.append(fact)

    # Coverage score
    visible_coverage = len(matched_visible) / total_facts if total_facts > 0 else 1.0

    # Detect placeholders in visible text
    placeholder_leaks = []
    for pattern in _PLACEHOLDER_PATTERNS:
        for m in pattern.finditer(visible_text):
            placeholder_leaks.append(m.group(0))

    # Detect markdown leaks
    markdown_leaks = []
    for pattern in _MARKDOWN_PATTERNS:
        for m in pattern.finditer(visible_text):
            markdown_leaks.append(m.group(0))

    # --- Hallucination detection ---
    # Numbers/percentages in visible text not in user input
    input_numbers = set()
    for m in _NUMBER_RE.finditer(user_input):
        input_numbers.add(m.group(0))

    hallucinated = []

    # Numbers
    for m in _NUMBER_RE.finditer(visible_text):
        num = m.group(0)
        if num not in input_numbers and num not in _COMMON_NUMBERS:
            if re.match(r"^\d+(?:st|nd|rd|th)$", num):
                continue
            hallucinated.append(num)

    # Dollar amounts in output not in input
    input_dollars = {m.group(0) for m in _DOLLAR_RE.finditer(user_input)}
    for m in _DOLLAR_RE.finditer(visible_text):
        val = m.group(0)
        if val not in input_dollars:
            hallucinated.append(val)

    # Quoted phrases in output not in input
    input_quotes = {m.group(1) for m in _QUOTED_RE.finditer(user_input)}
    for m in _QUOTED_RE.finditer(visible_text):
        phrase = m.group(1)
        if phrase not in input_quotes:
            hallucinated.append(f'"{phrase}"')

    # Dates in output not in input
    input_dates = _extract_dates(user_input)
    output_dates = _extract_dates(visible_text)
    for d in output_dates - input_dates:
        hallucinated.append(d)

    # Proper nouns in output not in input
    input_nouns = set()
    for word in user_input.split():
        clean = re.sub(r'[^\w]', '', word)
        if clean and clean[0].isupper() and len(clean) >= 3:
            input_nouns.add(clean)

    output_nouns = _extract_output_proper_nouns(visible_text)
    for noun in output_nouns - input_nouns:
        hallucinated.append(noun)

    # Deduplicate while preserving order
    hallucinated = list(dict.fromkeys(hallucinated))
    placeholder_leaks = list(dict.fromkeys(placeholder_leaks))
    markdown_leaks = list(dict.fromkeys(markdown_leaks))

    return {
        "visible_coverage_score": round(visible_coverage, 2),
        "notes_only_fact_count": len(matched_notes_only),
        "total_facts": total_facts,
        "matched_visible_facts": len(matched_visible),
        "matched_notes_only_facts": matched_notes_only,
        "dropped_facts": dropped,
        "hallucinated_specifics": hallucinated,
        "placeholder_leaks": placeholder_leaks,
        "markdown_leaks": markdown_leaks,
    }
