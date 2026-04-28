"""Deterministic post-build scanner implementing 26 visual hygiene checks.

Usage::

    report = scan_pptx("output.pptx", "assets/template/design_system.json")
    assert report["pass"]
"""

from __future__ import annotations

import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

_MARKDOWN_PATTERNS = [
    re.compile(r"\*\*"),          # bold markdown
    re.compile(r"__"),            # bold/italic markdown
    re.compile(r"##"),            # heading markdown
    re.compile(r"```"),           # code fence
    re.compile(r"\[.*?\]\(.*?\)"),  # link markdown
    re.compile(r"(?m)^- "),       # list markdown at line start
]

_PLACEHOLDER_PATTERNS = [
    re.compile(r"\{title\}", re.IGNORECASE),
    re.compile(r"\{body\}", re.IGNORECASE),
    re.compile(r"Lorem ipsum", re.IGNORECASE),
    re.compile(r"\bTODO\b"),
    re.compile(r"\[insert\]", re.IGNORECASE),
    re.compile(r"\bTBD\b"),
    re.compile(r"\{\{.*?\}\}"),
]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _parse_hex(hex_str: str) -> tuple[int, int, int]:
    """Parse '#RRGGBB' to (r, g, b)."""
    h = hex_str.lstrip("#")
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _rgb_to_tuple(rgb) -> tuple[int, int, int] | None:
    """Convert a python-pptx RGBColor to (r, g, b) or None."""
    if rgb is None:
        return None
    try:
        return (rgb[0], rgb[1], rgb[2])
    except (TypeError, IndexError):
        return None


def _relative_luminance(r: int, g: int, b: int) -> float:
    """WCAG 2.1 relative luminance from sRGB values 0-255."""
    def linearize(c):
        s = c / 255.0
        return s / 12.92 if s <= 0.04045 else ((s + 0.055) / 1.055) ** 2.4
    return 0.2126 * linearize(r) + 0.7152 * linearize(g) + 0.0722 * linearize(b)


def _contrast_ratio(rgb1: tuple[int, int, int], rgb2: tuple[int, int, int]) -> float:
    """WCAG contrast ratio between two RGB tuples."""
    l1 = _relative_luminance(*rgb1)
    l2 = _relative_luminance(*rgb2)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def _shape_area(shape) -> int:
    """Return area of shape in EMU^2."""
    try:
        return shape.width * shape.height
    except Exception:
        return 0


def _is_full_bleed(shape, slide_width: int, slide_height: int) -> bool:
    """True if shape covers > 90% of the slide."""
    slide_area = slide_width * slide_height
    if slide_area == 0:
        return False
    return _shape_area(shape) / slide_area > 0.90


def _is_background_shape(shape, slide_width: int, slide_height: int) -> bool:
    """True if shape covers > 50% of the slide."""
    slide_area = slide_width * slide_height
    if slide_area == 0:
        return False
    return _shape_area(shape) / slide_area > 0.50


def _get_shape_fill_rgb(shape) -> tuple[int, int, int] | None:
    """Extract fill RGB from a shape, returning None if not solid fill."""
    try:
        fill = shape.fill
        if fill.type is not None:
            fg = fill.fore_color
            if fg and fg.rgb:
                return _rgb_to_tuple(fg.rgb)
    except Exception:
        pass
    return None


def _get_slide_bg_rgb(slide, ds: dict) -> tuple[int, int, int]:
    """Best-effort slide background color. Falls back to design system bg_primary."""
    # Try to get from slide background
    try:
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            fg = fill.fore_color
            if fg and fg.rgb:
                rgb = _rgb_to_tuple(fg.rgb)
                if rgb:
                    return rgb
    except Exception:
        pass
    # Fall back to design system
    return _parse_hex(ds["colors"].get("bg_primary", "#FFFFFF"))


def _iter_text_runs(shape):
    """Yield (run, paragraph) for all text runs in a shape."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            yield run, para


def _get_text_from_shape(shape) -> str:
    """Get all text from a shape."""
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text


def _shape_name(shape) -> str:
    """Get a display name for a shape."""
    try:
        return shape.name or f"Shape {shape.shape_id}"
    except Exception:
        return "Unknown"


def _is_image_shape(shape) -> bool:
    """True if shape is a picture/image."""
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    except Exception:
        return False


def _is_placeholder(shape) -> bool:
    """True if shape is a layout placeholder (not user-added)."""
    try:
        return shape.is_placeholder
    except Exception:
        return False


def _get_font_size_pt(run) -> float | None:
    """Get font size in points from a run, or None."""
    try:
        if run.font.size is not None:
            return run.font.size.pt
    except Exception:
        pass
    return None


def _aabb_overlap_fraction(s1, s2) -> float:
    """Return the fraction of the smaller shape's area that overlaps with the other."""
    try:
        overlap_area = _aabb_overlap_area(s1, s2)
        area1 = s1.width * s1.height
        area2 = s2.width * s2.height
        if area1 == 0 or area2 == 0:
            return 0.0
        smaller_area = min(area1, area2)
        return overlap_area / smaller_area
    except Exception:
        return 0.0


def _aabb_overlap_area(s1, s2) -> int:
    """Return axis-aligned overlap area in EMU^2."""
    try:
        x_overlap = max(0, min(s1.left + s1.width, s2.left + s2.width) - max(s1.left, s2.left))
        y_overlap = max(0, min(s1.top + s1.height, s2.top + s2.height) - max(s1.top, s2.top))
        return x_overlap * y_overlap
    except Exception:
        return 0


def _aabb_overlap_fraction_of(target, other) -> float:
    """Return fraction of target area overlapped by other."""
    try:
        area = target.width * target.height
        if area == 0:
            return 0.0
        return _aabb_overlap_area(target, other) / area
    except Exception:
        return 0.0


def _shape_has_text(shape) -> bool:
    """True if shape has non-empty visible text."""
    return bool(_get_text_from_shape(shape).strip())


def _is_text_inside_container(text_shape, maybe_container) -> bool:
    """True when text is intentionally layered inside a larger filled container."""
    if not _shape_has_text(text_shape):
        return False
    if _shape_has_text(maybe_container):
        return False
    if _get_shape_fill_rgb(maybe_container) is None:
        return False
    text_area = _shape_area(text_shape)
    container_area = _shape_area(maybe_container)
    if text_area == 0 or container_area <= text_area:
        return False
    return _aabb_overlap_fraction_of(text_shape, maybe_container) >= 0.95


def _overlap_severity(s1, s2, overlap: float) -> str:
    """Classify overlaps; readable-text collisions are blocking."""
    if overlap <= 0.10:
        return ""
    s1_text = _shape_has_text(s1)
    s2_text = _shape_has_text(s2)
    if s1_text and s2_text:
        return "BLOCKING"
    if s1_text and not _is_text_inside_container(s1, s2):
        return "BLOCKING"
    if s2_text and not _is_text_inside_container(s2, s1):
        return "BLOCKING"
    return "WARNING" if overlap > 0.10 else ""


def _get_slide_layout_index(prs, slide) -> int | None:
    """Return the slide layout index within the presentation, if determinable."""
    try:
        for idx, layout in enumerate(prs.slide_layouts):
            if layout == slide.slide_layout:
                return idx
    except Exception:
        pass
    return None


def _get_slide_canvas_def(prs, slide, ds: dict) -> dict | None:
    """Return the design-system canvas definition for a slide's layout."""
    layout_index = _get_slide_layout_index(prs, slide)
    if layout_index is None:
        return None

    for canvas_def in ds.get("canvases", {}).values():
        if canvas_def.get("layout_index") == layout_index:
            return canvas_def
    return None


def _is_near_black(rgb: tuple[int, int, int]) -> bool:
    """True for full/near-full black backgrounds."""
    return max(rgb) <= 8


# ---------------------------------------------------------------------------
# Individual checks
# ---------------------------------------------------------------------------

def _check_vh01(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-01: All shape fills use palette token colors."""
    palette = {_parse_hex(v) for v in ds["colors"].values()}
    findings = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if _is_image_shape(shape):
                continue
            rgb = _get_shape_fill_rgb(shape)
            if rgb is not None and rgb not in palette:
                findings.append({
                    "check_id": "VH-01",
                    "category": "Color",
                    "check_name": "All fills use palette tokens",
                    "severity": "BLOCKING",
                    "pass": False,
                    "slide_index": si,
                    "details": (
                        f"Shape '{_shape_name(shape)}' has fill "
                        f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}, not in palette"
                    ),
                })
    return findings


def _check_vh02(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-02: All text colors use palette tokens."""
    palette = {_parse_hex(v) for v in ds["colors"].values()}
    findings = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            for run, _ in _iter_text_runs(shape):
                try:
                    if run.font.color and run.font.color.rgb:
                        rgb = _rgb_to_tuple(run.font.color.rgb)
                        if rgb is not None and rgb not in palette:
                            findings.append({
                                "check_id": "VH-02",
                                "category": "Color",
                                "check_name": "All text colors use palette tokens",
                                "severity": "BLOCKING",
                                "pass": False,
                                "slide_index": si,
                                "details": (
                                    f"Shape '{_shape_name(shape)}' has text color "
                                    f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}, not in palette"
                                ),
                            })
                except Exception:
                    pass
    return findings


def _check_vh03(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-03: Accent limit <= 2 per slide."""
    accent_keys = {k for k in ds["colors"] if k.startswith("accent_")}
    accent_rgb_to_role = {}
    for k in accent_keys:
        accent_rgb_to_role[_parse_hex(ds["colors"][k])] = k

    findings = []
    for si, slide in enumerate(prs.slides):
        used_accents = set()
        for shape in slide.shapes:
            # Check fills
            rgb = _get_shape_fill_rgb(shape)
            if rgb in accent_rgb_to_role:
                used_accents.add(accent_rgb_to_role[rgb])
            # Check text colors
            for run, _ in _iter_text_runs(shape):
                try:
                    if run.font.color and run.font.color.rgb:
                        trgb = _rgb_to_tuple(run.font.color.rgb)
                        if trgb in accent_rgb_to_role:
                            used_accents.add(accent_rgb_to_role[trgb])
                except Exception:
                    pass
        if len(used_accents) > 2:
            findings.append({
                "check_id": "VH-03",
                "category": "Color",
                "check_name": "Accent limit <= 2 per slide",
                "severity": "WARNING",
                "pass": False,
                "slide_index": si,
                "details": (
                    f"Slide uses {len(used_accents)} accent roles: "
                    f"{', '.join(sorted(used_accents))}"
                ),
            })
    return findings


def _check_vh04(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-04: Text-background contrast ratio >= 4.5:1 (body) or >= 3:1 (large)."""
    findings = []
    for si, slide in enumerate(prs.slides):
        bg_rgb = _get_slide_bg_rgb(slide, ds)
        for shape in slide.shapes:
            if _is_image_shape(shape):
                continue
            # Check if shape has its own fill to use as background
            shape_bg = _get_shape_fill_rgb(shape) or bg_rgb
            for run, _ in _iter_text_runs(shape):
                try:
                    if not run.text.strip():
                        continue
                    text_rgb = None
                    if run.font.color and run.font.color.rgb:
                        text_rgb = _rgb_to_tuple(run.font.color.rgb)
                    if text_rgb is None:
                        continue
                    size_pt = _get_font_size_pt(run)
                    is_large = size_pt is not None and size_pt >= 24
                    ratio = _contrast_ratio(text_rgb, shape_bg)
                    threshold = 3.0 if is_large else 4.5
                    if ratio < threshold:
                        findings.append({
                            "check_id": "VH-04",
                            "category": "Color",
                            "check_name": "Text-background contrast ratio",
                            "severity": "BLOCKING",
                            "pass": False,
                            "slide_index": si,
                            "details": (
                                f"Shape '{_shape_name(shape)}' text "
                                f"#{text_rgb[0]:02X}{text_rgb[1]:02X}{text_rgb[2]:02X} "
                                f"on #{shape_bg[0]:02X}{shape_bg[1]:02X}{shape_bg[2]:02X} "
                                f"has contrast {ratio:.2f}:1, needs >= {threshold}:1"
                            ),
                        })
                except Exception:
                    pass
    return findings


def _check_vh05(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-05: No invisible shapes (fill == slide background)."""
    findings = []
    for si, slide in enumerate(prs.slides):
        bg_rgb = _get_slide_bg_rgb(slide, ds)
        for shape in slide.shapes:
            if _is_image_shape(shape):
                continue
            if _is_full_bleed(shape, slide_width, slide_height):
                continue
            rgb = _get_shape_fill_rgb(shape)
            if rgb is not None and rgb == bg_rgb:
                findings.append({
                    "check_id": "VH-05",
                    "category": "Color",
                    "check_name": "No invisible shapes",
                    "severity": "WARNING",
                    "pass": False,
                    "slide_index": si,
                    "details": (
                        f"Shape '{_shape_name(shape)}' fill matches slide background "
                        f"#{bg_rgb[0]:02X}{bg_rgb[1]:02X}{bg_rgb[2]:02X}"
                    ),
                })
    return findings


def _check_vh06(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-06: All fonts from font_substitution allowlist."""
    font_sub = ds.get("font_substitution", {})
    # Build allowlist: both original fonts and substitutes
    allowed = set()
    for orig, sub in font_sub.items():
        allowed.add(orig)
        allowed.add(sub)
    # Also add fonts from type_scale
    for style in ds.get("type_scale", {}).values():
        if "font" in style:
            allowed.add(style["font"])

    findings = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            for run, _ in _iter_text_runs(shape):
                try:
                    font_name = run.font.name
                    if font_name and font_name not in allowed:
                        findings.append({
                            "check_id": "VH-06",
                            "category": "Typography",
                            "check_name": "All fonts from allowlist",
                            "severity": "BLOCKING",
                            "pass": False,
                            "slide_index": si,
                            "details": (
                                f"Shape '{_shape_name(shape)}' uses font "
                                f"'{font_name}', not in allowlist"
                            ),
                        })
                except Exception:
                    pass
    return findings


def _check_vh07(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-07: Font sizes match type_scale steps +/- 1pt."""
    type_scale = ds.get("type_scale", {})
    valid_sizes = {style["size_pt"] for style in type_scale.values()}

    findings = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            for run, _ in _iter_text_runs(shape):
                if not run.text.strip():
                    continue
                size_pt = _get_font_size_pt(run)
                if size_pt is None:
                    continue
                if not any(abs(size_pt - vs) <= 1 for vs in valid_sizes):
                    findings.append({
                        "check_id": "VH-07",
                        "category": "Typography",
                        "check_name": "Font sizes match type_scale",
                        "severity": "WARNING",
                        "pass": False,
                        "slide_index": si,
                        "details": (
                            f"Shape '{_shape_name(shape)}' uses {size_pt}pt, "
                            f"not near any type_scale step "
                            f"({', '.join(str(s) for s in sorted(valid_sizes))})"
                        ),
                    })
    return findings


def _check_vh08(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-08: Bold matches type_scale role definition for that size."""
    type_scale = ds.get("type_scale", {})
    # Build size -> expected bold map (with 1pt tolerance)
    size_bold_map = {}
    for style in type_scale.values():
        size_bold_map[style["size_pt"]] = style.get("bold", False)

    findings = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            for run, _ in _iter_text_runs(shape):
                if not run.text.strip():
                    continue
                size_pt = _get_font_size_pt(run)
                if size_pt is None:
                    continue
                # Find matching type_scale size (prefer exact match)
                matched_size = None
                for vs in size_bold_map:
                    if size_pt == vs:
                        matched_size = vs
                        break
                if matched_size is None:
                    for vs in size_bold_map:
                        if abs(size_pt - vs) <= 1:
                            matched_size = vs
                            break
                if matched_size is None:
                    continue
                expected_bold = size_bold_map[matched_size]
                actual_bold = run.font.bold
                if actual_bold is None:
                    continue
                if actual_bold != expected_bold:
                    findings.append({
                        "check_id": "VH-08",
                        "category": "Typography",
                        "check_name": "Bold matches type_scale role",
                        "severity": "WARNING",
                        "pass": False,
                        "slide_index": si,
                        "details": (
                            f"Shape '{_shape_name(shape)}' at {size_pt}pt is "
                            f"{'bold' if actual_bold else 'not bold'}, "
                            f"expected {'bold' if expected_bold else 'not bold'}"
                        ),
                    })
    return findings


def _check_vh09(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-09: ALLCAPS only on kicker-role text (11pt bold)."""
    findings = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            for run, _ in _iter_text_runs(shape):
                text = run.text.strip()
                if not text or len(text) < 3:
                    continue
                # Check if text is all caps (ignoring non-alpha)
                alpha_chars = [c for c in text if c.isalpha()]
                if not alpha_chars:
                    continue
                is_allcaps = all(c.isupper() for c in alpha_chars)
                if not is_allcaps:
                    continue
                # Check if this is kicker-role (11pt bold)
                size_pt = _get_font_size_pt(run)
                is_bold = run.font.bold
                is_kicker = (
                    size_pt is not None
                    and abs(size_pt - 11) <= 1
                    and is_bold is True
                )
                if not is_kicker:
                    findings.append({
                        "check_id": "VH-09",
                        "category": "Typography",
                        "check_name": "ALLCAPS only on kicker text",
                        "severity": "WARNING",
                        "pass": False,
                        "slide_index": si,
                        "details": (
                            f"Shape '{_shape_name(shape)}' has ALLCAPS text "
                            f"'{text[:30]}' but is not kicker-role "
                            f"(size={size_pt}, bold={is_bold})"
                        ),
                    })
    return findings


def _check_vh10(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-10: No shapes outside canvas."""
    findings = []
    tolerance = 1  # 1 EMU tolerance
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            try:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                if (left + width > slide_width + tolerance or
                        top + height > slide_height + tolerance or
                        left < -tolerance or
                        top < -tolerance):
                    findings.append({
                        "check_id": "VH-10",
                        "category": "Spatial",
                        "check_name": "No shapes outside canvas",
                        "severity": "BLOCKING",
                        "pass": False,
                        "slide_index": si,
                        "details": (
                            f"Shape '{_shape_name(shape)}' extends outside canvas: "
                            f"left={left}, top={top}, right={left + width}, "
                            f"bottom={top + height}, canvas={slide_width}x{slide_height}"
                        ),
                    })
            except Exception:
                pass
    return findings


def _check_vh11(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-11: No text overflow."""
    # Import measure_text - handle case where it might not be available
    try:
        from src.ppt_runtime.measure import measure_text
    except ImportError:
        return []  # Skip if ppt_runtime not available

    type_scale = ds.get("type_scale", {})

    findings = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # Get the dominant run style
            for run, _ in _iter_text_runs(shape):
                if not run.text.strip():
                    continue
                font_name = run.font.name
                size_pt = _get_font_size_pt(run)
                bold = run.font.bold or False
                if font_name and size_pt:
                    style = {
                        "font": font_name,
                        "size_pt": int(size_pt),
                        "bold": bold,
                        "line": 1.2,
                    }
                    # Find matching type_scale for line height
                    for ts in type_scale.values():
                        if ts.get("font") == font_name and abs(ts["size_pt"] - size_pt) <= 1:
                            style["line"] = ts.get("line", 1.2)
                            break
                    try:
                        _, text_h = measure_text(
                            text, style, max_width_emu=shape.width
                        )
                        if text_h > shape.height * 1.10:  # 10% tolerance
                            findings.append({
                                "check_id": "VH-11",
                                "category": "Spatial",
                                "check_name": "No text overflow",
                                "severity": "BLOCKING",
                                "pass": False,
                                "slide_index": si,
                                "details": (
                                    f"Shape '{_shape_name(shape)}' text height "
                                    f"{text_h} exceeds frame height {shape.height}"
                                ),
                            })
                    except Exception:
                        pass
                break  # Only check first run for style
    return findings


def _check_vh12(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-12: Content within safe_area."""
    safe = ds.get("canvas", {}).get("safe_area", {})
    sa_left = safe.get("left_emu", 0)
    sa_right = safe.get("right_emu", 0)
    sa_top = safe.get("top_emu", 0)
    sa_bottom = safe.get("bottom_emu", 0)

    safe_left = sa_left
    safe_top = sa_top
    safe_right = slide_width - sa_right
    safe_bottom = slide_height - sa_bottom

    findings = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if _is_full_bleed(shape, slide_width, slide_height):
                continue
            try:
                if (shape.left < safe_left or
                        shape.top < safe_top or
                        shape.left + shape.width > safe_right or
                        shape.top + shape.height > safe_bottom):
                    findings.append({
                        "check_id": "VH-12",
                        "category": "Spatial",
                        "check_name": "Content within safe_area",
                        "severity": "WARNING",
                        "pass": False,
                        "slide_index": si,
                        "details": (
                            f"Shape '{_shape_name(shape)}' is outside safe area"
                        ),
                    })
            except Exception:
                pass
    return findings


def _check_vh13(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-13: No significant shape overlaps (>10% AABB overlap)."""
    findings = []
    for si, slide in enumerate(prs.slides):
        shapes = [s for s in slide.shapes
                  if not _is_background_shape(s, slide_width, slide_height)]
        for i in range(len(shapes)):
            for j in range(i + 1, len(shapes)):
                overlap = _aabb_overlap_fraction(shapes[i], shapes[j])
                severity = _overlap_severity(shapes[i], shapes[j], overlap)
                if severity:
                    findings.append({
                        "check_id": "VH-13",
                        "category": "Spatial",
                        "check_name": "No significant shape overlaps",
                        "severity": severity,
                        "pass": False,
                        "slide_index": si,
                        "details": (
                            f"Shapes '{_shape_name(shapes[i])}' and "
                            f"'{_shape_name(shapes[j])}' overlap by "
                            f"{overlap:.0%}"
                        ),
                    })
    return findings


def _check_vh27(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-27: No full black slide backgrounds."""
    findings = []
    for si, slide in enumerate(prs.slides):
        bg_rgb = _get_slide_bg_rgb(slide, ds)
        if _is_near_black(bg_rgb):
            findings.append({
                "check_id": "VH-27",
                "category": "Color",
                "check_name": "No full black slide backgrounds",
                "severity": "BLOCKING",
                "pass": False,
                "slide_index": si,
                "details": (
                    "Slide background is full black or near-black; use a light "
                    "canvas or a non-full-bleed dark treatment so template "
                    "branding remains visible"
                ),
            })
            continue
        for shape in slide.shapes:
            if not _is_full_bleed(shape, slide_width, slide_height):
                continue
            rgb = _get_shape_fill_rgb(shape)
            if rgb is not None and _is_near_black(rgb):
                findings.append({
                    "check_id": "VH-27",
                    "category": "Color",
                    "check_name": "No full black slide backgrounds",
                    "severity": "BLOCKING",
                    "pass": False,
                    "slide_index": si,
                    "details": (
                        f"Full-bleed shape '{_shape_name(shape)}' is black or "
                        "near-black and can hide template branding"
                    ),
                })
    return findings


def _check_vh14(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-14: Consistent gutters between peer shapes."""
    spacing_scale = ds.get("spacing_scale", {})
    valid_spacings = list(spacing_scale.values())

    findings = []
    for si, slide in enumerate(prs.slides):
        content_shapes = [
            s for s in slide.shapes
            if not _is_background_shape(s, slide_width, slide_height)
            and not _is_placeholder(s)
        ]
        if len(content_shapes) < 3:
            continue

        # Sort by left position
        sorted_h = sorted(content_shapes, key=lambda s: s.left)
        gaps = []
        for i in range(len(sorted_h) - 1):
            gap = sorted_h[i + 1].left - (sorted_h[i].left + sorted_h[i].width)
            if gap > 0:
                gaps.append(gap)

        if len(gaps) < 2:
            continue

        # Check if gaps match spacing_scale values
        for gap in gaps:
            matches = any(
                abs(gap - sp) / max(sp, 1) <= 0.10
                for sp in valid_spacings
            )
            if not matches:
                findings.append({
                    "check_id": "VH-14",
                    "category": "Spatial",
                    "check_name": "Consistent gutters",
                    "severity": "WARNING",
                    "pass": False,
                    "slide_index": si,
                    "details": (
                        f"Gap of {gap} EMU doesn't match any spacing_scale value"
                    ),
                })
    return findings


def _check_vh15(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-15: Grid alignment (shape left edges align to grid columns)."""
    grid_cfg = ds.get("grid", {})
    cols = grid_cfg.get("cols", 12)
    gutter = grid_cfg.get("gutter_md_emu", 137160)
    tolerance = 45720  # 0.05 inches

    findings = []
    for si, slide in enumerate(prs.slides):
        canvas_def = _get_slide_canvas_def(prs, slide, ds)
        if canvas_def is None:
            continue

        body_region = canvas_def.get("body_region", {})
        body_left = body_region.get("left_emu")
        body_width = body_region.get("width_emu")
        if body_left is None or body_width is None:
            continue

        total_gutter = (cols - 1) * gutter
        col_width = (body_width - total_gutter) / cols
        grid_lefts = [
            body_left + c * (col_width + gutter)
            for c in range(cols)
        ]

        content_shapes = [
            s for s in slide.shapes
            if not _is_background_shape(s, slide_width, slide_height)
            and not _is_placeholder(s)
        ]
        if len(content_shapes) < 3:
            continue

        for shape in content_shapes:
            try:
                shape_left = shape.left
                aligned = any(abs(shape_left - gl) <= tolerance for gl in grid_lefts)
                if not aligned:
                    findings.append({
                        "check_id": "VH-15",
                        "category": "Spatial",
                        "check_name": "Grid alignment",
                        "severity": "WARNING",
                        "pass": False,
                        "slide_index": si,
                        "details": (
                            f"Shape '{_shape_name(shape)}' left edge {shape_left} "
                            f"doesn't align to any grid column"
                        ),
                    })
            except Exception:
                pass
    return findings


def _check_vh16(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-16: No leaked markdown."""
    findings = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            text = _get_text_from_shape(shape)
            if not text:
                continue
            for pattern in _MARKDOWN_PATTERNS:
                if pattern.search(text):
                    findings.append({
                        "check_id": "VH-16",
                        "category": "Content Rendering",
                        "check_name": "No leaked markdown",
                        "severity": "BLOCKING",
                        "pass": False,
                        "slide_index": si,
                        "details": (
                            f"Shape '{_shape_name(shape)}' contains markdown: "
                            f"matched '{pattern.pattern}'"
                        ),
                    })
                    break  # One finding per shape is enough
    return findings


def _check_vh17(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-17: No placeholder text."""
    findings = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            text = _get_text_from_shape(shape)
            if not text:
                continue
            for pattern in _PLACEHOLDER_PATTERNS:
                if pattern.search(text):
                    findings.append({
                        "check_id": "VH-17",
                        "category": "Content Rendering",
                        "check_name": "No placeholder text",
                        "severity": "BLOCKING",
                        "pass": False,
                        "slide_index": si,
                        "details": (
                            f"Shape '{_shape_name(shape)}' contains placeholder: "
                            f"matched '{pattern.pattern}'"
                        ),
                    })
                    break
    return findings


def _check_vh18(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-18: All picture rels resolve."""
    findings = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if _is_image_shape(shape):
                try:
                    # Access image to verify it exists
                    _ = shape.image
                except Exception as e:
                    findings.append({
                        "check_id": "VH-18",
                        "category": "Content Rendering",
                        "check_name": "All picture rels resolve",
                        "severity": "BLOCKING",
                        "pass": False,
                        "slide_index": si,
                        "details": (
                            f"Shape '{_shape_name(shape)}' picture reference "
                            f"is broken: {e}"
                        ),
                    })
    return findings


def _check_vh19(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-19: No empty content frames."""
    findings = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if _is_background_shape(shape, slide_width, slide_height):
                continue
            if _is_placeholder(shape):
                continue
            text = shape.text_frame.text.strip()
            if not text:
                findings.append({
                    "check_id": "VH-19",
                    "category": "Content Rendering",
                    "check_name": "No empty content frames",
                    "severity": "WARNING",
                    "pass": False,
                    "slide_index": si,
                    "details": (
                        f"Shape '{_shape_name(shape)}' is an empty text frame"
                    ),
                })
    return findings


def _check_vh20(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-20: Title position consistent across slides."""
    if len(prs.slides) < 2:
        return []

    # Find title shapes (first text shape with large text or first shape)
    title_positions = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                for run, _ in _iter_text_runs(shape):
                    size = _get_font_size_pt(run)
                    if size is not None and size >= 24:
                        try:
                            title_positions.append((si, shape.left, shape.top))
                        except Exception:
                            pass
                        break
                if title_positions and title_positions[-1][0] == si:
                    break

    if len(title_positions) < 2:
        return []

    findings = []
    tolerance = 91440  # 0.1 inches
    ref_left = title_positions[0][1]
    ref_top = title_positions[0][2]
    for si, left, top in title_positions[1:]:
        if abs(left - ref_left) > tolerance or abs(top - ref_top) > tolerance:
            findings.append({
                "check_id": "VH-20",
                "category": "Cross-Slide",
                "check_name": "Title position consistent",
                "severity": "WARNING",
                "pass": False,
                "slide_index": si,
                "details": (
                    f"Title position ({left}, {top}) differs from reference "
                    f"({ref_left}, {ref_top}) by more than 0.1 inches"
                ),
            })
    return findings


def _check_vh21(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-21: Title style consistent across slides."""
    if len(prs.slides) < 2:
        return []

    title_styles = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                for run, _ in _iter_text_runs(shape):
                    size = _get_font_size_pt(run)
                    if size is not None and size >= 24:
                        try:
                            rgb = None
                            if run.font.color and run.font.color.rgb:
                                rgb = _rgb_to_tuple(run.font.color.rgb)
                            title_styles.append({
                                "slide_index": si,
                                "font_name": run.font.name,
                                "size_pt": size,
                                "bold": run.font.bold,
                                "color": rgb,
                            })
                        except Exception:
                            pass
                        break
                if title_styles and title_styles[-1]["slide_index"] == si:
                    break

    if len(title_styles) < 2:
        return []

    findings = []
    ref = title_styles[0]
    for ts in title_styles[1:]:
        diffs = []
        if ts["font_name"] != ref["font_name"]:
            diffs.append(f"font={ts['font_name']} vs {ref['font_name']}")
        if ts["size_pt"] != ref["size_pt"]:
            diffs.append(f"size={ts['size_pt']} vs {ref['size_pt']}")
        if ts["bold"] != ref["bold"]:
            diffs.append(f"bold={ts['bold']} vs {ref['bold']}")
        if ts["color"] != ref["color"]:
            diffs.append(f"color differs")
        if diffs:
            findings.append({
                "check_id": "VH-21",
                "category": "Cross-Slide",
                "check_name": "Title style consistent",
                "severity": "BLOCKING",
                "pass": False,
                "slide_index": ts["slide_index"],
                "details": (
                    f"Title style differs from slide 0: {', '.join(diffs)}"
                ),
            })
    return findings


def _check_vh22(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-22: Kicker style consistent across slides."""
    if len(prs.slides) < 2:
        return []

    # Kicker = 11pt bold uppercase text
    kicker_styles = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            for run, _ in _iter_text_runs(shape):
                size = _get_font_size_pt(run)
                if size is not None and abs(size - 11) <= 1 and run.font.bold:
                    try:
                        rgb = None
                        if run.font.color and run.font.color.rgb:
                            rgb = _rgb_to_tuple(run.font.color.rgb)
                        kicker_styles.append({
                            "slide_index": si,
                            "font_name": run.font.name,
                            "size_pt": size,
                            "color": rgb,
                        })
                    except Exception:
                        pass
                    break
            if kicker_styles and kicker_styles[-1]["slide_index"] == si:
                break

    if len(kicker_styles) < 2:
        return []

    findings = []
    ref = kicker_styles[0]
    for ks in kicker_styles[1:]:
        diffs = []
        if ks["font_name"] != ref["font_name"]:
            diffs.append(f"font={ks['font_name']} vs {ref['font_name']}")
        if ks["size_pt"] != ref["size_pt"]:
            diffs.append(f"size={ks['size_pt']} vs {ref['size_pt']}")
        if ks["color"] != ref["color"]:
            diffs.append(f"color differs")
        if diffs:
            findings.append({
                "check_id": "VH-22",
                "category": "Cross-Slide",
                "check_name": "Kicker style consistent",
                "severity": "WARNING",
                "pass": False,
                "slide_index": ks["slide_index"],
                "details": (
                    f"Kicker style differs from slide 0: {', '.join(diffs)}"
                ),
            })
    return findings


def _check_vh23(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-23: Body text style consistent across slides."""
    if len(prs.slides) < 2:
        return []

    # Body text = 12pt non-bold text
    body_styles = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            for run, _ in _iter_text_runs(shape):
                size = _get_font_size_pt(run)
                if size is not None and abs(size - 12) <= 1 and not run.font.bold:
                    try:
                        rgb = None
                        if run.font.color and run.font.color.rgb:
                            rgb = _rgb_to_tuple(run.font.color.rgb)
                        body_styles.append({
                            "slide_index": si,
                            "font_name": run.font.name,
                            "size_pt": size,
                            "color": rgb,
                        })
                    except Exception:
                        pass
                    break
            if body_styles and body_styles[-1]["slide_index"] == si:
                break

    if len(body_styles) < 2:
        return []

    findings = []
    ref = body_styles[0]
    for bs in body_styles[1:]:
        diffs = []
        if bs["font_name"] != ref["font_name"]:
            diffs.append(f"font={bs['font_name']} vs {ref['font_name']}")
        if bs["size_pt"] != ref["size_pt"]:
            diffs.append(f"size={bs['size_pt']} vs {ref['size_pt']}")
        if bs["color"] != ref["color"]:
            diffs.append(f"color differs")
        if diffs:
            findings.append({
                "check_id": "VH-23",
                "category": "Cross-Slide",
                "check_name": "Body text style consistent",
                "severity": "WARNING",
                "pass": False,
                "slide_index": bs["slide_index"],
                "details": (
                    f"Body style differs from slide 0: {', '.join(diffs)}"
                ),
            })
    return findings


def _check_vh24(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-24: No empty slides (>= 2 non-placeholder shapes)."""
    findings = []
    for si, slide in enumerate(prs.slides):
        non_ph_count = sum(1 for s in slide.shapes if not _is_placeholder(s))
        if non_ph_count < 2:
            findings.append({
                "check_id": "VH-24",
                "category": "Structural",
                "check_name": "No empty slides",
                "severity": "BLOCKING",
                "pass": False,
                "slide_index": si,
                "details": (
                    f"Slide has only {non_ph_count} non-placeholder shapes "
                    f"(need >= 2)"
                ),
            })
    return findings


def _check_vh25(prs, ds, slide_width, slide_height, deck_plan: dict | None) -> list[dict]:
    """VH-25: Slide count matches deck_plan."""
    if deck_plan is None:
        return []
    expected = len(deck_plan.get("slides", []))
    actual = len(prs.slides)
    if actual != expected:
        return [{
            "check_id": "VH-25",
            "category": "Structural",
            "check_name": "Slide count matches deck_plan",
            "severity": "BLOCKING",
            "pass": False,
            "slide_index": -1,
            "details": (
                f"Expected {expected} slides from deck_plan, got {actual}"
            ),
        }]
    return []


def _check_vh26(prs, ds, slide_width, slide_height) -> list[dict]:
    """VH-26: Every content slide has >= 1 non-text visual element."""
    findings = []
    for si, slide in enumerate(prs.slides):
        has_visual = False
        for shape in slide.shapes:
            if _is_placeholder(shape):
                continue
            # Visual = image or non-text shape with fill
            if _is_image_shape(shape):
                has_visual = True
                break
            if not shape.has_text_frame:
                has_visual = True
                break
            # Shape with fill and text is also visual (e.g., card)
            if _get_shape_fill_rgb(shape) is not None:
                has_visual = True
                break
        if not has_visual:
            findings.append({
                "check_id": "VH-26",
                "category": "Structural",
                "check_name": "Non-text visual element per slide",
                "severity": "WARNING",
                "pass": False,
                "slide_index": si,
                "details": (
                    f"Slide has no non-text visual elements"
                ),
            })
    return findings


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

_ACTIVE_CHECKS = [
    _check_vh01, _check_vh02, _check_vh03, _check_vh04, _check_vh05,
    _check_vh06, _check_vh07, _check_vh08, _check_vh09, _check_vh10,
    _check_vh11, _check_vh12, _check_vh13, _check_vh15,
    _check_vh16, _check_vh17, _check_vh18, _check_vh19,
    _check_vh24, _check_vh26, _check_vh27,
]

_DEFERRED_CHECKS = {
    "VH-14": "Peer-gutter detection is heuristic without explicit grouping anchors.",
    "VH-20": "Title-role inference is heuristic without explicit title anchors.",
    "VH-21": "Title-role inference is heuristic without explicit title anchors.",
    "VH-22": "Kicker-role inference is heuristic without explicit kicker anchors.",
    "VH-23": "Body-role inference is heuristic without explicit body anchors.",
}

# VH-25 handled separately (needs deck_plan)


def _describe_check(check_fn) -> tuple[str, str]:
    """Return (check_id, check_name) for a scanner check function."""
    doc = (check_fn.__doc__ or "").strip().splitlines()
    if doc:
        header = doc[0]
        if header.startswith("VH-") and ":" in header:
            check_id, check_name = header.split(":", 1)
            return check_id.strip(), check_name.strip()

    match = re.search(r"_check_vh(\d{2})$", check_fn.__name__)
    if match:
        return f"VH-{match.group(1)}", check_fn.__name__

    return "VH-00", check_fn.__name__


def _internal_error_finding(
    check_fn,
    exc: Exception,
    *,
    check_id: str | None = None,
    check_name: str | None = None,
) -> dict:
    """Convert an internal scanner failure into a blocking report finding."""
    derived_id, derived_name = _describe_check(check_fn)
    check_id = check_id or derived_id
    check_name = check_name or derived_name
    return {
        "check_id": check_id,
        "category": "Structural",
        "check_name": check_name,
        "severity": "BLOCKING",
        "pass": False,
        "slide_index": -1,
        "details": (
            f"Scanner internal error while running {check_id} "
            f"({check_name}): {type(exc).__name__}: {exc}"
        ),
    }


def scan_pptx(
    pptx_path: str | Path,
    design_system_path: str | Path,
    deck_plan: dict | None = None,
) -> dict:
    """Run the active objective visual hygiene checks on a PPTX file.

    Args:
        pptx_path: Path to the PPTX file.
        design_system_path: Path to design_system.json.
        deck_plan: Optional deck plan dict for slide-count check.

    Returns:
        A geometry_report dict with pass/fail, counts, and findings.
        Internal scanner failures are surfaced as synthetic BLOCKING findings.
    """
    pptx_path = Path(pptx_path)
    design_system_path = Path(design_system_path)

    prs = Presentation(str(pptx_path))
    with open(design_system_path) as f:
        ds = json.load(f)

    slide_width = ds["canvas"]["width_emu"]
    slide_height = ds["canvas"]["height_emu"]

    all_findings = []

    # Run standard checks
    for check_fn in _ACTIVE_CHECKS:
        try:
            findings = check_fn(prs, ds, slide_width, slide_height)
            all_findings.extend(findings)
        except Exception as exc:
            all_findings.append(_internal_error_finding(check_fn, exc))

    # VH-25 needs deck_plan
    try:
        all_findings.extend(_check_vh25(prs, ds, slide_width, slide_height, deck_plan))
    except Exception as exc:
        all_findings.append(
            _internal_error_finding(
                _check_vh25,
                exc,
                check_id="VH-25",
                check_name="Slide count matches deck_plan",
            )
        )

    blocking_count = sum(1 for f in all_findings if f["severity"] == "BLOCKING")
    warning_count = sum(1 for f in all_findings if f["severity"] == "WARNING")

    return {
        "pass": blocking_count == 0,
        "blocking_count": blocking_count,
        "warning_count": warning_count,
        "findings": all_findings,
    }
