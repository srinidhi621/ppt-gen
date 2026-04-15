"""Low-level shape helpers that wrap python-pptx primitives.

Every helper takes a *slide* and a :class:`Rect` (or explicit EMU
coordinates) and returns the created shape.  Fill colours come from
``tokens.color()``; font properties come from ``tokens.type()``.
"""

from __future__ import annotations

from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt


_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}

_ANCHOR_MAP = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}


def add_rect(
    slide,
    rect,
    fill: RGBColor,
    line: RGBColor | None = None,
    line_width_pt: float = 0.75,
):
    """Add a filled rectangle.  Returns the shape."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(rect.left), Emu(rect.top), Emu(rect.width), Emu(rect.height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_width_pt)
    shp.shadow.inherit = False
    return shp


def add_text(
    slide,
    rect,
    text: str,
    *,
    type_style: dict | None = None,
    font_name: str | None = None,
    font_size_pt: int | None = None,
    bold: bool | None = None,
    color: RGBColor | None = None,
    fill: RGBColor | None = None,
    align: str = "left",
    anchor: str = "top",
    word_wrap: bool = True,
    upper: bool | None = None,
):
    """Add a text box.  Returns the textbox shape.

    Accepts either a *type_style* dict (from ``tokens.type()``) or
    explicit *font_name* / *font_size_pt* / *bold*.  Explicit values
    override the style dict.  *fill* sets the text-box background.
    """
    tb = slide.shapes.add_textbox(
        Emu(rect.left), Emu(rect.top), Emu(rect.width), Emu(rect.height),
    )
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = _ANCHOR_MAP.get(anchor, MSO_ANCHOR.TOP)

    p = tf.paragraphs[0]
    p.alignment = _ALIGN_MAP.get(align, PP_ALIGN.LEFT)

    # Resolve font properties
    _name = font_name
    _size = font_size_pt
    _bold = bold
    if type_style:
        _name = _name or type_style.get("font")
        _size = _size if _size is not None else type_style.get("size_pt")
        _bold = _bold if _bold is not None else type_style.get("bold", False)

    # Apply upper-case if style requests it
    display_text = text
    should_upper = upper
    if should_upper is None and type_style:
        should_upper = type_style.get("upper")
    if should_upper:
        display_text = text.upper()

    r = p.add_run()
    r.text = display_text
    if _name:
        r.font.name = _name
    if _size is not None:
        r.font.size = Pt(_size)
    if _bold is not None:
        r.font.bold = _bold
    if color is not None:
        r.font.color.rgb = color

    return tb


def add_image(slide, rect, path: str | Path):
    """Add a picture.  Returns the picture shape."""
    return slide.shapes.add_picture(
        str(path),
        Emu(rect.left), Emu(rect.top), Emu(rect.width), Emu(rect.height),
    )


def add_line(
    slide,
    start_x: int,
    start_y: int,
    end_x: int,
    end_y: int,
    color: RGBColor,
    width_pt: float = 0.75,
):
    """Add a straight line connector.  Returns the connector shape."""
    conn = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Emu(start_x), Emu(start_y),
        Emu(end_x), Emu(end_y),
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    return conn
