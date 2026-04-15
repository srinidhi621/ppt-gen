"""Template loading and canvas management.

A Canvas wraps a python-pptx Presentation and the design-system
definition. It tracks the current canvas (layout) so that body-region
properties reflect the most recently added slide.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation

from .errors import CanvasNotFoundError


class Canvas:
    """Wrapper around a Presentation + design system.

    Usage::

        canvas = load_template("assets/template/template.pptx")
        slide  = canvas.add_slide("header_light")
        # canvas.body_left / body_top / body_width / body_height
        # now reflect the header_light body region.
        canvas.save("output.pptx")
    """

    def __init__(self, prs: Presentation, design_system: dict) -> None:
        self._prs = prs
        self._ds = design_system
        self._current_canvas_def: dict | None = None

    # -- public properties --------------------------------------------------

    @property
    def design_system(self) -> dict:
        return self._ds

    @property
    def presentation(self) -> Presentation:
        return self._prs

    @property
    def slide_width(self) -> int:
        return self._ds["canvas"]["width_emu"]

    @property
    def slide_height(self) -> int:
        return self._ds["canvas"]["height_emu"]

    @property
    def body_left(self) -> int:
        self._require_canvas()
        return self._current_canvas_def["body_region"]["left_emu"]

    @property
    def body_top(self) -> int:
        self._require_canvas()
        return self._current_canvas_def["body_region"]["top_emu"]

    @property
    def body_width(self) -> int:
        self._require_canvas()
        return self._current_canvas_def["body_region"]["width_emu"]

    @property
    def body_height(self) -> int:
        self._require_canvas()
        return self._current_canvas_def["body_region"]["height_emu"]

    # -- slide management ---------------------------------------------------

    def add_slide(self, canvas_name: str = "header_light"):
        """Add a slide using the named canvas and return the pptx Slide.

        Updates body-region properties to reflect the chosen canvas.
        """
        canvas_def = self._ds["canvases"].get(canvas_name)
        if canvas_def is None:
            available = ", ".join(sorted(self._ds["canvases"]))
            raise CanvasNotFoundError(
                f"Canvas '{canvas_name}' not found. Available: {available}"
            )

        self._current_canvas_def = canvas_def
        layout_index = canvas_def["layout_index"]
        layout = self._prs.slide_layouts[layout_index]
        slide = self._prs.slides.add_slide(layout)
        _remove_placeholders(slide)
        return slide

    def save(self, output_path: str | Path) -> None:
        self._prs.save(str(output_path))

    # -- internals ----------------------------------------------------------

    def _require_canvas(self) -> None:
        if self._current_canvas_def is None:
            raise CanvasNotFoundError(
                "No canvas set. Call add_slide() before accessing body properties."
            )


def load_template(
    template_path: str | Path,
    design_system_path: str | Path | None = None,
    *,
    keep_template_slides: bool = False,
) -> Canvas:
    """Load a branded template and its design system.

    If *design_system_path* is not given, looks for ``design_system.json``
    in the same directory as the template.
    """
    template_path = Path(template_path)
    if design_system_path is None:
        design_system_path = template_path.parent / "design_system.json"
    else:
        design_system_path = Path(design_system_path)

    prs = Presentation(str(template_path))
    if not keep_template_slides:
        _drop_existing_slides(prs)

    with open(design_system_path, encoding="utf-8") as f:
        ds = json.load(f)

    return Canvas(prs, ds)


def _drop_existing_slides(prs: Presentation) -> None:
    """Start from the template's theme/layouts, not its sample slides."""
    for slide_id in list(prs.slides._sldIdLst):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)


def _remove_placeholders(slide) -> None:
    """Generated slides should not inherit empty layout placeholders."""
    for shape in list(slide.shapes):
        if not shape.is_placeholder:
            continue
        element = shape._element
        element.getparent().remove(element)
