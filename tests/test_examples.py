"""Tests for the example library.

Discovers all examples in examples/ and verifies:
- build.py exists and has valid Python syntax
- metadata.json exists and has required fields
- Running build.py produces a valid PPTX
- The PPTX passes the objective scanner (no BLOCKING findings)
- Shape count stays within declared density bounds (min and max)
- The runner script discovers examples correctly
"""

import json
import subprocess
import sys
from pathlib import Path

import pytest

EXAMPLES_DIR = Path(__file__).resolve().parent.parent / "examples"
PROJECT_ROOT = Path(__file__).resolve().parent.parent
DESIGN_SYSTEM = PROJECT_ROOT / "assets" / "template" / "design_system.json"

# Required fields in metadata.json
REQUIRED_METADATA_FIELDS = {
    "archetype",
    "source",
    "designer",
    "intent",
    "invariants",
    "variables",
    "canvas",
    "grid",
    "density",
}

REQUIRED_DENSITY_FIELDS = {"min_shapes", "max_shapes", "text_density"}
REQUIRED_GRID_FIELDS = {"cols", "gutter"}


def discover_example_dirs():
    """Find all example directories (those containing build.py)."""
    build_files = sorted(EXAMPLES_DIR.glob("*/*/build.py"))
    return [
        f.parent for f in build_files
        if "__pycache__" not in str(f.relative_to(EXAMPLES_DIR))
    ]


EXAMPLE_DIRS = discover_example_dirs()
EXAMPLE_IDS = [
    str(d.relative_to(EXAMPLES_DIR)) for d in EXAMPLE_DIRS
]


# ---------------------------------------------------------------------------
# Parametrized tests per example
# ---------------------------------------------------------------------------

@pytest.fixture(params=EXAMPLE_DIRS, ids=EXAMPLE_IDS)
def example_dir(request):
    return request.param


class TestExampleStructure:
    """Verify structural correctness of each example."""

    def test_build_py_exists(self, example_dir):
        build_py = example_dir / "build.py"
        assert build_py.exists(), f"build.py missing in {example_dir}"

    def test_build_py_valid_syntax(self, example_dir):
        build_py = example_dir / "build.py"
        source = build_py.read_text()
        try:
            compile(source, str(build_py), "exec")
        except SyntaxError as e:
            pytest.fail(f"Syntax error in {build_py}: {e}")

    def test_build_py_no_inline_fonts(self, example_dir):
        """Build files must use tokens.type(...), not inline font_name/font_size_pt."""
        build_py = example_dir / "build.py"
        source = build_py.read_text()
        for line_num, line in enumerate(source.splitlines(), 1):
            # Skip comments
            stripped = line.lstrip()
            if stripped.startswith("#"):
                continue
            assert "font_name=" not in line, (
                f"{build_py}:{line_num} uses inline font_name=. "
                f"Use tokens.type(...) instead."
            )
            assert "font_size_pt=" not in line, (
                f"{build_py}:{line_num} uses inline font_size_pt=. "
                f"Use tokens.type(...) instead."
            )

    def test_metadata_json_exists(self, example_dir):
        meta = example_dir / "metadata.json"
        assert meta.exists(), f"metadata.json missing in {example_dir}"

    def test_metadata_json_valid(self, example_dir):
        meta = example_dir / "metadata.json"
        try:
            data = json.loads(meta.read_text())
        except json.JSONDecodeError as e:
            pytest.fail(f"Invalid JSON in {meta}: {e}")

        missing = REQUIRED_METADATA_FIELDS - set(data.keys())
        assert not missing, f"Missing metadata fields: {missing}"

        # Check sub-objects
        assert isinstance(data["invariants"], list), "invariants must be a list"
        assert len(data["invariants"]) >= 1, "invariants must have at least 1 entry"
        assert isinstance(data["variables"], list), "variables must be a list"
        assert len(data["variables"]) >= 1, "variables must have at least 1 entry"

        grid = data["grid"]
        missing_grid = REQUIRED_GRID_FIELDS - set(grid.keys())
        assert not missing_grid, f"Missing grid fields: {missing_grid}"

        density = data["density"]
        missing_density = REQUIRED_DENSITY_FIELDS - set(density.keys())
        assert not missing_density, f"Missing density fields: {missing_density}"
        assert density["min_shapes"] > 0, "min_shapes must be > 0"
        assert density["max_shapes"] >= density["min_shapes"], (
            "max_shapes must be >= min_shapes"
        )

    def test_metadata_archetype_matches_directory(self, example_dir):
        """The archetype in metadata should match the parent directory name."""
        meta = example_dir / "metadata.json"
        if not meta.exists():
            pytest.skip("no metadata.json")
        data = json.loads(meta.read_text())
        archetype_dir = example_dir.parent.name
        assert data["archetype"] == archetype_dir, (
            f"archetype '{data['archetype']}' does not match "
            f"directory name '{archetype_dir}'"
        )


class TestExampleExecution:
    """Verify each example actually runs and produces a valid PPTX."""

    def test_build_produces_pptx(self, example_dir):
        """Run build.py and check that output.pptx is created."""
        build_py = example_dir / "build.py"
        output_pptx = example_dir / "output.pptx"

        # Remove existing output to ensure a fresh build
        if output_pptx.exists():
            output_pptx.unlink()

        result = subprocess.run(
            [sys.executable, str(build_py)],
            cwd=str(PROJECT_ROOT),
            env={**__import__("os").environ, "PYTHONPATH": str(PROJECT_ROOT)},
            capture_output=True,
            text=True,
            timeout=60,
        )

        assert result.returncode == 0, (
            f"build.py failed:\nstdout: {result.stdout}\nstderr: {result.stderr}"
        )
        assert output_pptx.exists(), "output.pptx was not created"
        assert output_pptx.stat().st_size > 0, "output.pptx is empty"

    def test_pptx_has_at_least_one_slide(self, example_dir):
        """The output PPTX should have at least one slide."""
        output_pptx = example_dir / "output.pptx"
        if not output_pptx.exists():
            pytest.skip("output.pptx not found; run test_build_produces_pptx first")

        from pptx import Presentation
        prs = Presentation(str(output_pptx))
        assert len(prs.slides) >= 1, (
            f"Expected at least 1 slide, got {len(prs.slides)}"
        )

    def test_pptx_passes_objective_scanner(self, example_dir):
        """The output PPTX must pass the active objective scanner with zero BLOCKING."""
        output_pptx = example_dir / "output.pptx"
        if not output_pptx.exists():
            pytest.skip("output.pptx not found")

        from src.scan.scanner import scan_pptx

        report = scan_pptx(str(output_pptx), str(DESIGN_SYSTEM))
        blocking = [
            f for f in report["findings"] if f["severity"] == "BLOCKING"
        ]
        assert report["pass"], (
            f"Scanner found {len(blocking)} BLOCKING finding(s):\n"
            + "\n".join(
                f"  {f['check_id']}: {f['details']}" for f in blocking
            )
        )

    def test_pptx_within_density_bounds(self, example_dir):
        """Shape count on the built slide must be within [min_shapes, max_shapes]."""
        output_pptx = example_dir / "output.pptx"
        meta_path = example_dir / "metadata.json"

        if not output_pptx.exists():
            pytest.skip("output.pptx not found")
        if not meta_path.exists():
            pytest.skip("metadata.json not found")

        from pptx import Presentation
        prs = Presentation(str(output_pptx))
        last_slide = prs.slides[-1]
        shape_count = len(last_slide.shapes)

        meta = json.loads(meta_path.read_text())
        min_shapes = meta["density"]["min_shapes"]
        max_shapes = meta["density"]["max_shapes"]

        assert shape_count >= min_shapes, (
            f"Last slide has {shape_count} shapes, "
            f"metadata requires at least {min_shapes}"
        )
        assert shape_count <= max_shapes, (
            f"Last slide has {shape_count} shapes, "
            f"metadata allows at most {max_shapes}"
        )


class TestRunnerScript:
    """Verify the run_all.py runner script."""

    def test_runner_script_exists(self):
        runner = EXAMPLES_DIR / "run_all.py"
        assert runner.exists(), "run_all.py missing from examples/"

    def test_runner_script_valid_syntax(self):
        runner = EXAMPLES_DIR / "run_all.py"
        source = runner.read_text()
        try:
            compile(source, str(runner), "exec")
        except SyntaxError as e:
            pytest.fail(f"Syntax error in run_all.py: {e}")

    def test_runner_discovers_examples(self):
        """The runner's discover function should find all examples."""
        sys.path.insert(0, str(EXAMPLES_DIR))
        try:
            from run_all import discover_examples
            examples = discover_examples()
            assert len(examples) >= 6, (
                f"Runner found only {len(examples)} examples, expected at least 6"
            )
        finally:
            sys.path.pop(0)


class TestExtractScript:
    """Verify the extraction script exists and has valid syntax."""

    def test_extract_script_exists(self):
        script = EXAMPLES_DIR / "extract_designer_slides.py"
        assert script.exists(), "extract_designer_slides.py missing from examples/"

    def test_extract_script_valid_syntax(self):
        script = EXAMPLES_DIR / "extract_designer_slides.py"
        source = script.read_text()
        try:
            compile(source, str(script), "exec")
        except SyntaxError as e:
            pytest.fail(f"Syntax error in extract_designer_slides.py: {e}")
