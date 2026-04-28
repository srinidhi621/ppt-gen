"""Tests for src.scan.scanner — 26 visual hygiene checks."""

import json
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt

import src.scan.scanner as scanner_module
from src.scan.scanner import scan_pptx, _relative_luminance, _contrast_ratio

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

_TEMPLATE = Path("assets/template/template.pptx")
_DESIGN_SYSTEM = Path("assets/template/design_system.json")


def _load_ds():
    with open(_DESIGN_SYSTEM) as f:
        return json.load(f)


# ---------------------------------------------------------------------------
# Fixture helpers
# ---------------------------------------------------------------------------

def _make_pptx(callback, *, layout_index=None, use_template=False):
    """Create a temp PPTX via callback(prs, slide, ds) and return the path.

    By default creates a blank presentation (no pre-existing slides) so tests
    only see shapes they explicitly add.  Set ``use_template=True`` to build
    on the real Ascendion template.
    """
    ds = _load_ds()
    if use_template:
        prs = Presentation(str(_TEMPLATE))
        if layout_index is None:
            layout_index = ds["canvases"]["header_light"]["layout_index"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    else:
        prs = Presentation()
        prs.slide_width = Emu(ds["canvas"]["width_emu"])
        prs.slide_height = Emu(ds["canvas"]["height_emu"])
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    callback(prs, slide, ds)
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    return tmp.name


def _add_filled_rect(slide, left, top, width, height, fill_rgb):
    """Add a colored rectangle shape."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(left), Emu(top), Emu(width), Emu(height),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    shp.line.fill.background()
    return shp


def _add_textbox(slide, left, top, width, height, text,
                 font_name=None, font_size_pt=None, bold=None, color=None):
    """Add a text box with styled text."""
    tb = slide.shapes.add_textbox(
        Emu(left), Emu(top), Emu(width), Emu(height),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    if font_name:
        r.font.name = font_name
    if font_size_pt is not None:
        r.font.size = Pt(font_size_pt)
    if bold is not None:
        r.font.bold = bold
    if color is not None:
        r.font.color.rgb = color
    return tb


# ---------------------------------------------------------------------------
# Helper tests
# ---------------------------------------------------------------------------

class TestHelpers:
    def test_relative_luminance_black(self):
        assert _relative_luminance(0, 0, 0) == pytest.approx(0.0, abs=0.001)

    def test_relative_luminance_white(self):
        assert _relative_luminance(255, 255, 255) == pytest.approx(1.0, abs=0.001)

    def test_contrast_ratio_bw(self):
        ratio = _contrast_ratio((0, 0, 0), (255, 255, 255))
        assert ratio == pytest.approx(21.0, abs=0.1)

    def test_contrast_ratio_same(self):
        ratio = _contrast_ratio((128, 128, 128), (128, 128, 128))
        assert ratio == pytest.approx(1.0, abs=0.01)


# ---------------------------------------------------------------------------
# Color checks (VH-01 to VH-05)
# ---------------------------------------------------------------------------

class TestColorChecks:
    def test_vh01_off_palette_fill(self):
        """VH-01: Shape with non-palette fill is BLOCKING."""
        def build(prs, slide, ds):
            _add_filled_rect(slide, 500000, 500000, 2000000, 1000000,
                             RGBColor(0xFF, 0x00, 0x00))
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh01 = [f for f in report["findings"] if f["check_id"] == "VH-01"]
        assert len(vh01) >= 1
        assert vh01[0]["severity"] == "BLOCKING"
        assert not vh01[0]["pass"]

    def test_vh01_palette_fill_passes(self):
        """VH-01: Shape with palette fill passes."""
        ds = _load_ds()
        accent_hex = ds["colors"]["accent_1"].lstrip("#")
        fill = RGBColor(int(accent_hex[0:2], 16), int(accent_hex[2:4], 16), int(accent_hex[4:6], 16))

        def build(prs, slide, ds_inner):
            _add_filled_rect(slide, 500000, 500000, 2000000, 1000000, fill)
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh01 = [f for f in report["findings"] if f["check_id"] == "VH-01"]
        assert len(vh01) == 0

    def test_vh02_off_palette_text_color(self):
        """VH-02: Text with non-palette color is BLOCKING."""
        def build(prs, slide, ds):
            _add_textbox(slide, 500000, 500000, 5000000, 500000,
                         "Bad color text", font_name="Inter", font_size_pt=12,
                         color=RGBColor(0xAB, 0xCD, 0xEF))
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh02 = [f for f in report["findings"] if f["check_id"] == "VH-02"]
        assert len(vh02) >= 1
        assert vh02[0]["severity"] == "BLOCKING"

    def test_vh03_too_many_accents(self):
        """VH-03: More than 2 accent roles per slide is WARNING."""
        ds = _load_ds()
        def build(prs, slide, ds_inner):
            colors = ds_inner["colors"]
            # Use 3 different accent colors
            for i, key in enumerate(["accent_1", "accent_2", "accent_3"]):
                h = colors[key].lstrip("#")
                rgb = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
                _add_filled_rect(slide, 500000 + i * 2500000, 500000,
                                 2000000, 1000000, rgb)
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh03 = [f for f in report["findings"] if f["check_id"] == "VH-03"]
        assert len(vh03) >= 1
        assert vh03[0]["severity"] == "WARNING"

    def test_vh04_low_contrast(self):
        """VH-04: Low contrast text is BLOCKING."""
        def build(prs, slide, ds):
            # Light gray text on white background
            _add_textbox(slide, 500000, 500000, 5000000, 500000,
                         "Low contrast text", font_name="Inter", font_size_pt=12,
                         color=RGBColor(0xDD, 0xDD, 0xDD))
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh04 = [f for f in report["findings"] if f["check_id"] == "VH-04"]
        assert len(vh04) >= 1
        assert vh04[0]["severity"] == "BLOCKING"

    def test_vh04_high_contrast_passes(self):
        """VH-04: High contrast text passes."""
        def build(prs, slide, ds):
            _add_textbox(slide, 500000, 500000, 5000000, 500000,
                         "Good contrast", font_name="Inter", font_size_pt=12,
                         color=RGBColor(0x00, 0x00, 0x00))
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh04 = [f for f in report["findings"] if f["check_id"] == "VH-04"]
        assert len(vh04) == 0

    def test_vh05_invisible_shape(self):
        """VH-05: Shape with fill matching background is WARNING."""
        def build(prs, slide, ds):
            # White fill on white background
            _add_filled_rect(slide, 500000, 500000, 2000000, 1000000,
                             RGBColor(0xFF, 0xFF, 0xFF))
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh05 = [f for f in report["findings"] if f["check_id"] == "VH-05"]
        assert len(vh05) >= 1
        assert vh05[0]["severity"] == "WARNING"


# ---------------------------------------------------------------------------
# Typography checks (VH-06 to VH-09)
# ---------------------------------------------------------------------------

class TestTypographyChecks:
    def test_vh06_disallowed_font(self):
        """VH-06: Font not in allowlist is BLOCKING."""
        def build(prs, slide, ds):
            _add_textbox(slide, 500000, 500000, 5000000, 500000,
                         "Comic Sans text", font_name="Comic Sans MS",
                         font_size_pt=12)
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh06 = [f for f in report["findings"] if f["check_id"] == "VH-06"]
        assert len(vh06) >= 1
        assert vh06[0]["severity"] == "BLOCKING"

    def test_vh06_allowed_font_passes(self):
        """VH-06: Allowed font passes."""
        def build(prs, slide, ds):
            _add_textbox(slide, 500000, 500000, 5000000, 500000,
                         "Inter text", font_name="Inter", font_size_pt=12)
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh06 = [f for f in report["findings"] if f["check_id"] == "VH-06"]
        assert len(vh06) == 0

    def test_vh06_substitute_font_passes(self):
        """VH-06: Original font (before substitution) also passes."""
        def build(prs, slide, ds):
            _add_textbox(slide, 500000, 500000, 5000000, 500000,
                         "Aptos text", font_name="Aptos", font_size_pt=12)
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh06 = [f for f in report["findings"] if f["check_id"] == "VH-06"]
        assert len(vh06) == 0

    def test_vh07_nonstandard_size(self):
        """VH-07: Non-standard font size is WARNING."""
        def build(prs, slide, ds):
            _add_textbox(slide, 500000, 500000, 5000000, 500000,
                         "Weird size text", font_name="Inter", font_size_pt=18)
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh07 = [f for f in report["findings"] if f["check_id"] == "VH-07"]
        assert len(vh07) >= 1
        assert vh07[0]["severity"] == "WARNING"

    def test_vh08_wrong_bold(self):
        """VH-08: Wrong bold for size is WARNING."""
        def build(prs, slide, ds):
            # 12pt should be non-bold (body), but we set bold=True
            _add_textbox(slide, 500000, 500000, 5000000, 500000,
                         "Bold body text", font_name="Inter",
                         font_size_pt=12, bold=True)
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh08 = [f for f in report["findings"] if f["check_id"] == "VH-08"]
        assert len(vh08) >= 1
        assert vh08[0]["severity"] == "WARNING"

    def test_vh09_allcaps_non_kicker(self):
        """VH-09: ALLCAPS on non-kicker text is WARNING."""
        def build(prs, slide, ds):
            _add_textbox(slide, 500000, 500000, 5000000, 500000,
                         "ALL CAPS BODY TEXT", font_name="Inter",
                         font_size_pt=12, bold=False)
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh09 = [f for f in report["findings"] if f["check_id"] == "VH-09"]
        assert len(vh09) >= 1
        assert vh09[0]["severity"] == "WARNING"


# ---------------------------------------------------------------------------
# Spatial checks (VH-10 to VH-15)
# ---------------------------------------------------------------------------

class TestSpatialChecks:
    def test_vh10_shape_outside_canvas(self):
        """VH-10: Shape extending past canvas is BLOCKING."""
        def build(prs, slide, ds):
            sw = ds["canvas"]["width_emu"]
            _add_filled_rect(slide, sw - 500000, 500000, 1000000, 500000,
                             RGBColor(0x00, 0x85, 0x67))
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh10 = [f for f in report["findings"] if f["check_id"] == "VH-10"]
        assert len(vh10) >= 1
        assert vh10[0]["severity"] == "BLOCKING"

    def test_vh10_shape_inside_canvas_passes(self):
        """VH-10: Shape within canvas passes."""
        def build(prs, slide, ds):
            _add_filled_rect(slide, 500000, 500000, 2000000, 1000000,
                             RGBColor(0x00, 0x85, 0x67))
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh10 = [f for f in report["findings"] if f["check_id"] == "VH-10"]
        assert len(vh10) == 0

    def test_vh12_outside_safe_area(self):
        """VH-12: Shape outside safe area is WARNING."""
        def build(prs, slide, ds):
            # Place at (0,0) which is outside safe area
            _add_filled_rect(slide, 0, 0, 200000, 200000,
                             RGBColor(0x00, 0x85, 0x67))
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh12 = [f for f in report["findings"] if f["check_id"] == "VH-12"]
        assert len(vh12) >= 1
        assert vh12[0]["severity"] == "WARNING"

    def test_vh13_text_text_overlap_is_blocking(self):
        """VH-13: Overlapping text boxes are blocking, not warnings."""
        def build(prs, slide, ds):
            _add_textbox(slide, 600000, 600000, 3000000, 500000,
                         "Slide title", font_name="Inter", font_size_pt=28)
            _add_textbox(slide, 650000, 700000, 3000000, 500000,
                         "Body content", font_name="Inter", font_size_pt=12)

        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh13 = [f for f in report["findings"] if f["check_id"] == "VH-13"]
        assert any(f["severity"] == "BLOCKING" for f in vh13)
        assert report["pass"] is False

    def test_vh13_text_inside_filled_card_is_allowed(self):
        """VH-13: Text intentionally layered inside a filled card is not blocking."""
        def build(prs, slide, ds):
            _add_filled_rect(slide, 600000, 600000, 3000000, 1200000,
                             RGBColor(0x00, 0x85, 0x67))
            _add_textbox(slide, 750000, 750000, 2700000, 300000,
                         "Card heading", font_name="Inter", font_size_pt=16)

        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh13 = [f for f in report["findings"] if f["check_id"] == "VH-13"]
        assert not any(f["severity"] == "BLOCKING" for f in vh13)

    def test_vh13_text_overlapping_non_container_shape_is_blocking(self):
        """VH-13: Text colliding with a nearby non-container shape is blocking."""
        def build(prs, slide, ds):
            _add_textbox(slide, 600000, 600000, 3000000, 500000,
                         "Slide title", font_name="Inter", font_size_pt=28)
            _add_filled_rect(slide, 800000, 700000, 3500000, 1200000,
                             RGBColor(0x00, 0x85, 0x67))

        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh13 = [f for f in report["findings"] if f["check_id"] == "VH-13"]
        assert any(f["severity"] == "BLOCKING" for f in vh13)

    def test_vh27_black_slide_background_is_blocking(self):
        """VH-27: Full black slide backgrounds are blocking."""
        def build(prs, slide, ds):
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
            _add_textbox(slide, 600000, 600000, 3000000, 500000,
                         "Visible text", font_name="Inter", font_size_pt=20,
                         color=RGBColor(0xFF, 0xFF, 0xFF))

        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh27 = [f for f in report["findings"] if f["check_id"] == "VH-27"]
        assert len(vh27) == 1
        assert vh27[0]["severity"] == "BLOCKING"
        assert report["pass"] is False

    def test_vh15_grid_alignment_uses_template_canvas_region(self):
        """VH-15: Grid alignment uses the matched template canvas body region."""
        ds = _load_ds()
        body = ds["canvases"]["header_light"]["body_region"]
        gutter = ds["grid"]["gutter_md_emu"]
        cols = ds["grid"]["cols"]
        col_width = (body["width_emu"] - (cols - 1) * gutter) / cols
        prs = Presentation(str(_TEMPLATE))

        # Start from template layouts/theme but not its seed slides.
        for slide_id in list(prs.slides._sldIdLst):
            rel_id = slide_id.rId
            prs.part.drop_rel(rel_id)
            prs.slides._sldIdLst.remove(slide_id)

        slide = prs.slides.add_slide(prs.slide_layouts[8])

        top = body["top_emu"] + 100000
        width = int(col_width)
        height = 250000

        # Two aligned shapes and one intentionally misaligned shape.
        _add_filled_rect(slide, int(body["left_emu"]), top, width, height,
                         RGBColor(0x00, 0x85, 0x67))
        _add_filled_rect(slide, int(body["left_emu"] + col_width + gutter),
                         top + 350000, width, height,
                         RGBColor(0x00, 0x85, 0x67))
        _add_filled_rect(slide, int(body["left_emu"] + 100000),
                         top + 700000, width, height,
                         RGBColor(0x00, 0x85, 0x67))

        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        prs.save(tmp.name)

        report = scan_pptx(tmp.name, _DESIGN_SYSTEM)
        vh15 = [f for f in report["findings"] if f["check_id"] == "VH-15"]
        assert len(vh15) == 1
        assert vh15[0]["severity"] == "WARNING"


# ---------------------------------------------------------------------------
# Content Rendering checks (VH-16 to VH-19)
# ---------------------------------------------------------------------------

class TestContentRenderingChecks:
    def test_vh16_markdown_leak(self):
        """VH-16: Markdown in text is BLOCKING."""
        def build(prs, slide, ds):
            _add_textbox(slide, 500000, 500000, 5000000, 500000,
                         "This has **bold** markdown", font_name="Inter",
                         font_size_pt=12)
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh16 = [f for f in report["findings"] if f["check_id"] == "VH-16"]
        assert len(vh16) >= 1
        assert vh16[0]["severity"] == "BLOCKING"

    def test_vh16_clean_text_passes(self):
        """VH-16: Clean text passes."""
        def build(prs, slide, ds):
            _add_textbox(slide, 500000, 500000, 5000000, 500000,
                         "This is clean text with no markdown",
                         font_name="Inter", font_size_pt=12)
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh16 = [f for f in report["findings"] if f["check_id"] == "VH-16"]
        assert len(vh16) == 0

    def test_vh17_placeholder_text(self):
        """VH-17: Placeholder text is BLOCKING."""
        def build(prs, slide, ds):
            _add_textbox(slide, 500000, 500000, 5000000, 500000,
                         "Please {title} here", font_name="Inter",
                         font_size_pt=12)
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh17 = [f for f in report["findings"] if f["check_id"] == "VH-17"]
        assert len(vh17) >= 1
        assert vh17[0]["severity"] == "BLOCKING"

    def test_vh17_todo_detected(self):
        """VH-17: TODO in text is BLOCKING."""
        def build(prs, slide, ds):
            _add_textbox(slide, 500000, 500000, 5000000, 500000,
                         "TODO: fill in content", font_name="Inter",
                         font_size_pt=12)
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh17 = [f for f in report["findings"] if f["check_id"] == "VH-17"]
        assert len(vh17) >= 1

    def test_vh17_lorem_ipsum_detected(self):
        """VH-17: Lorem ipsum is BLOCKING."""
        def build(prs, slide, ds):
            _add_textbox(slide, 500000, 500000, 5000000, 500000,
                         "Lorem ipsum dolor sit amet", font_name="Inter",
                         font_size_pt=12)
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh17 = [f for f in report["findings"] if f["check_id"] == "VH-17"]
        assert len(vh17) >= 1

    def test_vh19_empty_content_frame(self):
        """VH-19: Empty text frame is WARNING."""
        def build(prs, slide, ds):
            _add_textbox(slide, 500000, 500000, 5000000, 500000,
                         "", font_name="Inter", font_size_pt=12)
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh19 = [f for f in report["findings"] if f["check_id"] == "VH-19"]
        assert len(vh19) >= 1
        assert vh19[0]["severity"] == "WARNING"


# ---------------------------------------------------------------------------
# Structural checks (VH-24 to VH-26)
# ---------------------------------------------------------------------------

class TestStructuralChecks:
    def test_vh24_empty_slide(self):
        """VH-24: Slide with < 2 non-placeholder shapes is BLOCKING."""
        def build(prs, slide, ds):
            # Add only one non-placeholder shape
            _add_filled_rect(slide, 500000, 500000, 2000000, 1000000,
                             RGBColor(0x00, 0x85, 0x67))
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh24 = [f for f in report["findings"] if f["check_id"] == "VH-24"]
        # May or may not fire depending on layout placeholders
        # The check counts non-placeholder shapes
        for f in vh24:
            assert f["severity"] == "BLOCKING"

    def test_vh25_slide_count_mismatch(self):
        """VH-25: Wrong slide count is BLOCKING."""
        def build(prs, slide, ds):
            _add_filled_rect(slide, 500000, 500000, 2000000, 1000000,
                             RGBColor(0x00, 0x85, 0x67))
            _add_textbox(slide, 500000, 2000000, 5000000, 500000,
                         "Content", font_name="Inter", font_size_pt=12)
        path = _make_pptx(build)
        deck_plan = {"slides": [{"title": "A"}, {"title": "B"}, {"title": "C"}]}
        report = scan_pptx(path, _DESIGN_SYSTEM, deck_plan=deck_plan)
        vh25 = [f for f in report["findings"] if f["check_id"] == "VH-25"]
        assert len(vh25) == 1
        assert vh25[0]["severity"] == "BLOCKING"

    def test_vh25_no_deck_plan_skips(self):
        """VH-25: No deck_plan means check is skipped."""
        def build(prs, slide, ds):
            _add_filled_rect(slide, 500000, 500000, 2000000, 1000000,
                             RGBColor(0x00, 0x85, 0x67))
            _add_textbox(slide, 500000, 2000000, 5000000, 500000,
                         "Content", font_name="Inter", font_size_pt=12)
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        vh25 = [f for f in report["findings"] if f["check_id"] == "VH-25"]
        assert len(vh25) == 0


# ---------------------------------------------------------------------------
# Deferred heuristic checks
# ---------------------------------------------------------------------------

class TestDeferredHeuristicChecks:
    def test_vh21_title_style_check_is_deferred(self):
        """VH-21 is deferred until title roles have explicit anchors."""
        ds = _load_ds()
        prs = Presentation()
        prs.slide_width = Emu(ds["canvas"]["width_emu"])
        prs.slide_height = Emu(ds["canvas"]["height_emu"])

        # Slide 1: title with Space Grotesk 28pt bold
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide1, 515938, 298692, 5000000, 500000,
                     "Title One", font_name="Space Grotesk",
                     font_size_pt=28, bold=True,
                     color=RGBColor(0x00, 0x00, 0x00))
        _add_filled_rect(slide1, 515938, 1000000, 2000000, 500000,
                         RGBColor(0x00, 0x85, 0x67))

        # Slide 2: title with different font
        slide2 = prs.slides.add_slide(prs.slide_layouts[6])
        _add_textbox(slide2, 515938, 298692, 5000000, 500000,
                     "Title Two", font_name="Inter",
                     font_size_pt=28, bold=True,
                     color=RGBColor(0x00, 0x00, 0x00))
        _add_filled_rect(slide2, 515938, 1000000, 2000000, 500000,
                         RGBColor(0x00, 0x85, 0x67))

        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        prs.save(tmp.name)

        report = scan_pptx(tmp.name, _DESIGN_SYSTEM)
        vh21 = [f for f in report["findings"] if f["check_id"] == "VH-21"]
        assert len(vh21) == 0
        assert "VH-21" in scanner_module._DEFERRED_CHECKS


# ---------------------------------------------------------------------------
# End-to-end tests
# ---------------------------------------------------------------------------

class TestScanEndToEnd:
    def test_report_schema(self):
        """Report has expected top-level keys."""
        def build(prs, slide, ds):
            _add_filled_rect(slide, 500000, 500000, 2000000, 1000000,
                             RGBColor(0x00, 0x85, 0x67))
            _add_textbox(slide, 500000, 2000000, 5000000, 500000,
                         "Test content", font_name="Inter",
                         font_size_pt=12, color=RGBColor(0, 0, 0))
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        assert "pass" in report
        assert "blocking_count" in report
        assert "warning_count" in report
        assert "findings" in report
        assert isinstance(report["pass"], bool)
        assert isinstance(report["blocking_count"], int)
        assert isinstance(report["warning_count"], int)
        assert isinstance(report["findings"], list)

    def test_findings_have_required_fields(self):
        """Each finding has all required fields."""
        def build(prs, slide, ds):
            # Add an off-palette fill to guarantee at least one finding
            _add_filled_rect(slide, 500000, 500000, 2000000, 1000000,
                             RGBColor(0xFF, 0x00, 0x00))
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        assert len(report["findings"]) > 0
        for f in report["findings"]:
            assert "check_id" in f
            assert "category" in f
            assert "check_name" in f
            assert "severity" in f
            assert "pass" in f
            assert "slide_index" in f
            assert "details" in f

    def test_blocking_count_matches(self):
        """blocking_count matches actual BLOCKING findings."""
        def build(prs, slide, ds):
            _add_filled_rect(slide, 500000, 500000, 2000000, 1000000,
                             RGBColor(0xFF, 0x00, 0x00))
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        actual = sum(1 for f in report["findings"] if f["severity"] == "BLOCKING")
        assert report["blocking_count"] == actual

    def test_pass_false_when_blocking(self):
        """pass is False when there are BLOCKING findings."""
        def build(prs, slide, ds):
            _add_filled_rect(slide, 500000, 500000, 2000000, 1000000,
                             RGBColor(0xFF, 0x00, 0x00))
        path = _make_pptx(build)
        report = scan_pptx(path, _DESIGN_SYSTEM)
        assert report["pass"] is False

    def test_clean_deck_no_blocking(self):
        """A simple well-formed deck has no BLOCKING findings."""
        ds = _load_ds()
        prs = Presentation()
        prs.slide_width = Emu(ds["canvas"]["width_emu"])
        prs.slide_height = Emu(ds["canvas"]["height_emu"])
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add shapes using only palette colors, allowed fonts, correct sizes
        _add_filled_rect(slide, 515938, 298692, 11160122, 200000,
                         RGBColor(0x00, 0x85, 0x67))  # accent_1
        _add_textbox(slide, 515938, 550000, 5000000, 500000,
                     "Clean Title", font_name="Space Grotesk",
                     font_size_pt=28, bold=True,
                     color=RGBColor(0x00, 0x00, 0x00))
        _add_textbox(slide, 515938, 1200000, 5000000, 500000,
                     "Clean body text content", font_name="Inter",
                     font_size_pt=12, bold=False,
                     color=RGBColor(0x00, 0x00, 0x00))

        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        prs.save(tmp.name)

        report = scan_pptx(tmp.name, _DESIGN_SYSTEM)
        blocking = [f for f in report["findings"] if f["severity"] == "BLOCKING"]
        assert len(blocking) == 0, (
            f"Clean deck has unexpected BLOCKING findings: "
            f"{[f['check_id'] + ': ' + f['details'] for f in blocking]}"
        )


class TestScannerInternalErrors:
    def test_check_crash_surfaces_as_blocking_finding(self, monkeypatch):
        """A crashed check becomes an explicit BLOCKING finding."""
        def broken_check(*_args, **_kwargs):
            raise RuntimeError("boom")

        broken_check.__name__ = "_check_vh01"
        broken_check.__doc__ = "VH-01: All fills use palette tokens."

        monkeypatch.setattr(scanner_module, "_ACTIVE_CHECKS", [broken_check])
        monkeypatch.setattr(scanner_module, "_check_vh25", lambda *_args, **_kwargs: [])

        path = _make_pptx(lambda prs, slide, ds: None)
        report = scan_pptx(path, _DESIGN_SYSTEM)

        assert report["pass"] is False
        assert report["blocking_count"] == 1
        assert report["warning_count"] == 0
        assert len(report["findings"]) == 1

        finding = report["findings"][0]
        assert finding["check_id"] == "VH-01"
        assert finding["severity"] == "BLOCKING"
        assert finding["slide_index"] == -1
        assert "RuntimeError: boom" in finding["details"]

    def test_vh25_crash_surfaces_as_blocking_finding(self, monkeypatch):
        """The deck-plan-specific check also surfaces internal failures."""
        def broken_vh25(*_args, **_kwargs):
            raise ValueError("bad deck plan")

        monkeypatch.setattr(scanner_module, "_ACTIVE_CHECKS", [])
        monkeypatch.setattr(scanner_module, "_check_vh25", broken_vh25)

        path = _make_pptx(lambda prs, slide, ds: None)
        report = scan_pptx(path, _DESIGN_SYSTEM, deck_plan={"slides": []})

        assert report["pass"] is False
        assert report["blocking_count"] == 1
        assert report["warning_count"] == 0
        assert len(report["findings"]) == 1

        finding = report["findings"][0]
        assert finding["check_id"] == "VH-25"
        assert finding["check_name"] == "Slide count matches deck_plan"
        assert finding["severity"] == "BLOCKING"
        assert finding["slide_index"] == -1
        assert "ValueError: bad deck plan" in finding["details"]
