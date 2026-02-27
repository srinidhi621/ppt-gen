"""Renderer tests."""

import json
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from src.config import load_config
from src.models.deck_ir import AssetRef, DeckIR, DeckSlide
from src.render.renderer import Renderer, parse_markdown_runs


class TestRenderer(unittest.TestCase):
    def test_render_writes_pptx_and_map(self) -> None:
        config = load_config()
        renderer = Renderer(
            Path(config.template_path),
            Path(config.layout_catalog_path),
            Path(config.icons_json_path),
        )
        deck = DeckIR(
            deck_id="deck_1",
            run_id="run_1",
            template_id="template_1",
            title="Title",
            slides=[
                DeckSlide(
                    slide_id="slide_1",
                    layout_id="one_content_light",
                    fields={"ph_title": "Hello", "ph_body": ["One", "Two"]},
                    speaker_notes="Notes",
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "out.pptx"
            render_map = renderer.render(deck, output_path)
            self.assertTrue(output_path.exists())
            self.assertIn("slide_1", render_map.entries)

            prs = Presentation(str(output_path))
            self.assertEqual(len(prs.slides), 1)
            entry = render_map.entries["slide_1"]
            self.assertIn("ph_title", entry.field_keys)
            self.assertIn("ph_body", entry.field_keys)

    def test_icon_index_includes_external_registry_and_resolves_paths(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            assets_dir = root / "assets"
            template_path = assets_dir / "template" / "template.pptx"
            layout_catalog_path = assets_dir / "layout" / "layout_catalog.json"
            icons_json_path = assets_dir / "icons" / "icons.json"
            local_icon_path = assets_dir / "icons" / "png" / "icon_001.png"
            external_icon_path = assets_dir / "external_assets" / "tabler" / "svg" / "cloud.svg"
            external_registry_path = assets_dir / "external_assets" / "registry.manifest.json"

            template_path.parent.mkdir(parents=True, exist_ok=True)
            layout_catalog_path.parent.mkdir(parents=True, exist_ok=True)
            icons_json_path.parent.mkdir(parents=True, exist_ok=True)
            local_icon_path.parent.mkdir(parents=True, exist_ok=True)
            external_icon_path.parent.mkdir(parents=True, exist_ok=True)
            external_registry_path.parent.mkdir(parents=True, exist_ok=True)

            template_path.write_bytes(b"")
            layout_catalog_path.write_text('{"layouts":[]}', encoding="utf-8")
            local_icon_path.write_bytes(b"local")
            external_icon_path.write_text(
                '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 24 24"></svg>',
                encoding="utf-8",
            )
            icons_json_path.write_text(
                '{"icons":[{"icon_id":"icon_001","filename":"icon_001.png"}]}',
                encoding="utf-8",
            )
            external_registry_path.write_text(
                json.dumps(
                    {
                        "icons": [
                            {
                                "id": "tabler:cloud",
                                "pack": "tabler",
                                "svg_path": "svg/cloud.svg",
                            }
                        ]
                    }
                ),
                encoding="utf-8",
            )

            renderer = Renderer(template_path, layout_catalog_path, icons_json_path)
            icon_index = renderer._load_icon_index()
            self.assertIn("icon_001", icon_index)
            self.assertIn("tabler:cloud", icon_index)

            local_asset = AssetRef(
                asset_type="icon",
                asset_id="icon_001",
                target_field_key="ph_image",
            )
            external_asset = AssetRef(
                asset_type="icon",
                asset_id="tabler:cloud",
                target_field_key="ph_image",
            )
            self.assertEqual(
                renderer._resolve_asset_path(local_asset, icon_index),
                local_icon_path,
            )
            with self.assertRaises(ValueError):
                renderer._resolve_asset_path(external_asset, icon_index)

    def test_icon_index_prefers_png_mapping_over_external_svg_for_same_id(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            assets_dir = root / "assets"
            template_path = assets_dir / "template" / "template.pptx"
            layout_catalog_path = assets_dir / "layout" / "layout_catalog.json"
            icons_json_path = assets_dir / "icons" / "icons.json"
            local_icon_path = assets_dir / "icons" / "png" / "external" / "tabler" / "cloud.png"
            external_icon_path = assets_dir / "external_assets" / "tabler" / "svg" / "cloud.svg"
            external_registry_path = assets_dir / "external_assets" / "registry.manifest.json"

            template_path.parent.mkdir(parents=True, exist_ok=True)
            layout_catalog_path.parent.mkdir(parents=True, exist_ok=True)
            icons_json_path.parent.mkdir(parents=True, exist_ok=True)
            local_icon_path.parent.mkdir(parents=True, exist_ok=True)
            external_icon_path.parent.mkdir(parents=True, exist_ok=True)
            external_registry_path.parent.mkdir(parents=True, exist_ok=True)

            template_path.write_bytes(b"")
            layout_catalog_path.write_text('{"layouts":[]}', encoding="utf-8")
            local_icon_path.write_bytes(b"local")
            external_icon_path.write_text(
                '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 24 24"></svg>',
                encoding="utf-8",
            )
            icons_json_path.write_text(
                '{"icons":[{"icon_id":"tabler:cloud","filename":"external/tabler/cloud.png"}]}',
                encoding="utf-8",
            )
            external_registry_path.write_text(
                json.dumps(
                    {
                        "icons": [
                            {
                                "id": "tabler:cloud",
                                "pack": "tabler",
                                "svg_path": "svg/cloud.svg",
                            }
                        ]
                    }
                ),
                encoding="utf-8",
            )

            renderer = Renderer(template_path, layout_catalog_path, icons_json_path)
            icon_index = renderer._load_icon_index()
            resolved = renderer._resolve_asset_path(
                AssetRef(asset_type="icon", asset_id="tabler:cloud", target_field_key="ph_image"),
                icon_index,
            )
            self.assertEqual(resolved, local_icon_path)


class TestMarkdownRendering(unittest.TestCase):
    def test_parse_bold(self) -> None:
        runs = parse_markdown_runs("**Data Integration**")
        self.assertEqual(len(runs), 1)
        self.assertEqual(runs[0], ("Data Integration", True, False))

    def test_parse_italic(self) -> None:
        runs = parse_markdown_runs("*emphasis here*")
        self.assertEqual(len(runs), 1)
        self.assertEqual(runs[0], ("emphasis here", False, True))

    def test_parse_bold_italic(self) -> None:
        runs = parse_markdown_runs("***bold and italic***")
        self.assertEqual(len(runs), 1)
        self.assertEqual(runs[0], ("bold and italic", True, True))

    def test_parse_mixed(self) -> None:
        runs = parse_markdown_runs("**Phase 1**: Discovery and design")
        self.assertEqual(len(runs), 2)
        self.assertEqual(runs[0], ("Phase 1", True, False))
        self.assertEqual(runs[1], (": Discovery and design", False, False))

    def test_plain_text_unchanged(self) -> None:
        runs = parse_markdown_runs("No formatting here")
        self.assertEqual(runs, [("No formatting here", False, False)])

    def test_multiple_bold_segments(self) -> None:
        runs = parse_markdown_runs("**A** and **B** end")
        self.assertEqual(len(runs), 4)
        self.assertEqual(runs[0], ("A", True, False))
        self.assertEqual(runs[1], (" and ", False, False))
        self.assertEqual(runs[2], ("B", True, False))
        self.assertEqual(runs[3], (" end", False, False))

    def test_render_bold_in_pptx(self) -> None:
        config = load_config()
        renderer = Renderer(
            Path(config.template_path),
            Path(config.layout_catalog_path),
            Path(config.icons_json_path),
        )
        deck = DeckIR(
            deck_id="md_test",
            run_id="run_md",
            template_id="t1",
            title="Test",
            slides=[
                DeckSlide(
                    slide_id="s1",
                    layout_id="one_content_light",
                    fields={
                        "ph_title": "**Bold Title**",
                        "ph_body": ["**Phase 1**: Discovery", "Plain bullet"],
                    },
                    speaker_notes="",
                )
            ],
        )
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "md_test.pptx"
            renderer.render(deck, output_path)
            prs = Presentation(str(output_path))
            slide = prs.slides[0]
            # Check that no literal ** appears in any shape text
            for shape in slide.shapes:
                if shape.has_text_frame:
                    full_text = shape.text_frame.text
                    self.assertNotIn("**", full_text, f"Literal ** found in shape text: {full_text}")
            # Verify bold runs exist
            found_bold = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.bold:
                                found_bold = True
            self.assertTrue(found_bold, "Expected at least one bold run in rendered PPTX")


if __name__ == "__main__":
    unittest.main()
