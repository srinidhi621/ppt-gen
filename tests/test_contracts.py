"""Tests for src.contracts.validator — JSON schema + artifact validation."""

import json
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu

from src.contracts.validator import (
    validate,
    validate_sandbox_to_scanner,
    validate_scanner_to_reviewer,
    PENDING_HANDOFFS,
)

_TEMPLATE = Path("assets/template/template.pptx")
_DESIGN_SYSTEM = Path("assets/template/design_system.json")


def _load_ds():
    with open(_DESIGN_SYSTEM) as f:
        return json.load(f)


def _make_simple_pptx(slide_count=1):
    """Create a minimal valid PPTX with the given slide count."""
    ds = _load_ds()
    prs = Presentation(str(_TEMPLATE))
    # Strip template seed slides
    for sid in list(prs.slides._sldIdLst):
        rel_id = sid.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(sid)

    layout_idx = ds["canvases"]["header_light"]["layout_index"]
    for _ in range(slide_count):
        prs.slides.add_slide(prs.slide_layouts[layout_idx])

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    return tmp.name


# ---------------------------------------------------------------------------
# Geometry report schema
# ---------------------------------------------------------------------------

class TestGeometryReportSchema:
    def test_valid_passing_report(self):
        report = {
            "pass": True,
            "blocking_count": 0,
            "warning_count": 0,
            "findings": [],
        }
        ok, errors = validate(report, "geometry_report")
        assert ok, f"Valid report failed: {errors}"

    def test_valid_failing_report(self):
        report = {
            "pass": False,
            "blocking_count": 1,
            "warning_count": 2,
            "findings": [
                {
                    "check_id": "VH-01",
                    "category": "Color",
                    "check_name": "All fills use palette tokens",
                    "severity": "BLOCKING",
                    "pass": False,
                    "slide_index": 0,
                    "details": "Shape 'Rect 1' has fill #FF0000",
                },
            ],
        }
        ok, errors = validate(report, "geometry_report")
        assert ok, f"Valid report failed: {errors}"

    def test_missing_pass_field(self):
        report = {
            "blocking_count": 0,
            "warning_count": 0,
            "findings": [],
        }
        ok, errors = validate(report, "geometry_report")
        assert not ok
        assert any("pass" in e for e in errors)

    def test_missing_findings(self):
        report = {
            "pass": True,
            "blocking_count": 0,
            "warning_count": 0,
        }
        ok, errors = validate(report, "geometry_report")
        assert not ok

    def test_invalid_severity(self):
        report = {
            "pass": False,
            "blocking_count": 1,
            "warning_count": 0,
            "findings": [
                {
                    "check_id": "VH-01",
                    "category": "Color",
                    "check_name": "test",
                    "severity": "CRITICAL",  # Invalid
                    "pass": False,
                    "slide_index": 0,
                    "details": "test",
                },
            ],
        }
        ok, errors = validate(report, "geometry_report")
        assert not ok

    def test_invalid_check_id_format(self):
        report = {
            "pass": False,
            "blocking_count": 1,
            "warning_count": 0,
            "findings": [
                {
                    "check_id": "INVALID",  # Invalid format
                    "category": "Color",
                    "check_name": "test",
                    "severity": "BLOCKING",
                    "pass": False,
                    "slide_index": 0,
                    "details": "test",
                },
            ],
        }
        ok, errors = validate(report, "geometry_report")
        assert not ok

    def test_extra_fields_rejected(self):
        report = {
            "pass": True,
            "blocking_count": 0,
            "warning_count": 0,
            "findings": [],
            "extra_field": "not allowed",
        }
        ok, errors = validate(report, "geometry_report")
        assert not ok


# ---------------------------------------------------------------------------
# Content fidelity report schema
# ---------------------------------------------------------------------------

class TestContentFidelityReportSchema:
    def test_valid_report(self):
        report = {
            "visible_coverage_score": 0.88,
            "notes_only_fact_count": 1,
            "total_facts": 12,
            "matched_visible_facts": 10,
            "matched_notes_only_facts": ["Q3 revenue grew 14%"],
            "dropped_facts": ["Q3 revenue grew 14%"],
            "hallucinated_specifics": [],
            "placeholder_leaks": [],
            "markdown_leaks": [],
        }
        ok, errors = validate(report, "content_fidelity_report")
        assert ok, f"Valid report failed: {errors}"

    def test_missing_coverage_score(self):
        report = {
            "notes_only_fact_count": 1,
            "total_facts": 12,
            "matched_visible_facts": 10,
            "matched_notes_only_facts": [],
            "dropped_facts": [],
            "hallucinated_specifics": [],
            "placeholder_leaks": [],
            "markdown_leaks": [],
        }
        ok, errors = validate(report, "content_fidelity_report")
        assert not ok

    def test_coverage_score_out_of_range(self):
        report = {
            "visible_coverage_score": 1.5,  # > 1
            "notes_only_fact_count": 0,
            "total_facts": 0,
            "matched_visible_facts": 0,
            "matched_notes_only_facts": [],
            "dropped_facts": [],
            "hallucinated_specifics": [],
            "placeholder_leaks": [],
            "markdown_leaks": [],
        }
        ok, errors = validate(report, "content_fidelity_report")
        assert not ok


# ---------------------------------------------------------------------------
# Deck plan schema
# ---------------------------------------------------------------------------

class TestDeckPlanSchema:
    def test_valid_deck_plan(self):
        plan = {
            "deck_title": "Q3 Review",
            "slides": [
                {
                    "slide_id": "cover",
                    "archetype": "hero_title",
                    "purpose": "introduce",
                    "audience_takeaway": "This is the Q3 review deck.",
                    "headline": "Q3 Business Review",
                    "body": "Company performance overview",
                },
                {
                    "slide_id": "kpis",
                    "archetype": "kpi_grid",
                    "purpose": "present_evidence",
                    "audience_takeaway": "Revenue grew 14% in Q3.",
                    "headline": "Key Metrics",
                    "metrics": [
                        {"value": "14%", "label": "Revenue Growth"},
                    ],
                },
            ],
        }
        ok, errors = validate(plan, "deck_plan")
        assert ok, f"Valid plan failed: {errors}"

    def test_missing_slides(self):
        plan = {"deck_title": "Test"}
        ok, errors = validate(plan, "deck_plan")
        assert not ok

    def test_empty_slides(self):
        plan = {"slides": []}
        ok, errors = validate(plan, "deck_plan")
        assert not ok

    def test_slide_missing_required_fields(self):
        plan = {
            "slides": [
                {"archetype": "hero_title"},
            ],
        }
        ok, errors = validate(plan, "deck_plan")
        assert not ok


# ---------------------------------------------------------------------------
# Normalized content schema
# ---------------------------------------------------------------------------

class TestNormalizedContentSchema:
    def test_valid_content(self):
        content = {
            "title": "Q3 Review",
            "sections": [
                {"heading": "Overview", "body": "Q3 was strong"},
            ],
        }
        ok, errors = validate(content, "normalized_content")
        assert ok, f"Valid content failed: {errors}"

    def test_missing_title(self):
        content = {
            "sections": [
                {"heading": "Overview", "body": "Q3 was strong"},
            ],
        }
        ok, errors = validate(content, "normalized_content")
        assert not ok

    def test_empty_title(self):
        content = {
            "title": "",
            "sections": [
                {"heading": "Overview", "body": "Q3 was strong"},
            ],
        }
        ok, errors = validate(content, "normalized_content")
        assert not ok

    def test_missing_sections(self):
        content = {"title": "Test"}
        ok, errors = validate(content, "normalized_content")
        assert not ok


# ---------------------------------------------------------------------------
# Review feedback schema
# ---------------------------------------------------------------------------

class TestReviewFeedbackSchema:
    def test_valid_feedback(self):
        feedback = {
            "overall_score": 8.5,
            "pass": True,
            "issues": [],
            "strengths": ["Clear layout", "Good use of color"],
            "summary": "Well-designed deck",
        }
        ok, errors = validate(feedback, "review_feedback")
        assert ok, f"Valid feedback failed: {errors}"

    def test_missing_overall_score(self):
        feedback = {
            "issues": [],
        }
        ok, errors = validate(feedback, "review_feedback")
        assert not ok

    def test_score_out_of_range(self):
        feedback = {
            "overall_score": 11,  # > 10
            "issues": [],
        }
        ok, errors = validate(feedback, "review_feedback")
        assert not ok

    def test_valid_issue(self):
        feedback = {
            "overall_score": 6.0,
            "issues": [
                {
                    "slide_index": 2,
                    "severity": "major",
                    "description": "Title too long",
                    "category": "content",
                    "suggested_fix": "Shorten to < 8 words",
                },
            ],
        }
        ok, errors = validate(feedback, "review_feedback")
        assert ok, f"Valid feedback failed: {errors}"


# ---------------------------------------------------------------------------
# Build exec report schema
# ---------------------------------------------------------------------------

class TestBuildExecReportSchema:
    def test_valid_report(self):
        report = {
            "success": True,
            "slides_built": 5,
            "pptx_path": "/tmp/output.pptx",
            "errors": [],
            "warnings": [],
            "build_time_seconds": 2.34,
        }
        ok, errors = validate(report, "build_exec_report")
        assert ok, f"Valid report failed: {errors}"

    def test_missing_success(self):
        report = {
            "slides_built": 5,
            "pptx_path": "/tmp/output.pptx",
        }
        ok, errors = validate(report, "build_exec_report")
        assert not ok

    def test_missing_pptx_path(self):
        report = {
            "success": True,
            "slides_built": 5,
        }
        ok, errors = validate(report, "build_exec_report")
        assert not ok


# ---------------------------------------------------------------------------
# Validator edge cases
# ---------------------------------------------------------------------------

class TestValidatorEdgeCases:
    def test_unknown_schema(self):
        ok, errors = validate({}, "nonexistent_schema")
        assert not ok
        assert len(errors) > 0

    def test_schema_name_with_extension(self):
        """Can use full filename as schema name."""
        report = {
            "pass": True,
            "blocking_count": 0,
            "warning_count": 0,
            "findings": [],
        }
        ok, errors = validate(report, "geometry_report.schema.json")
        assert ok, f"Valid report with extension failed: {errors}"


# ---------------------------------------------------------------------------
# Sandbox → Scanner artifact validator
# ---------------------------------------------------------------------------

class TestSandboxToScanner:
    def test_valid_pptx_passes(self):
        path = _make_simple_pptx(slide_count=2)
        ok, errors = validate_sandbox_to_scanner(path)
        assert ok, f"Valid PPTX failed: {errors}"

    def test_missing_file_fails(self):
        ok, errors = validate_sandbox_to_scanner("/tmp/nonexistent.pptx")
        assert not ok
        assert any("not found" in e for e in errors)

    def test_empty_file_fails(self):
        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        tmp.close()
        ok, errors = validate_sandbox_to_scanner(tmp.name)
        assert not ok
        assert any("empty" in e for e in errors)

    def test_invalid_pptx_fails(self):
        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        tmp.write(b"not a pptx file")
        tmp.close()
        ok, errors = validate_sandbox_to_scanner(tmp.name)
        assert not ok
        assert any("Cannot open" in e for e in errors)

    def test_slide_count_matches_plan(self):
        path = _make_simple_pptx(slide_count=3)
        plan = {"slides": [{"title": "a"}, {"title": "b"}, {"title": "c"}]}
        ok, errors = validate_sandbox_to_scanner(path, deck_plan=plan)
        assert ok, f"Matching slide count failed: {errors}"

    def test_slide_count_mismatch_fails(self):
        path = _make_simple_pptx(slide_count=2)
        plan = {"slides": [{"title": "a"}, {"title": "b"}, {"title": "c"}]}
        ok, errors = validate_sandbox_to_scanner(path, deck_plan=plan)
        assert not ok
        assert any("Slide count" in e for e in errors)

    def test_exec_report_success(self):
        path = _make_simple_pptx(slide_count=1)
        report = {
            "success": True,
            "slides_built": 1,
            "pptx_path": path,
            "errors": [],
            "warnings": [],
            "build_time_seconds": 1.0,
        }
        ok, errors = validate_sandbox_to_scanner(path, exec_report=report)
        assert ok, f"Valid exec report failed: {errors}"

    def test_exec_report_failure(self):
        path = _make_simple_pptx(slide_count=1)
        report = {
            "success": False,
            "slides_built": 0,
            "pptx_path": path,
            "errors": [{"slide_index": 0, "error": "build crashed"}],
            "warnings": [],
            "build_time_seconds": 0.5,
        }
        ok, errors = validate_sandbox_to_scanner(path, exec_report=report)
        assert not ok
        assert any("failure" in e for e in errors)

    def test_exec_report_invalid_schema(self):
        path = _make_simple_pptx(slide_count=1)
        bad_report = {"not_valid": True}
        ok, errors = validate_sandbox_to_scanner(path, exec_report=bad_report)
        assert not ok
        assert any("exec_report" in e for e in errors)


# ---------------------------------------------------------------------------
# Scanner / Fidelity → Reviewer artifact validator
# ---------------------------------------------------------------------------

class TestScannerToReviewer:
    def test_clean_report_passes(self):
        geo = {
            "pass": True,
            "blocking_count": 0,
            "warning_count": 1,
            "findings": [
                {
                    "check_id": "VH-12",
                    "category": "Spatial",
                    "check_name": "test",
                    "severity": "WARNING",
                    "pass": False,
                    "slide_index": 0,
                    "details": "minor overlap",
                },
            ],
        }
        ok, errors = validate_scanner_to_reviewer(geo)
        assert ok, f"Clean report failed: {errors}"

    def test_blocking_findings_fail(self):
        geo = {
            "pass": False,
            "blocking_count": 1,
            "warning_count": 0,
            "findings": [
                {
                    "check_id": "VH-01",
                    "category": "Color",
                    "check_name": "test",
                    "severity": "BLOCKING",
                    "pass": False,
                    "slide_index": 0,
                    "details": "bad color",
                },
            ],
        }
        ok, errors = validate_scanner_to_reviewer(geo)
        assert not ok
        assert any("BLOCKING" in e for e in errors)

    def test_invalid_geometry_schema_fails(self):
        bad = {"not_valid": True}
        ok, errors = validate_scanner_to_reviewer(bad)
        assert not ok
        assert any("geometry_report" in e for e in errors)

    def test_with_valid_fidelity_report(self):
        geo = {
            "pass": True,
            "blocking_count": 0,
            "warning_count": 0,
            "findings": [],
        }
        fidelity = {
            "visible_coverage_score": 0.95,
            "notes_only_fact_count": 0,
            "total_facts": 10,
            "matched_visible_facts": 10,
            "matched_notes_only_facts": [],
            "dropped_facts": [],
            "hallucinated_specifics": [],
            "placeholder_leaks": [],
            "markdown_leaks": [],
        }
        ok, errors = validate_scanner_to_reviewer(geo, fidelity)
        assert ok, f"Valid fidelity report failed: {errors}"

    def test_with_invalid_fidelity_report(self):
        geo = {
            "pass": True,
            "blocking_count": 0,
            "warning_count": 0,
            "findings": [],
        }
        bad_fidelity = {"not_valid": True}
        ok, errors = validate_scanner_to_reviewer(geo, bad_fidelity)
        assert not ok
        assert any("content_fidelity_report" in e for e in errors)


# ---------------------------------------------------------------------------
# Pending handoffs coverage
# ---------------------------------------------------------------------------

class TestPendingHandoffs:
    def test_pending_handoffs_documented(self):
        """All not-yet-wired handoffs are explicitly listed."""
        assert "normalize_to_planner" in PENDING_HANDOFFS
        assert "planner_to_feasibility" in PENDING_HANDOFFS
        assert "feasibility_to_builder" in PENDING_HANDOFFS
        assert "builder_to_sandbox" in PENDING_HANDOFFS
        assert "reviewer_to_repair" in PENDING_HANDOFFS
        assert "repair_to_accept" in PENDING_HANDOFFS
