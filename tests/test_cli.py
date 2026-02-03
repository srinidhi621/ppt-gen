"""CLI tests."""

import json
import shutil
import tempfile
import unittest
from pathlib import Path

from src.cli import cmd_render, cmd_smoke, cmd_validate, build_parser
from src.config import load_config


class MockArgs:
    """Mock argparse.Namespace for testing."""
    def __init__(self, **kwargs):
        for k, v in kwargs.items():
            setattr(self, k, v)


class TestCLI(unittest.TestCase):
    def test_validate_command_passes(self) -> None:
        args = MockArgs(project_root=None)
        result = cmd_validate(args)
        self.assertEqual(result, 0)

    def test_render_command_missing_deckir(self) -> None:
        args = MockArgs(project_root=None, deckir="/nonexistent/path.json", run_id=None)
        result = cmd_render(args)
        self.assertEqual(result, 1)

    def test_smoke_command_with_sample_deckir(self) -> None:
        config = load_config()
        sample_deckir = Path(config.inputs_dir) / "sample_deckir.json"
        
        if not sample_deckir.exists():
            self.skipTest("sample_deckir.json not found")
        
        with tempfile.TemporaryDirectory() as temp_dir:
            # Create a temporary runs directory
            temp_runs_dir = Path(temp_dir) / "runs"
            temp_runs_dir.mkdir()
            
            # We need to run with actual paths
            args = MockArgs(
                project_root=None,
                deckir=str(sample_deckir),
                run_id="test_smoke_run",
            )
            result = cmd_smoke(args)
            self.assertEqual(result, 0)
            
            # Check that artifacts were created
            run_dir = Path(config.runs_dir) / "test_smoke_run"
            self.assertTrue((run_dir / "deckir_v1.json").exists())
            self.assertTrue((run_dir / "deckir_v1_1.json").exists())
            self.assertTrue((run_dir / "validation_report.json").exists())
            self.assertTrue((run_dir / "render_map.json").exists())
            self.assertTrue((run_dir / "deck_v1.pptx").exists())
            self.assertTrue((run_dir / "run_log.jsonl").exists())
            
            # Clean up
            shutil.rmtree(run_dir)

    def test_parser_structure(self) -> None:
        parser = build_parser()
        # Test that subcommands exist
        args = parser.parse_args(["validate"])
        self.assertEqual(args.command, "validate")
        
        args = parser.parse_args(["render", "--deckir", "test.json"])
        self.assertEqual(args.command, "render")
        self.assertEqual(args.deckir, "test.json")
        
        args = parser.parse_args(["smoke"])
        self.assertEqual(args.command, "smoke")


class TestSmokeIntegration(unittest.TestCase):
    """Integration test for the full smoke pipeline."""
    
    def test_smoke_generates_valid_pptx(self) -> None:
        """Run smoke test and verify the PPTX is valid."""
        from pptx import Presentation
        
        config = load_config()
        sample_deckir = Path(config.inputs_dir) / "sample_deckir.json"
        
        if not sample_deckir.exists():
            self.skipTest("sample_deckir.json not found")
        
        args = MockArgs(
            project_root=None,
            deckir=str(sample_deckir),
            run_id="integration_test_run",
        )
        result = cmd_smoke(args)
        self.assertEqual(result, 0)
        
        # Verify PPTX is valid
        run_dir = Path(config.runs_dir) / "integration_test_run"
        pptx_path = run_dir / "deck_v1.pptx"
        
        prs = Presentation(str(pptx_path))
        self.assertGreater(len(prs.slides), 0)
        
        # Verify render map matches slide count
        with open(run_dir / "render_map.json") as f:
            render_map = json.load(f)
        self.assertEqual(len(render_map["entries"]), len(prs.slides))
        
        # Clean up
        shutil.rmtree(run_dir)


if __name__ == "__main__":
    unittest.main()
