"""Regression tests for alternate-approach/build_v3.py."""

import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation


REPO_ROOT = Path(__file__).resolve().parent.parent
BUILD_V3 = REPO_ROOT / "alternate-approach" / "build_v3.py"


def _slide_texts(slide):
    texts = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
        if text:
            texts.append(text)
    return texts


class TestBuildV3Script(unittest.TestCase):
    def test_build_v3_runs_standalone_and_outputs_clean_six_slide_deck(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            output_path = Path(tmpdir) / "runtime.pptx"
            subprocess.run(
                [sys.executable, str(BUILD_V3), str(output_path)],
                cwd=REPO_ROOT,
                check=True,
            )

            prs = Presentation(output_path)
            self.assertEqual(len(prs.slides), 6)

            for slide in prs.slides:
                placeholders = [shape for shape in slide.shapes if shape.is_placeholder]
                self.assertEqual(placeholders, [])

            self.assertIn(
                "Internal nominations run from day one and enter at Stage 2. Highest-yield source.",
                _slide_texts(prs.slides[1]),
            )


if __name__ == "__main__":
    unittest.main()
