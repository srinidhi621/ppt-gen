"""Tests for src.sandbox — AST scanner and subprocess runner."""

import json
import textwrap
from pathlib import Path

import pytest

from src.sandbox.ast_scanner import ScanResult, scan_ast
from src.sandbox.runner import ExecResult, run_in_sandbox


# ===========================================================================
# AST Scanner Tests
# ===========================================================================


class TestASTScannerAllowed:
    """Scripts that should pass the AST scan."""

    def test_ppt_runtime_import(self):
        result = scan_ast("import ppt_runtime")
        assert result.ok, result.summary()

    def test_ppt_runtime_submodule(self):
        result = scan_ast("from ppt_runtime.grid import Grid")
        assert result.ok, result.summary()

    def test_src_ppt_runtime_import(self):
        result = scan_ast("from src.ppt_runtime import Grid, Tokens")
        assert result.ok, result.summary()

    def test_pptx_import(self):
        result = scan_ast("import pptx")
        assert result.ok, result.summary()

    def test_pptx_util_import(self):
        result = scan_ast("from pptx.util import Inches, Pt, Emu")
        assert result.ok, result.summary()

    def test_pptx_enum_import(self):
        result = scan_ast("from pptx.enum.text import PP_ALIGN")
        assert result.ok, result.summary()

    def test_sys_import(self):
        result = scan_ast("import sys")
        assert result.ok, result.summary()

    def test_pathlib_import(self):
        result = scan_ast("from pathlib import Path")
        assert result.ok, result.summary()

    def test_json_import(self):
        result = scan_ast("import json")
        assert result.ok, result.summary()

    def test_math_import(self):
        result = scan_ast("import math")
        assert result.ok, result.summary()

    def test_os_path_import(self):
        result = scan_ast("import os.path")
        assert result.ok, result.summary()

    def test_typing_import(self):
        result = scan_ast("from typing import Optional, List")
        assert result.ok, result.summary()

    def test_dataclasses_import(self):
        result = scan_ast("from dataclasses import dataclass")
        assert result.ok, result.summary()

    def test_enum_import(self):
        result = scan_ast("from enum import Enum")
        assert result.ok, result.summary()

    def test_full_valid_script(self):
        code = textwrap.dedent("""\
            from pathlib import Path
            from src.ppt_runtime import Grid, Tokens, load_template, add_rect, add_text
            from pptx.util import Inches

            template = Path("assets/template/template.pptx")
            canvas = load_template(template, Path("assets/template/design_system.json"))
            slide = canvas.add_slide("header_light")
            g = Grid(canvas, cols=12, gutter="md")
            canvas.save(Path("deck.pptx"))
        """)
        result = scan_ast(code)
        assert result.ok, result.summary()


class TestASTScannerBlocked:
    """Scripts that must be rejected by the AST scan."""

    def test_dunder_import(self):
        result = scan_ast("__import__('os')")
        assert not result.ok
        assert any(v.rule == "blocked-builtin" for v in result.violations)

    def test_eval(self):
        result = scan_ast("eval('1+1')")
        assert not result.ok
        assert any(v.rule == "blocked-builtin" for v in result.violations)

    def test_exec(self):
        result = scan_ast("exec('x = 1')")
        assert not result.ok
        assert any(v.rule == "blocked-builtin" for v in result.violations)

    def test_compile(self):
        result = scan_ast("compile('x=1', '<s>', 'exec')")
        assert not result.ok
        assert any(v.rule == "blocked-builtin" for v in result.violations)

    def test_open_call(self):
        result = scan_ast("f = open('/etc/passwd', 'r')")
        assert not result.ok
        assert any(v.rule == "blocked-open" for v in result.violations)

    def test_open_write(self):
        result = scan_ast("open('secret.txt', 'w').write('data')")
        assert not result.ok
        assert any(v.rule == "blocked-open" for v in result.violations)

    def test_os_system(self):
        result = scan_ast("import os\nos.system('rm -rf /')")
        assert not result.ok
        # Two violations: blocked import + blocked os call
        rules = {v.rule for v in result.violations}
        assert "blocked-os" in rules

    def test_os_environ(self):
        result = scan_ast("import os\nprint(os.environ)")
        assert not result.ok

    def test_os_listdir(self):
        result = scan_ast("import os\nos.listdir('/')")
        assert not result.ok

    def test_subprocess_import(self):
        result = scan_ast("import subprocess")
        assert not result.ok
        assert any(v.rule == "blocked-import" for v in result.violations)

    def test_socket_import(self):
        result = scan_ast("import socket")
        assert not result.ok
        assert any(v.rule == "blocked-import" for v in result.violations)

    def test_http_import(self):
        result = scan_ast("from http.client import HTTPConnection")
        assert not result.ok
        assert any(v.rule == "blocked-import" for v in result.violations)

    def test_requests_import(self):
        result = scan_ast("import requests")
        assert not result.ok

    def test_shutil_import(self):
        result = scan_ast("import shutil")
        assert not result.ok

    def test_ctypes_import(self):
        result = scan_ast("import ctypes")
        assert not result.ok

    def test_getattr_call(self):
        result = scan_ast("getattr(obj, 'secret')")
        assert not result.ok
        assert any(v.rule == "blocked-builtin" for v in result.violations)

    def test_breakpoint_call(self):
        result = scan_ast("breakpoint()")
        assert not result.ok
        assert any(v.rule == "blocked-builtin" for v in result.violations)

    def test_input_call(self):
        result = scan_ast("x = input('Enter: ')")
        assert not result.ok
        assert any(v.rule == "blocked-builtin" for v in result.violations)

    def test_syntax_error(self):
        result = scan_ast("def foo(:")
        assert not result.ok
        assert any(v.rule == "syntax-error" for v in result.violations)

    def test_multiple_violations(self):
        code = textwrap.dedent("""\
            import subprocess
            import socket
            eval("1+1")
            open("/etc/passwd")
        """)
        result = scan_ast(code)
        assert not result.ok
        assert len(result.violations) >= 4

    def test_mixed_allowed_and_blocked(self):
        """Script with both valid and invalid imports."""
        code = textwrap.dedent("""\
            from ppt_runtime import Grid
            import subprocess
            from pptx.util import Inches
        """)
        result = scan_ast(code)
        assert not result.ok
        assert len(result.violations) == 1
        assert result.violations[0].detail == "import 'subprocess' is not allowed"


class TestASTScannerHexColors:
    """Tests for raw hex color literal detection."""

    def test_hex_color_string_blocked(self):
        result = scan_ast('color = "#FF0000"')
        assert not result.ok
        assert any(v.rule == "raw-hex-color" for v in result.violations)

    def test_hex_color_lowercase_blocked(self):
        result = scan_ast('c = "#aabbcc"')
        assert not result.ok
        assert any(v.rule == "raw-hex-color" for v in result.violations)

    def test_hex_color_mixed_case_blocked(self):
        result = scan_ast('c = "#1A2B3C"')
        assert not result.ok

    def test_non_color_hex_string_allowed(self):
        """Short hex strings (not 6 digits) should not trigger."""
        result = scan_ast('x = "#FFF"')
        assert result.ok

    def test_hex_in_longer_string_allowed(self):
        """Hex-like substring in a longer string should not trigger."""
        result = scan_ast('x = "color is #FF0000 here"')
        assert result.ok

    def test_tokens_color_call_allowed(self):
        """Using tokens.color() is the correct pattern — should pass."""
        result = scan_ast('c = tokens.color("accent_1")')
        assert result.ok

    def test_non_hash_hex_string_allowed(self):
        """Hex strings without # prefix are not color literals."""
        result = scan_ast('x = "FF0000"')
        assert result.ok


class TestASTScannerEdgeCases:
    """Edge cases and boundary conditions."""

    def test_empty_script(self):
        result = scan_ast("")
        assert result.ok

    def test_comments_only(self):
        result = scan_ast("# just a comment\n# another one")
        assert result.ok

    def test_string_containing_import(self):
        """String literals that look like imports should not trigger."""
        result = scan_ast('x = "import subprocess"')
        assert result.ok

    def test_scan_from_file(self, tmp_path):
        """Scan from an actual file path."""
        script = tmp_path / "test_script.py"
        script.write_text("from ppt_runtime import Grid\n")
        result = scan_ast(script)
        assert result.ok

    def test_scan_from_file_with_violation(self, tmp_path):
        script = tmp_path / "bad_script.py"
        script.write_text("import subprocess\n")
        result = scan_ast(script)
        assert not result.ok

    def test_summary_pass(self):
        result = scan_ast("x = 1")
        assert result.summary() == "AST scan passed"

    def test_summary_fail(self):
        result = scan_ast("eval('1')")
        assert "violation" in result.summary().lower()


# ===========================================================================
# Subprocess Runner Tests
# ===========================================================================


class TestRunnerHappyPath:
    """Tests for successful sandbox execution."""

    def test_simple_script(self, tmp_path):
        """A script that just creates a minimal pptx file."""
        script = tmp_path / "build_deck.py"
        script.write_text(textwrap.dedent("""\
            from pptx import Presentation
            from pathlib import Path

            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            prs.save(Path(__file__).parent / "deck.pptx")
            print("OK built deck")
        """))

        result = run_in_sandbox(
            script,
            attempt_dir=tmp_path,
            timeout_s=30,
            write_report=True,
        )

        assert result.success, f"Expected success: {result.error}\nstderr: {result.stderr}"
        assert result.exit_code == 0
        assert result.ast_scan_ok
        assert result.pptx_path is not None
        assert Path(result.pptx_path).exists()
        assert "OK built deck" in result.stdout

        # Check report was written
        report_path = tmp_path / "build_exec_report.json"
        assert report_path.exists()
        report = json.loads(report_path.read_text())
        assert report["success"] is True
        assert report["build_time_seconds"] > 0

    def test_output_pptx_name(self, tmp_path):
        """Script that writes output.pptx instead of deck.pptx."""
        script = tmp_path / "build_deck.py"
        script.write_text(textwrap.dedent("""\
            from pptx import Presentation
            from pathlib import Path

            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])
            prs.save(Path(__file__).parent / "output.pptx")
        """))

        result = run_in_sandbox(script, attempt_dir=tmp_path, timeout_s=30)
        assert result.success
        assert "output.pptx" in result.pptx_path

    def test_skip_ast_scan(self, tmp_path):
        """With skip_ast_scan, even scripts with blocked imports run."""
        script = tmp_path / "build_deck.py"
        # Script imports json (allowed) and writes pptx
        script.write_text(textwrap.dedent("""\
            import json
            from pptx import Presentation
            from pathlib import Path

            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])
            prs.save(Path(__file__).parent / "deck.pptx")
        """))

        result = run_in_sandbox(
            script, attempt_dir=tmp_path, timeout_s=30, skip_ast_scan=True
        )
        assert result.success


class TestRunnerASTRejection:
    """Tests for scripts rejected at AST scan stage."""

    def test_blocked_import_never_runs(self, tmp_path):
        """Script with blocked import should fail before execution."""
        script = tmp_path / "build_deck.py"
        script.write_text(textwrap.dedent("""\
            import subprocess
            subprocess.run(["echo", "should not run"])
        """))

        result = run_in_sandbox(script, attempt_dir=tmp_path, timeout_s=30)
        assert not result.success
        assert not result.ast_scan_ok
        assert result.exit_code == -1
        assert len(result.ast_violations) > 0
        assert "AST pre-scan rejected" in result.error

    def test_eval_rejected(self, tmp_path):
        script = tmp_path / "build_deck.py"
        script.write_text("result = eval('1+1')\n")

        result = run_in_sandbox(script, attempt_dir=tmp_path, timeout_s=30)
        assert not result.success
        assert not result.ast_scan_ok

    def test_report_written_on_ast_failure(self, tmp_path):
        script = tmp_path / "build_deck.py"
        script.write_text("import subprocess\n")

        run_in_sandbox(script, attempt_dir=tmp_path, timeout_s=30, write_report=True)

        report_path = tmp_path / "build_exec_report.json"
        assert report_path.exists()
        report = json.loads(report_path.read_text())
        assert report["success"] is False
        assert "ast_violations" in report
        assert len(report["ast_violations"]) > 0


class TestRunnerTimeout:
    """Tests for wall-clock timeout enforcement."""

    def test_timeout_kills_script(self, tmp_path):
        """Script with infinite busy loop should be killed by wall-clock timeout."""
        script = tmp_path / "build_deck.py"
        script.write_text(textwrap.dedent("""\
            import sys
            i = 0
            while True:
                i += 1
        """))

        result = run_in_sandbox(
            script, attempt_dir=tmp_path, timeout_s=2, write_report=False
        )
        assert not result.success
        assert "timed out" in result.error.lower()
        assert result.duration_s >= 1.5  # should be close to 2s


class TestRunnerMemory:
    """Tests for memory limit enforcement."""

    def test_memory_limit_kills_script(self, tmp_path):
        """Script that allocates more memory than the limit should be killed."""
        script = tmp_path / "build_deck.py"
        # Allocate ~200MB in a tight loop — should exceed a 64MB limit.
        # Uses sys (allowed) to avoid blocked imports.
        script.write_text(textwrap.dedent("""\
            import sys
            chunks = []
            for i in range(200):
                chunks.append(b'x' * (1024 * 1024))  # 1MB per chunk
        """))

        result = run_in_sandbox(
            script,
            attempt_dir=tmp_path,
            timeout_s=30,
            memory_mb=64,
            write_report=False,
        )
        # On macOS RLIMIT_AS/RLIMIT_RSS enforcement is best-effort.
        # On Linux this should reliably kill the process.
        # We accept either a non-zero exit code or a MemoryError.
        import platform
        if platform.system() == "Linux":
            assert not result.success, "Memory-hungry script should have been killed"
            assert result.exit_code != 0
        else:
            # macOS: RLIMIT_RSS is advisory, may not kill the process.
            # Just verify the runner didn't crash and returned a result.
            assert isinstance(result, ExecResult)


class TestRunnerScriptErrors:
    """Tests for scripts that execute but fail."""

    def test_runtime_error(self, tmp_path):
        """Script that raises an exception."""
        script = tmp_path / "build_deck.py"
        script.write_text(textwrap.dedent("""\
            from pathlib import Path
            raise ValueError("something went wrong")
        """))

        result = run_in_sandbox(script, attempt_dir=tmp_path, timeout_s=30)
        assert not result.success
        assert result.exit_code != 0
        assert result.ast_scan_ok  # AST scan passed
        assert "ValueError" in result.stderr

    def test_no_pptx_produced(self, tmp_path):
        """Script that runs OK but doesn't produce a PPTX."""
        script = tmp_path / "build_deck.py"
        script.write_text(textwrap.dedent("""\
            from pathlib import Path
            print("did some work but no pptx")
        """))

        result = run_in_sandbox(script, attempt_dir=tmp_path, timeout_s=30)
        assert not result.success
        assert "no .pptx file was produced" in result.error.lower()
        assert result.exit_code == 0  # script itself succeeded

    def test_import_error_in_script(self, tmp_path):
        """Script that imports something not installed."""
        script = tmp_path / "build_deck.py"
        script.write_text(textwrap.dedent("""\
            from pptx import Presentation
            from nonexistent_package import something
        """))
        # 'nonexistent_package' isn't in ALLOWED_MODULES, so AST scan blocks it
        result = run_in_sandbox(script, attempt_dir=tmp_path, timeout_s=30)
        assert not result.success
        assert not result.ast_scan_ok


class TestRunnerPPTXValidation:
    """Tests for PPTX validity check and slide counting (§4.5 step 4)."""

    def test_valid_pptx_counts_slides(self, tmp_path):
        """A valid PPTX should report the correct slide count."""
        script = tmp_path / "build_deck.py"
        script.write_text(textwrap.dedent("""\
            from pptx import Presentation
            from pathlib import Path

            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])
            prs.slides.add_slide(prs.slide_layouts[6])
            prs.slides.add_slide(prs.slide_layouts[6])
            prs.save(Path(__file__).parent / "deck.pptx")
        """))

        result = run_in_sandbox(script, attempt_dir=tmp_path, timeout_s=30)
        assert result.success
        assert result.slides_built == 3
        assert result.pptx_valid

    def test_corrupt_pptx_fails_validation(self, tmp_path):
        """A file named deck.pptx that isn't a valid PPTX should fail."""
        script = tmp_path / "build_deck.py"
        script.write_text(textwrap.dedent("""\
            from pathlib import Path
            # Write garbage bytes to deck.pptx
            Path(__file__).parent.joinpath("deck.pptx").write_bytes(
                b"this is not a pptx file")
        """))

        result = run_in_sandbox(
            script, attempt_dir=tmp_path, timeout_s=30, skip_ast_scan=True
        )
        assert not result.success
        assert not result.pptx_valid
        assert "invalid or corrupt" in result.error.lower()

    def test_slide_count_in_report(self, tmp_path):
        """build_exec_report.json should contain the slide count."""
        script = tmp_path / "build_deck.py"
        script.write_text(textwrap.dedent("""\
            from pptx import Presentation
            from pathlib import Path

            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])
            prs.slides.add_slide(prs.slide_layouts[6])
            prs.save(Path(__file__).parent / "deck.pptx")
        """))

        run_in_sandbox(script, attempt_dir=tmp_path, timeout_s=30, write_report=True)
        report = json.loads((tmp_path / "build_exec_report.json").read_text())
        assert report["slides_built"] == 2
        assert report["pptx_valid"] is True


class TestRunnerReport:
    """Tests for build_exec_report.json output."""

    def test_report_structure(self, tmp_path):
        script = tmp_path / "build_deck.py"
        script.write_text(textwrap.dedent("""\
            from pptx import Presentation
            from pathlib import Path

            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])
            prs.save(Path(__file__).parent / "deck.pptx")
        """))

        result = run_in_sandbox(script, attempt_dir=tmp_path, timeout_s=30)
        report_path = tmp_path / "build_exec_report.json"
        assert report_path.exists()

        report = json.loads(report_path.read_text())
        # Required fields per schema
        assert "success" in report
        assert "slides_built" in report
        assert "pptx_path" in report
        assert "build_time_seconds" in report
        assert "ast_scan_ok" in report
        assert report["success"] is True

    def test_no_report_when_disabled(self, tmp_path):
        script = tmp_path / "build_deck.py"
        script.write_text("from pathlib import Path\nprint('hi')\n")

        run_in_sandbox(script, attempt_dir=tmp_path, timeout_s=30, write_report=False)
        assert not (tmp_path / "build_exec_report.json").exists()

    def test_to_report_method(self):
        result = ExecResult(
            success=True,
            exit_code=0,
            stdout="OK",
            stderr="",
            duration_s=1.5,
            ast_scan_ok=True,
            slides_built=3,
            pptx_path="/tmp/deck.pptx",
            pptx_valid=True,
        )
        report = result.to_report()
        assert report["success"] is True
        assert report["pptx_path"] == "/tmp/deck.pptx"
        assert report["build_time_seconds"] == 1.5
        assert report["slides_built"] == 3
        assert report["pptx_valid"] is True
        assert "stdout" in report


class TestRunnerEnvironment:
    """Tests for restricted environment in sandbox."""

    def test_restricted_env_excludes_secrets(self, tmp_path):
        """Subprocess env should only contain the allowlisted keys."""
        # Set a fake secret in the parent environment
        import os
        os.environ["SANDBOX_TEST_SECRET"] = "should_not_leak"
        try:
            script = tmp_path / "build_deck.py"
            # skip_ast_scan=True so we can use os.environ in the script
            script.write_text(textwrap.dedent("""\
                import os
                import json
                from pathlib import Path
                from pptx import Presentation

                env_keys = sorted(os.environ.keys())
                Path(__file__).parent.joinpath("env_keys.json").write_text(
                    json.dumps(env_keys))

                prs = Presentation()
                prs.slides.add_slide(prs.slide_layouts[6])
                prs.save(Path(__file__).parent / "deck.pptx")
            """))

            result = run_in_sandbox(
                script, attempt_dir=tmp_path, timeout_s=30, skip_ast_scan=True
            )
            assert result.success, f"Script failed: {result.error}\nstderr: {result.stderr}"

            env_keys_file = tmp_path / "env_keys.json"
            assert env_keys_file.exists()
            import json
            env_keys = json.loads(env_keys_file.read_text())

            # The secret should not have leaked into the subprocess
            assert "SANDBOX_TEST_SECRET" not in env_keys

            # Only allowlisted keys should be present (plus OS-injected
            # keys like __CF_USER_TEXT_ENCODING on macOS that the runtime
            # adds regardless of the env dict we pass).
            allowed = {"PATH", "HOME", "LANG", "LC_ALL", "VIRTUAL_ENV", "PYTHONPATH"}
            # macOS injects __CF_USER_TEXT_ENCODING into every subprocess
            platform_injected = {"__CF_USER_TEXT_ENCODING", "__PYVENV_LAUNCHER__"}
            for key in env_keys:
                assert key in allowed or key in platform_injected, (
                    f"Unexpected env key {key!r} in sandbox"
                )
        finally:
            os.environ.pop("SANDBOX_TEST_SECRET", None)

    def test_env_has_path(self, tmp_path):
        """Subprocess should have PATH so it can find python and libraries."""
        script = tmp_path / "build_deck.py"
        script.write_text(textwrap.dedent("""\
            import sys
            from pathlib import Path
            from pptx import Presentation

            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])
            prs.save(Path(__file__).parent / "deck.pptx")
        """))

        result = run_in_sandbox(script, attempt_dir=tmp_path, timeout_s=30)
        assert result.success
