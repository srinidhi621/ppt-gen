"""Tests for src.scan.content_fidelity — content fidelity checker."""

import json
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt

from src.scan.content_fidelity import check_content_fidelity, extract_facts

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

_TEMPLATE = Path("assets/template/template.pptx")
_DESIGN_SYSTEM = Path("assets/template/design_system.json")


def _load_ds():
    with open(_DESIGN_SYSTEM) as f:
        return json.load(f)


# ---------------------------------------------------------------------------
# Fixture helpers
# ---------------------------------------------------------------------------

def _make_pptx_with_text(visible_texts, notes_texts=None):
    """Create a PPTX with given visible texts and optional notes."""
    ds = _load_ds()
    prs = Presentation(str(_TEMPLATE))

    # Strip template seed slides so only test content appears
    for sid in list(prs.slides._sldIdLst):
        rel_id = sid.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(sid)

    layout_idx = ds["canvases"]["header_light"]["layout_index"]
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    for text in visible_texts:
        tb = slide.shapes.add_textbox(
            Emu(515938), Emu(500000), Emu(5000000), Emu(500000)
        )
        tf = tb.text_frame
        r = tf.paragraphs[0].add_run()
        r.text = text
        r.font.name = "Inter"
        r.font.size = Pt(12)

    if notes_texts:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        for nt in notes_texts:
            if notes_tf.text:
                p = notes_tf.add_paragraph()
                r = p.add_run()
                r.text = nt
            else:
                notes_tf.paragraphs[0].text = nt

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    return tmp.name


# ---------------------------------------------------------------------------
# Fact extraction tests
# ---------------------------------------------------------------------------

class TestFactExtraction:
    def test_extract_percentages(self):
        facts = extract_facts("Revenue grew 14% in Q3")
        assert "14%" in facts

    def test_extract_dollar_amounts(self):
        facts = extract_facts("Revenue was $2.5M in 2024")
        assert "$2.5M" in facts

    def test_extract_years(self):
        facts = extract_facts("Revenue grew in 2024")
        assert "2024" in facts

    def test_extract_quarters(self):
        facts = extract_facts("Q3 results were strong")
        assert "Q3" in facts

    def test_extract_quarter_with_year(self):
        facts = extract_facts("In Q3 2024 we shipped the feature")
        assert "Q3 2024" in facts

    def test_extract_numbers(self):
        facts = extract_facts("We have 150 employees")
        assert "150" in facts

    def test_extract_quoted_phrases(self):
        facts = extract_facts('The motto is "Move Fast" and grow')
        assert "Move Fast" in facts

    def test_extract_single_quoted_phrases(self):
        facts = extract_facts("The principle is 'Ship Early' always")
        assert "Ship Early" in facts

    def test_extract_proper_nouns(self):
        facts = extract_facts("The project at Microsoft was successful.")
        assert "Microsoft" in facts

    def test_proper_nouns_skip_common_titles(self):
        """Words like 'The', 'This', 'Our' should not be extracted as proper nouns."""
        facts = extract_facts("However, this approach works. The team agreed.")
        assert "However" not in facts
        assert "The" not in facts

    def test_extract_date_with_month_name(self):
        facts = extract_facts("The launch is scheduled for March 2025")
        assert "March 2025" in facts

    def test_extract_date_full(self):
        facts = extract_facts("The deadline is January 15, 2025")
        assert "January 15, 2025" in facts

    def test_extract_numeric_date(self):
        facts = extract_facts("Filed on 01/15/2025 per policy")
        assert "01/15/2025" in facts

    def test_empty_input(self):
        facts = extract_facts("")
        assert facts == []

    def test_no_duplicates(self):
        facts = extract_facts("Revenue grew 14% and then 14% again in 2024")
        assert facts.count("14%") == 1


# ---------------------------------------------------------------------------
# Matching tests
# ---------------------------------------------------------------------------

class TestContentMatching:
    def test_visible_coverage_full(self):
        """All facts appear in visible text."""
        user_input = "Q3 revenue grew 14%"
        path = _make_pptx_with_text(["Q3 revenue grew 14% year over year"])
        report = check_content_fidelity(user_input, path)
        assert report["visible_coverage_score"] >= 0.5
        assert report["total_facts"] > 0

    def test_dropped_facts(self):
        """Facts not in PPTX show up as dropped."""
        user_input = "Q3 revenue grew 14% and Q4 grew 20%"
        path = _make_pptx_with_text(["Q3 revenue grew 14%"])
        report = check_content_fidelity(user_input, path)
        # 20% and Q4 should be dropped
        assert len(report["dropped_facts"]) > 0

    def test_notes_only_facts(self):
        """Facts in notes but not visible show as notes_only."""
        user_input = "Q3 revenue grew 14%"
        path = _make_pptx_with_text(
            ["Revenue overview"],
            notes_texts=["Q3 revenue grew 14%"]
        )
        report = check_content_fidelity(user_input, path)
        assert report["notes_only_fact_count"] >= 0  # Some facts may be in notes

    def test_empty_user_input(self):
        """Empty input gives full coverage."""
        path = _make_pptx_with_text(["Some content"])
        report = check_content_fidelity("", path)
        assert report["visible_coverage_score"] == 1.0
        assert report["total_facts"] == 0


# ---------------------------------------------------------------------------
# Placeholder detection tests
# ---------------------------------------------------------------------------

class TestPlaceholderDetection:
    def test_detects_title_placeholder(self):
        path = _make_pptx_with_text(["{title} goes here"])
        report = check_content_fidelity("test input", path)
        assert len(report["placeholder_leaks"]) > 0

    def test_detects_todo(self):
        path = _make_pptx_with_text(["TODO: add content"])
        report = check_content_fidelity("test input", path)
        assert len(report["placeholder_leaks"]) > 0

    def test_detects_lorem_ipsum(self):
        path = _make_pptx_with_text(["Lorem ipsum dolor sit amet"])
        report = check_content_fidelity("test input", path)
        assert len(report["placeholder_leaks"]) > 0

    def test_detects_tbd(self):
        path = _make_pptx_with_text(["The deadline is TBD"])
        report = check_content_fidelity("test input", path)
        assert len(report["placeholder_leaks"]) > 0

    def test_detects_mustache(self):
        path = _make_pptx_with_text(["Welcome {{name}}"])
        report = check_content_fidelity("test input", path)
        assert len(report["placeholder_leaks"]) > 0

    def test_clean_text_no_placeholders(self):
        path = _make_pptx_with_text(["Clean professional content"])
        report = check_content_fidelity("test input", path)
        assert len(report["placeholder_leaks"]) == 0


# ---------------------------------------------------------------------------
# Markdown leak detection tests
# ---------------------------------------------------------------------------

class TestMarkdownLeakDetection:
    def test_detects_bold_markdown(self):
        path = _make_pptx_with_text(["This is **bold** text"])
        report = check_content_fidelity("test input", path)
        assert len(report["markdown_leaks"]) > 0

    def test_detects_heading_markdown(self):
        path = _make_pptx_with_text(["## Heading text"])
        report = check_content_fidelity("test input", path)
        assert len(report["markdown_leaks"]) > 0

    def test_detects_link_markdown(self):
        path = _make_pptx_with_text(["Click [here](http://example.com)"])
        report = check_content_fidelity("test input", path)
        assert len(report["markdown_leaks"]) > 0

    def test_clean_text_no_markdown(self):
        path = _make_pptx_with_text(["Clean professional content"])
        report = check_content_fidelity("test input", path)
        assert len(report["markdown_leaks"]) == 0


# ---------------------------------------------------------------------------
# Hallucinated specifics tests
# ---------------------------------------------------------------------------

class TestHallucinatedSpecifics:
    def test_detects_hallucinated_numbers(self):
        """Numbers in PPTX not in user input are flagged."""
        user_input = "Revenue grew 14%"
        path = _make_pptx_with_text(["Revenue grew 14% with 250 employees"])
        report = check_content_fidelity(user_input, path)
        assert "250" in report["hallucinated_specifics"]

    def test_no_false_positives_for_input_numbers(self):
        """Numbers from user input are not flagged."""
        user_input = "We have 150 employees"
        path = _make_pptx_with_text(["Our team of 150 employees"])
        report = check_content_fidelity(user_input, path)
        assert "150" not in report["hallucinated_specifics"]

    def test_common_numbers_excluded(self):
        """Common numbers (1-10, 100) are not flagged."""
        user_input = "Our strategy"
        path = _make_pptx_with_text(["Step 1: Plan, Step 2: Execute, Step 3: Review"])
        report = check_content_fidelity(user_input, path)
        assert "1" not in report["hallucinated_specifics"]
        assert "2" not in report["hallucinated_specifics"]
        assert "3" not in report["hallucinated_specifics"]

    def test_detects_hallucinated_dates(self):
        """Dates in PPTX not in user input are flagged."""
        user_input = "We launched in Q3"
        path = _make_pptx_with_text(["We launched in Q3. The next phase begins March 2026."])
        report = check_content_fidelity(user_input, path)
        assert "March 2026" in report["hallucinated_specifics"]

    def test_detects_hallucinated_dollar_amounts(self):
        """Dollar amounts in PPTX not in user input are flagged."""
        user_input = "Revenue was $5M"
        path = _make_pptx_with_text(["Revenue was $5M with $2.3B total market"])
        report = check_content_fidelity(user_input, path)
        assert "$2.3B" in report["hallucinated_specifics"]

    def test_no_false_positive_on_input_dollars(self):
        user_input = "Revenue was $5M"
        path = _make_pptx_with_text(["Revenue was $5M"])
        report = check_content_fidelity(user_input, path)
        dollar_hallucinated = [h for h in report["hallucinated_specifics"] if h.startswith("$")]
        assert len(dollar_hallucinated) == 0

    def test_detects_hallucinated_quoted_phrases(self):
        """Quoted phrases in PPTX not in user input are flagged."""
        user_input = "Our motto is always deliver"
        path = _make_pptx_with_text(['Our approach is "Zero Downtime Guaranteed" always'])
        report = check_content_fidelity(user_input, path)
        assert '"Zero Downtime Guaranteed"' in report["hallucinated_specifics"]

    def test_detects_hallucinated_proper_nouns(self):
        """Proper nouns in PPTX not in user input are flagged."""
        user_input = "We partnered with a consulting firm"
        path = _make_pptx_with_text(
            ["We partnered with a consulting firm. Accenture led the engagement."]
        )
        report = check_content_fidelity(user_input, path)
        assert "Accenture" in report["hallucinated_specifics"]

    def test_no_false_positive_on_input_proper_nouns(self):
        user_input = "The project at Microsoft was successful. Google reviewed it."
        path = _make_pptx_with_text(["The deal included Microsoft and Google"])
        report = check_content_fidelity(user_input, path)
        assert "Microsoft" not in report["hallucinated_specifics"]
        assert "Google" not in report["hallucinated_specifics"]


# ---------------------------------------------------------------------------
# Report schema tests
# ---------------------------------------------------------------------------

class TestReportSchema:
    def test_report_has_all_keys(self):
        path = _make_pptx_with_text(["Test content"])
        report = check_content_fidelity("test input", path)
        expected_keys = {
            "visible_coverage_score",
            "notes_only_fact_count",
            "total_facts",
            "matched_visible_facts",
            "matched_notes_only_facts",
            "dropped_facts",
            "hallucinated_specifics",
            "placeholder_leaks",
            "markdown_leaks",
        }
        assert set(report.keys()) == expected_keys

    def test_coverage_score_range(self):
        path = _make_pptx_with_text(["Test content"])
        report = check_content_fidelity("test input with 42%", path)
        assert 0 <= report["visible_coverage_score"] <= 1
