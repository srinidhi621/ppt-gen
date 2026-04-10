"""Tests for builder sandbox execution harness."""

from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from src.builder import execute_builder_harness
from src.config import load_config


class TestBuilderHarness(unittest.TestCase):
    def setUp(self) -> None:
        self.config = load_config()
        self.tmp = tempfile.TemporaryDirectory()
        self.run_dir = Path(self.tmp.name) / "run"
        self.builder_input = {
            "template_pptx": self.config.template_path,
            "run_id": "builder_test",
        }

    def tearDown(self) -> None:
        self.tmp.cleanup()

    def test_trivial_builder_success(self) -> None:
        code = """
from pptx import Presentation
prs = Presentation(BUILDER_INPUT[\"template_pptx\"])
prs.slides.add_slide(prs.slide_layouts[0])
prs.save(BUILDER_INPUT[\"output_pptx\"])
"""
        result = execute_builder_harness(
            run_dir=self.run_dir,
            builder_input=self.builder_input,
            candidate_codes=[code],
            max_attempts=1,
            timeout_seconds=30,
        )

        self.assertEqual(result["status"], "success")
        self.assertTrue(Path(result["output_pptx"]).exists())
        self.assertTrue((self.run_dir / "build_exec_report_v1.json").exists())

    def test_retry_syntax_then_success(self) -> None:
        bad_code = "def nope(:\n    pass\n"
        good_code = """
from pptx import Presentation
prs = Presentation(BUILDER_INPUT[\"template_pptx\"])
prs.slides.add_slide(prs.slide_layouts[0])
prs.save(BUILDER_INPUT[\"output_pptx\"])
"""
        result = execute_builder_harness(
            run_dir=self.run_dir,
            builder_input=self.builder_input,
            candidate_codes=[bad_code, good_code],
            max_attempts=2,
            timeout_seconds=30,
        )

        self.assertEqual(result["status"], "success")
        self.assertEqual(result["attempt_count"], 2)
        self.assertEqual(result["attempts"][0]["failure_type"], "syntax_error")

    def test_import_allowlist_rejection(self) -> None:
        code = "import subprocess\nprint('blocked')\n"
        result = execute_builder_harness(
            run_dir=self.run_dir,
            builder_input=self.builder_input,
            candidate_codes=[code],
            max_attempts=1,
            timeout_seconds=30,
        )

        self.assertEqual(result["status"], "failed")
        self.assertEqual(result["failure_type"], "import_error")

    def test_alternate_approach_rerun(self) -> None:
        source = Path("alternate-approach") / "build.py"
        code = source.read_text(encoding="utf-8").replace(
            "BASE_DIR = Path(__file__).resolve().parent",
            "BASE_DIR = Path(BUILDER_INPUT['alternate_base_dir']).resolve()",
        )
        builder_input = dict(self.builder_input)
        builder_input["alternate_base_dir"] = str(source.parent.resolve())
        result = execute_builder_harness(
            run_dir=self.run_dir,
            builder_input=builder_input,
            candidate_codes=[code],
            max_attempts=1,
            timeout_seconds=60,
            import_allowlist={
                "pathlib",
                "sys",
                "pptx",
            },
        )

        self.assertEqual(result["status"], "success")
        pptx = Path(result["output_pptx"])
        self.assertTrue(pptx.exists())


if __name__ == "__main__":
    unittest.main()
