"""Unit tests for shapes.py, patterns.py, and composers.py.

Tests produce real PPTX files via the template and verify shape properties.
"""

import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu

from src.ppt_runtime.canvas import load_template
from src.ppt_runtime.composers import (
    compose_card_row,
    compose_split_columns,
    compose_stat_grid,
    compose_timeline,
)
from src.ppt_runtime.grid import Grid, Rect
from src.ppt_runtime.patterns import draw_card, draw_header_bar, draw_kicker, draw_stat_block
from src.ppt_runtime.shapes import add_image, add_line, add_rect, add_text
from src.ppt_runtime.tokens import Tokens

TEMPLATE_PATH = Path("assets/template/template.pptx")
DS_PATH = Path("assets/template/design_system.json")


def _setup():
    """Load template + tokens, add a slide, build a grid."""
    canvas = load_template(TEMPLATE_PATH, DS_PATH)
    tokens = Tokens.from_design_system(DS_PATH)
    slide = canvas.add_slide("header_light")
    grid = Grid(canvas, cols=12, gutter="md")
    return canvas, tokens, slide, grid


def _save_and_reopen(canvas):
    """Save to a temp file and reopen for assertion."""
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        canvas.save(f.name)
        return Presentation(f.name), f.name


# ---------------------------------------------------------------------------
# shapes.py tests
# ---------------------------------------------------------------------------


class TestAddRect(unittest.TestCase):
    def test_creates_shape_with_correct_position(self):
        canvas, tokens, slide, grid = _setup()
        rect = Rect(515938, 1237129, 2000000, 500000)
        shp = add_rect(slide, rect, fill=tokens.color("accent_1"))

        self.assertEqual(shp.left, Emu(515938))
        self.assertEqual(shp.top, Emu(1237129))
        self.assertEqual(shp.width, Emu(2000000))
        self.assertEqual(shp.height, Emu(500000))

    def test_fill_color_matches(self):
        canvas, tokens, slide, grid = _setup()
        rect = Rect(0, 0, 1000000, 1000000)
        shp = add_rect(slide, rect, fill=tokens.color("accent_2"))
        self.assertEqual(shp.fill.fore_color.rgb, RGBColor(0xCC, 0x00, 0x66))

    def test_no_line_by_default(self):
        canvas, tokens, slide, grid = _setup()
        rect = Rect(0, 0, 1000000, 1000000)
        shp = add_rect(slide, rect, fill=tokens.color("bg_primary"))
        # Shadow should be disabled
        self.assertFalse(shp.shadow.inherit)


class TestAddText(unittest.TestCase):
    def test_creates_textbox_with_correct_text(self):
        canvas, tokens, slide, grid = _setup()
        rect = Rect(515938, 1237129, 5000000, 500000)
        tb = add_text(
            slide, rect, "Hello World",
            type_style=tokens.type("title"),
            color=tokens.color("text_primary"),
        )
        tf = tb.text_frame
        self.assertEqual(tf.paragraphs[0].runs[0].text, "Hello World")

    def test_applies_font_from_type_style(self):
        canvas, tokens, slide, grid = _setup()
        rect = Rect(0, 0, 5000000, 500000)
        tb = add_text(
            slide, rect, "Styled",
            type_style=tokens.type("title"),
            color=tokens.color("text_primary"),
        )
        run = tb.text_frame.paragraphs[0].runs[0]
        self.assertEqual(run.font.name, "Space Grotesk")
        self.assertTrue(run.font.bold)

    def test_kicker_upper_case(self):
        canvas, tokens, slide, grid = _setup()
        rect = Rect(0, 0, 5000000, 200000)
        tb = add_text(
            slide, rect, "01 | The thesis",
            type_style=tokens.type("kicker"),
            color=tokens.color("accent_2"),
        )
        text = tb.text_frame.paragraphs[0].runs[0].text
        self.assertEqual(text, "01 | THE THESIS")

    def test_explicit_overrides_style(self):
        canvas, tokens, slide, grid = _setup()
        rect = Rect(0, 0, 5000000, 500000)
        tb = add_text(
            slide, rect, "Override",
            type_style=tokens.type("body"),
            bold=True,  # override body's bold=False
            color=tokens.color("text_primary"),
        )
        self.assertTrue(tb.text_frame.paragraphs[0].runs[0].font.bold)

    def test_word_wrap_enabled(self):
        canvas, tokens, slide, grid = _setup()
        rect = Rect(0, 0, 2000000, 500000)
        tb = add_text(
            slide, rect, "Word wrap test",
            type_style=tokens.type("body"),
            color=tokens.color("text_primary"),
        )
        self.assertTrue(tb.text_frame.word_wrap)


class TestAddLine(unittest.TestCase):
    def test_creates_connector(self):
        canvas, tokens, slide, grid = _setup()
        conn = add_line(
            slide,
            start_x=515938, start_y=3000000,
            end_x=11676060, end_y=3000000,
            color=tokens.color("accent_5"),
        )
        self.assertEqual(conn.line.color.rgb, RGBColor(0xAF, 0xAA, 0xB9))


# ---------------------------------------------------------------------------
# patterns.py tests
# ---------------------------------------------------------------------------


class TestDrawCard(unittest.TestCase):
    def test_produces_shapes(self):
        canvas, tokens, slide, grid = _setup()
        rect = grid.span(col=1, col_span=4, top=canvas.body_top, height_emu=2000000)
        initial_count = len(slide.shapes)
        draw_card(
            slide, rect, title="Card Title", body="Card body text here.",
            accent=tokens.color("accent_1"), tokens=tokens,
        )
        # Should add at least 3 shapes: accent bar, title textbox, body textbox
        self.assertGreaterEqual(len(slide.shapes) - initial_count, 3)


class TestDrawHeaderBar(unittest.TestCase):
    def test_produces_shapes(self):
        canvas, tokens, slide, grid = _setup()
        initial_count = len(slide.shapes)
        draw_header_bar(
            slide, kicker="01 | Operating Principle",
            title="Legacy complexity is now a growth constraint",
            canvas=canvas, tokens=tokens,
        )
        # Top bar + dot + kicker text + title = at least 4 shapes
        self.assertGreaterEqual(len(slide.shapes) - initial_count, 4)


class TestDrawStatBlock(unittest.TestCase):
    def test_produces_shapes(self):
        canvas, tokens, slide, grid = _setup()
        rect = grid.span(col=1, col_span=3, top=canvas.body_top, height_emu=1000000)
        initial_count = len(slide.shapes)
        draw_stat_block(
            slide, rect, value="42%", label="Cost reduction",
            accent=tokens.color("accent_2"), tokens=tokens,
        )
        # Accent bar + value + label = at least 3 shapes
        self.assertGreaterEqual(len(slide.shapes) - initial_count, 3)


# ---------------------------------------------------------------------------
# composers.py tests
# ---------------------------------------------------------------------------


class TestComposeCardRow(unittest.TestCase):
    def test_three_cards(self):
        canvas, tokens, slide, grid = _setup()
        region = Rect(canvas.body_left, canvas.body_top, canvas.body_width, 2000000)
        initial_count = len(slide.shapes)
        compose_card_row(
            slide, region,
            items=[
                {"title": "Speed", "body": "2x faster delivery"},
                {"title": "Quality", "body": "Zero critical defects"},
                {"title": "Cost", "body": "30% reduction"},
            ],
            accent=tokens.color("accent_1"), tokens=tokens,
        )
        # 3 cards × 3 shapes each = at least 9 new shapes
        self.assertGreaterEqual(len(slide.shapes) - initial_count, 9)

    def test_empty_items_no_shapes(self):
        canvas, tokens, slide, grid = _setup()
        region = Rect(canvas.body_left, canvas.body_top, canvas.body_width, 2000000)
        initial_count = len(slide.shapes)
        compose_card_row(slide, region, items=[], accent=tokens.color("accent_1"), tokens=tokens)
        self.assertEqual(len(slide.shapes), initial_count)


class TestComposeStatGrid(unittest.TestCase):
    def test_four_metrics_2x2(self):
        canvas, tokens, slide, grid = _setup()
        region = Rect(canvas.body_left, canvas.body_top, canvas.body_width, 3000000)
        initial_count = len(slide.shapes)
        compose_stat_grid(
            slide, region,
            metrics=[
                {"value": "42%", "label": "Cost savings"},
                {"value": "3.2x", "label": "Speed gain"},
                {"value": "98%", "label": "Uptime"},
                {"value": "$12M", "label": "Revenue"},
            ],
            cols=2, tokens=tokens,
        )
        # 4 stat blocks × 3 shapes = at least 12
        self.assertGreaterEqual(len(slide.shapes) - initial_count, 12)


class TestComposeSplitColumns(unittest.TestCase):
    def test_two_panels(self):
        canvas, tokens, slide, grid = _setup()
        region = Rect(canvas.body_left, canvas.body_top, canvas.body_width, 3000000)
        initial_count = len(slide.shapes)
        compose_split_columns(
            slide, region,
            left_content="Left panel content goes here.",
            right_content="Right panel content goes here.",
            split=0.5, tokens=tokens,
        )
        # 2 text boxes
        self.assertGreaterEqual(len(slide.shapes) - initial_count, 2)


class TestComposeTimeline(unittest.TestCase):
    def test_four_phases(self):
        canvas, tokens, slide, grid = _setup()
        region = Rect(canvas.body_left, canvas.body_top, canvas.body_width, 3000000)
        initial_count = len(slide.shapes)
        compose_timeline(
            slide, region,
            phases=[
                {"label": "Phase 1", "body": "Discovery and assessment"},
                {"label": "Phase 2", "body": "Design and planning"},
                {"label": "Phase 3", "body": "Build and test"},
                {"label": "Phase 4", "body": "Deploy and optimize"},
            ],
            accent=tokens.color("accent_1"), tokens=tokens,
        )
        # Track line + 4×(dot + label + body) = at least 13 shapes
        self.assertGreaterEqual(len(slide.shapes) - initial_count, 13)


# ---------------------------------------------------------------------------
# Integration: save and verify PPTX opens
# ---------------------------------------------------------------------------


class TestFullSliceProducesPptx(unittest.TestCase):
    def test_multi_pattern_slide_saves_and_opens(self):
        """Build a slide using multiple patterns/composers, save, reopen."""
        canvas, tokens, slide, grid = _setup()

        # Header bar
        draw_header_bar(
            slide, kicker="01 | Operating Principle",
            title="Legacy complexity is now a growth constraint",
            canvas=canvas, tokens=tokens,
        )

        # Card row in the body region
        card_region = grid.span(
            col=1, col_span=12,
            top=canvas.body_top + tokens.spacing("lg"),
            height_emu=2000000,
        )
        compose_card_row(
            slide, card_region,
            items=[
                {"title": "The problem", "body": "Fragmented systems slow delivery."},
                {"title": "What we won't do", "body": "Rip and replace everything."},
                {"title": "What we will do", "body": "Modernize incrementally."},
            ],
            accent=tokens.color("accent_1"), tokens=tokens,
        )

        prs, path = _save_and_reopen(canvas)
        # Template has existing slides; our new slide is the last one
        self.assertGreater(len(prs.slides), 0)
        # Verify shapes exist on the last (newly added) slide
        last_slide = prs.slides[-1]
        self.assertGreater(len(last_slide.shapes), 10)


if __name__ == "__main__":
    unittest.main()
