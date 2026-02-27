"""Diagnostic script: parse a rendered PPTX and compare against its DeckIR source.

Usage:
    python scripts/diagnose_pptx.py runs/visual_polish_legacy_final

Outputs a per-slide report showing:
  - Layout used, placeholder shapes found, text content, images placed
  - Comparison with DeckIR expectations (fields, asset_refs)
  - Text overflow analysis against layout constraints
  - Missing/orphaned placeholders
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE


def emu_to_inches(emu: int) -> float:
    return round(emu / 914400, 2)


def read_alt_text(shape) -> str:
    """Extract alt-text / description from a shape."""
    elem = shape.element
    for tag in ("p:nvSpPr", "p:nvPicPr", "p:nvGrpSpPr"):
        nv = elem.find(f"{{{elem.nsmap.get('p', 'http://schemas.openxmlformats.org/presentationml/2006/main')}}}{tag.split(':')[1]}")
        if nv is None:
            # Try with full namespace
            from pptx.oxml.ns import qn
            nv = elem.find(qn(tag))
        if nv is not None:
            from pptx.oxml.ns import qn
            cNvPr = nv.find(qn("p:cNvPr"))
            if cNvPr is not None:
                return cNvPr.get("descr", "")
    return ""


def get_shape_text(shape) -> str:
    """Get all text from a shape's text frame."""
    if not shape.has_text_frame:
        return ""
    parts = []
    for para in shape.text_frame.paragraphs:
        parts.append(para.text)
    return "\n".join(parts)


def get_paragraph_details(shape) -> List[Dict[str, Any]]:
    """Get detailed paragraph info including run-level formatting."""
    if not shape.has_text_frame:
        return []
    details = []
    for para in shape.text_frame.paragraphs:
        runs = []
        for run in para.runs:
            runs.append({
                "text": run.text,
                "bold": run.font.bold,
                "italic": run.font.italic,
                "size_pt": round(run.font.size.pt, 1) if run.font.size else None,
                "font_name": run.font.name,
            })
        details.append({
            "text": para.text,
            "level": para.level,
            "runs": runs,
            "char_count": len(para.text),
        })
    return details


def analyze_shape(shape) -> Dict[str, Any]:
    """Analyze a single shape."""
    info: Dict[str, Any] = {
        "shape_id": shape.shape_id,
        "name": shape.name,
        "alt_text": read_alt_text(shape),
        "shape_type": str(shape.shape_type) if shape.shape_type else "unknown",
        "is_placeholder": shape.is_placeholder,
        "position": {
            "left_in": emu_to_inches(shape.left),
            "top_in": emu_to_inches(shape.top),
            "width_in": emu_to_inches(shape.width),
            "height_in": emu_to_inches(shape.height),
        },
    }
    if shape.is_placeholder:
        info["placeholder_idx"] = shape.placeholder_format.idx
        info["placeholder_type"] = str(shape.placeholder_format.type)

    # Check if it's a picture
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        info["is_picture"] = True
        info["image_content_type"] = shape.image.content_type if hasattr(shape, 'image') else "unknown"
    elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        # Check if placeholder contains an image
        try:
            if hasattr(shape, 'image'):
                info["has_image"] = True
                info["image_content_type"] = shape.image.content_type
            else:
                info["has_image"] = False
        except Exception:
            info["has_image"] = False

    # Text content
    if shape.has_text_frame:
        info["has_text"] = True
        info["text_preview"] = get_shape_text(shape)[:200]
        info["total_chars"] = len(get_shape_text(shape))
        info["paragraphs"] = get_paragraph_details(shape)
        info["paragraph_count"] = len(shape.text_frame.paragraphs)
    else:
        info["has_text"] = False

    return info


def load_deckir(run_dir: Path) -> Dict[str, Any] | None:
    deckir_path = run_dir / "deckir_v1_1.json"
    if not deckir_path.exists():
        deckir_path = run_dir / "deckir_v1.json"
    if not deckir_path.exists():
        return None
    with open(deckir_path, "r") as f:
        return json.load(f)


def load_layout_catalog() -> Dict[str, Any]:
    catalog_path = Path("assets/layout/layout_catalog.json")
    with open(catalog_path, "r") as f:
        catalog = json.load(f)
    return {entry["layout_id"]: entry for entry in catalog.get("layouts", [])}


def check_text_overflow(text: str, constraint_chars: int, label: str) -> Dict[str, Any]:
    """Check if text exceeds character budget."""
    char_count = len(text)
    overflow = char_count - constraint_chars
    return {
        "field": label,
        "chars": char_count,
        "budget": constraint_chars,
        "overflow": overflow if overflow > 0 else 0,
        "overflowing": overflow > 0,
    }


def diagnose_slide(slide, slide_index: int, deckir_slide: Dict[str, Any] | None,
                    layout_catalog: Dict[str, Any]) -> Dict[str, Any]:
    """Full diagnosis of a single slide."""
    report: Dict[str, Any] = {
        "slide_index": slide_index,
        "slide_layout_name": slide.slide_layout.name if slide.slide_layout else "unknown",
    }

    if deckir_slide:
        report["deckir_slide_id"] = deckir_slide.get("slide_id", "?")
        report["deckir_layout_id"] = deckir_slide.get("layout_id", "?")
        report["deckir_fields"] = list(deckir_slide.get("fields", {}).keys())
        report["deckir_asset_refs"] = deckir_slide.get("asset_refs", [])

    # Analyze all shapes
    shapes_report = []
    text_shapes = {}  # alt_text -> shape info
    image_shapes = []
    unbound_placeholders = []

    for shape in slide.shapes:
        shape_info = analyze_shape(shape)
        shapes_report.append(shape_info)

        alt = shape_info["alt_text"]
        if alt and shape_info.get("has_text"):
            text_shapes[alt] = shape_info
        if shape_info.get("is_picture") or shape_info.get("has_image"):
            image_shapes.append(shape_info)
        if shape_info["is_placeholder"] and not alt and not shape_info.get("is_picture"):
            unbound_placeholders.append(shape_info)

    report["total_shapes"] = len(shapes_report)
    report["text_shapes_by_alt"] = {k: {
        "text_preview": v["text_preview"],
        "total_chars": v["total_chars"],
        "paragraph_count": v["paragraph_count"],
    } for k, v in text_shapes.items()}
    report["image_shapes_count"] = len(image_shapes)
    report["image_shapes"] = image_shapes
    report["unbound_placeholders"] = len(unbound_placeholders)

    # Check text overflow against layout constraints
    if deckir_slide:
        layout_id = deckir_slide.get("layout_id", "")
        layout_entry = layout_catalog.get(layout_id, {})
        constraints = layout_entry.get("constraints", {})
        overflow_checks = []

        # Title overflow
        if "ph_title" in text_shapes:
            max_title = constraints.get("max_title_chars", 999)
            overflow_checks.append(
                check_text_overflow(text_shapes["ph_title"]["text_preview"], max_title, "ph_title")
            )

        # Body overflow - aggregate all body/col fields
        body_keys = [k for k in text_shapes if k.startswith("ph_body") or k.startswith("ph_col")]
        total_body_chars = sum(text_shapes[k]["total_chars"] for k in body_keys)
        max_body = constraints.get("max_total_body_chars", 9999)
        overflow_checks.append({
            "field": "total_body",
            "chars": total_body_chars,
            "budget": max_body,
            "overflow": max(0, total_body_chars - max_body),
            "overflowing": total_body_chars > max_body,
        })

        # Per-bullet check
        max_bullets = constraints.get("max_bullets", 99)
        for key in body_keys:
            if key in text_shapes:
                pcount = text_shapes[key]["paragraph_count"]
                if pcount > max_bullets and max_bullets > 0:
                    overflow_checks.append({
                        "field": f"{key}_bullet_count",
                        "chars": pcount,
                        "budget": max_bullets,
                        "overflow": pcount - max_bullets,
                        "overflowing": True,
                    })

        report["overflow_checks"] = overflow_checks
        report["has_overflow"] = any(c["overflowing"] for c in overflow_checks)

    # Compare expected vs actual images
    if deckir_slide:
        expected_assets = deckir_slide.get("asset_refs", [])
        report["expected_images"] = len(expected_assets)
        report["actual_images"] = len(image_shapes)
        report["image_gap"] = len(expected_assets) - len(image_shapes)

    # Speaker notes
    try:
        notes_text = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        report["has_speaker_notes"] = bool(notes_text.strip())
        report["speaker_notes_preview"] = notes_text[:150]
    except Exception:
        report["has_speaker_notes"] = False

    # Full shapes detail for deep dive
    report["all_shapes"] = shapes_report

    return report


def main():
    parser = argparse.ArgumentParser(description="Diagnose a rendered PPTX against DeckIR.")
    parser.add_argument("run_dir", type=str, help="Path to run directory containing deck_v1.pptx")
    parser.add_argument(
        "--pptx",
        type=str,
        default="deck_v1.pptx",
        help="PPTX filename inside run_dir (default: deck_v1.pptx)",
    )
    parser.add_argument(
        "--json-out",
        type=str,
        default=None,
        help="Optional path to write machine-readable diagnose JSON report",
    )
    args = parser.parse_args()

    run_dir = Path(args.run_dir)
    pptx_path = run_dir / args.pptx
    if not pptx_path.exists():
        print(f"ERROR: {pptx_path} not found")
        sys.exit(1)

    deckir = load_deckir(run_dir)
    layout_catalog = load_layout_catalog()

    prs = Presentation(str(pptx_path))
    deckir_slides = deckir.get("slides", []) if deckir else []

    print("=" * 80)
    print(f"PPTX DIAGNOSTIC: {pptx_path}")
    print(f"Slides in PPTX: {len(prs.slides)}")
    print(f"Slides in DeckIR: {len(deckir_slides)}")
    print("=" * 80)

    # Summary table
    print(f"\n{'Slide':>5} | {'Layout':>30} | {'Fields':>15} | {'Imgs':>4} | {'Overflow':>8} | Title")
    print("-" * 100)

    slide_reports = []
    for idx, slide in enumerate(prs.slides):
        deckir_slide = deckir_slides[idx] if idx < len(deckir_slides) else None
        report = diagnose_slide(slide, idx, deckir_slide, layout_catalog)
        slide_reports.append(report)

        layout_id = report.get("deckir_layout_id", "?")
        fields = ", ".join(report.get("deckir_fields", []))
        imgs = f"{report.get('actual_images', 0)}/{report.get('expected_images', 0)}"
        overflow = "YES" if report.get("has_overflow") else "ok"
        title = list(report.get("text_shapes_by_alt", {}).get("ph_title", {}).values())[0][:40] if "ph_title" in report.get("text_shapes_by_alt", {}) else "?"
        print(f"{idx+1:>5} | {layout_id:>30} | {fields:>15} | {imgs:>4} | {overflow:>8} | {title}")

    # Detailed per-slide reports
    print("\n" + "=" * 80)
    print("DETAILED SLIDE REPORTS")
    print("=" * 80)

    for report in slide_reports:
        idx = report["slide_index"]
        print(f"\n{'─' * 80}")
        print(f"SLIDE {idx + 1}: {report.get('deckir_slide_id', '?')}")
        print(f"  Layout: {report.get('deckir_layout_id', '?')} ({report.get('slide_layout_name', '?')})")
        print(f"  Total shapes: {report['total_shapes']}")
        print(f"  Images: {report.get('actual_images', 0)} actual / {report.get('expected_images', 0)} expected")

        # Text fields
        print(f"  Text fields bound:")
        for alt, info in report.get("text_shapes_by_alt", {}).items():
            print(f"    {alt}: {info['total_chars']} chars, {info['paragraph_count']} paragraphs")
            preview = info["text_preview"].replace("\n", " | ")
            print(f"      Preview: {preview[:100]}")

        # Overflow
        if report.get("overflow_checks"):
            print(f"  Overflow analysis:")
            for check in report["overflow_checks"]:
                status = "OVERFLOW" if check["overflowing"] else "ok"
                print(f"    {check['field']}: {check['chars']}/{check['budget']} [{status}]")

        # Image details
        if report.get("image_shapes"):
            print(f"  Image shapes:")
            for img in report["image_shapes"]:
                pos = img["position"]
                print(f"    {img['name']}: {pos['width_in']}x{pos['height_in']}in at ({pos['left_in']},{pos['top_in']})")
                if img.get("image_content_type"):
                    print(f"      Content type: {img['image_content_type']}")

        # Expected asset refs from DeckIR
        if report.get("deckir_asset_refs"):
            print(f"  Expected asset_refs from DeckIR:")
            for ref in report["deckir_asset_refs"]:
                print(f"    {ref.get('asset_type')}: {ref.get('asset_id')} -> {ref.get('target_field_key')}")

        # Unbound
        if report.get("unbound_placeholders"):
            print(f"  WARNING: {report['unbound_placeholders']} unbound placeholder(s)")

        # All shapes (compact)
        print(f"  All shapes:")
        for s in report.get("all_shapes", []):
            ph_info = f" [ph_idx={s.get('placeholder_idx')}]" if s.get("is_placeholder") else ""
            alt_info = f" alt='{s['alt_text']}'" if s.get("alt_text") else ""
            pic_info = " [PICTURE]" if s.get("is_picture") else ""
            text_info = f" [{s.get('total_chars', 0)}ch]" if s.get("has_text") else ""
            print(f"    #{s['shape_id']} {s['name']}{ph_info}{alt_info}{pic_info}{text_info}")

    # Final summary
    print("\n" + "=" * 80)
    print("SUMMARY")
    print("=" * 80)
    total_overflow = sum(1 for r in slide_reports if r.get("has_overflow"))
    total_image_gap = sum(r.get("image_gap", 0) for r in slide_reports)
    total_images = sum(r.get("actual_images", 0) for r in slide_reports)
    summary_payload = {
        "slides_with_text_overflow": total_overflow,
        "slides_total": len(slide_reports),
        "total_images_rendered": total_images,
        "total_image_gap": total_image_gap,
        "slides_with_speaker_notes": sum(1 for r in slide_reports if r.get("has_speaker_notes")),
    }
    print(f"  Slides with text overflow: {total_overflow}/{len(slide_reports)}")
    print(f"  Total images rendered: {total_images}")
    print(f"  Total image gap (expected - actual): {total_image_gap}")
    print(f"  Slides with speaker notes: {sum(1 for r in slide_reports if r.get('has_speaker_notes'))}")

    if args.json_out:
        json_out_path = Path(args.json_out)
        payload = {
            "run_dir": str(run_dir),
            "pptx_path": str(pptx_path),
            "slides_in_pptx": len(prs.slides),
            "slides_in_deckir": len(deckir_slides),
            "summary": summary_payload,
            "slides": slide_reports,
        }
        json_out_path.parent.mkdir(parents=True, exist_ok=True)
        json_out_path.write_text(
            json.dumps(payload, sort_keys=True, ensure_ascii=True, indent=2) + "\n",
            encoding="utf-8",
        )


if __name__ == "__main__":
    main()
