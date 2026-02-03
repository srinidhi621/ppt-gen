#!/usr/bin/env python3
"""
Add Alt-Text (field_key) to PowerPoint template placeholders.

This script programmatically sets the 'descr' attribute (alt-text) on each
placeholder in the template, enabling deterministic placeholder binding
during rendering.

Usage:
    python add_alt_text.py                    # Process all layouts, save changes
    python add_alt_text.py --dry-run          # Preview changes without saving
    python add_alt_text.py --mvp-only         # Only process MVP priority layouts
    python add_alt_text.py --no-backup        # Skip backup creation
"""

import argparse
import json
import shutil
from collections import defaultdict
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn


# Placeholder type constants (from python-pptx enum values)
PH_TYPE_MAP = {
    "TITLE (1)": "title",
    "CENTER_TITLE (3)": "title",
    "SUBTITLE (4)": "subtitle",
    "BODY (2)": "body",
    "OBJECT (7)": "content",  # Generic content placeholder
    "PICTURE (18)": "image",
    "DATE (16)": "date",
    "FOOTER (15)": "footer",
    "SLIDE_NUMBER (13)": "slide_number",
}

# MVP priority layout name patterns (case-insensitive matching)
MVP_LAYOUT_PATTERNS = [
    "title with image",
    "title with half image",
    "section break",
    "one content",
    "two content",
    "three content",
    "four content",
    "header only",
    "statement",
    "agenda",
    "only title",
]


def get_placeholder_type_key(placeholder_type_str: str) -> str:
    """Map placeholder type string to a base key."""
    return PH_TYPE_MAP.get(placeholder_type_str, "unknown")


def get_field_key(base_type: str, position_index: int, total_same_type: int) -> str:
    """
    Generate a unique field_key based on placeholder type and position.
    
    Args:
        base_type: The base type key (title, body, content, image, etc.)
        position_index: 0-based index of this placeholder among same-type placeholders
        total_same_type: Total count of placeholders of this type in the layout
    
    Returns:
        A unique field_key string (e.g., 'ph_title', 'ph_col1', 'ph_image_left')
    """
    prefix = "ph_"
    
    # Single placeholder of this type - use simple name
    if total_same_type == 1:
        if base_type == "content":
            return f"{prefix}body"
        return f"{prefix}{base_type}"
    
    # Multiple placeholders of same type - need disambiguation
    if base_type in ("content", "body"):
        # Content/body placeholders use column naming
        if total_same_type == 2:
            return f"{prefix}body_left" if position_index == 0 else f"{prefix}body_right"
        else:
            # 3+ columns
            return f"{prefix}col{position_index + 1}"
    
    elif base_type == "image":
        # Image placeholders
        if total_same_type == 2:
            return f"{prefix}image_left" if position_index == 0 else f"{prefix}image_right"
        elif total_same_type == 3:
            suffixes = ["left", "center", "right"]
            return f"{prefix}image_{suffixes[position_index]}"
        else:
            # Many images (like photo grid) - use numbered suffix
            return f"{prefix}image_{position_index + 1}"
    
    else:
        # Other types with multiple instances - use numbered suffix
        if total_same_type == 2:
            return f"{prefix}{base_type}_1" if position_index == 0 else f"{prefix}{base_type}_2"
        return f"{prefix}{base_type}_{position_index + 1}"


def get_placeholder_info(shape) -> dict:
    """Extract placeholder information including position."""
    info = {
        "shape_id": shape.shape_id,
        "name": shape.name,
        "is_placeholder": shape.is_placeholder,
    }
    
    if shape.is_placeholder:
        try:
            placeholder_format = shape.placeholder_format
            info["placeholder_type"] = str(placeholder_format.type)
            info["placeholder_idx"] = placeholder_format.idx
        except Exception:
            info["placeholder_type"] = "Unknown"
    
    # Get position for sorting
    try:
        info["left"] = shape.left.inches if shape.left else 999
        info["top"] = shape.top.inches if shape.top else 999
    except Exception:
        info["left"] = 999
        info["top"] = 999
    
    return info


def set_alt_text(shape, field_key: str) -> bool:
    """
    Set the alt-text (descr attribute) on a shape's cNvPr element.
    
    Returns True if successful, False otherwise.
    """
    try:
        elem = shape.element
        nvSpPr = elem.find(qn('p:nvSpPr'))
        if nvSpPr is not None:
            cNvPr = nvSpPr.find(qn('p:cNvPr'))
            if cNvPr is not None:
                cNvPr.set('descr', field_key)
                return True
        
        # Try picture placeholder structure (nvPicPr)
        nvPicPr = elem.find(qn('p:nvPicPr'))
        if nvPicPr is not None:
            cNvPr = nvPicPr.find(qn('p:cNvPr'))
            if cNvPr is not None:
                cNvPr.set('descr', field_key)
                return True
                
    except Exception as e:
        print(f"    Warning: Could not set alt-text on {shape.name}: {e}")
        return False
    
    return False


def is_mvp_layout(layout_name: str) -> bool:
    """Check if a layout name matches MVP priority patterns."""
    name_lower = layout_name.lower()
    return any(pattern in name_lower for pattern in MVP_LAYOUT_PATTERNS)


def process_layout(layout, master_idx: int, layout_idx: int, dry_run: bool = False) -> dict:
    """
    Process a single layout, adding alt-text to all placeholders.
    
    Returns a report dict with the layout info and field_key assignments.
    """
    layout_report = {
        "master_index": master_idx,
        "layout_index": layout_idx,
        "layout_name": layout.name,
        "placeholders": [],
    }
    
    # Collect placeholder info
    placeholders = []
    for shape in layout.shapes:
        if shape.is_placeholder:
            info = get_placeholder_info(shape)
            info["shape"] = shape  # Keep reference for modification
            placeholders.append(info)
    
    if not placeholders:
        return layout_report
    
    # Group placeholders by type
    by_type = defaultdict(list)
    for ph in placeholders:
        ph_type = ph.get("placeholder_type", "Unknown")
        base_type = get_placeholder_type_key(ph_type)
        by_type[base_type].append(ph)
    
    # Process each type group
    for base_type, ph_group in by_type.items():
        # Sort by position (left first, then top)
        ph_group.sort(key=lambda x: (x["left"], x["top"]))
        
        total = len(ph_group)
        for idx, ph in enumerate(ph_group):
            field_key = get_field_key(base_type, idx, total)
            
            ph_report = {
                "shape_name": ph["name"],
                "placeholder_type": ph.get("placeholder_type", "Unknown"),
                "field_key": field_key,
                "position": f"L:{ph['left']:.2f}, T:{ph['top']:.2f}",
            }
            
            if not dry_run:
                success = set_alt_text(ph["shape"], field_key)
                ph_report["success"] = success
            else:
                ph_report["success"] = None  # Dry run
            
            layout_report["placeholders"].append(ph_report)
    
    return layout_report


def create_backup(template_path: Path) -> Path:
    """Create a timestamped backup of the template."""
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_name = f"template_backup_{timestamp}.pptx"
    backup_path = template_path.parent / backup_name
    shutil.copy2(template_path, backup_path)
    return backup_path


def main():
    parser = argparse.ArgumentParser(
        description="Add Alt-Text field_keys to PowerPoint template placeholders"
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Preview changes without modifying the template"
    )
    parser.add_argument(
        "--mvp-only",
        action="store_true",
        help="Only process MVP priority layouts"
    )
    parser.add_argument(
        "--no-backup",
        action="store_true",
        help="Skip creating a backup of the template"
    )
    parser.add_argument(
        "--template",
        type=str,
        default=None,
        help="Path to template file (default: assets/template/template.pptx)"
    )
    args = parser.parse_args()
    
    # Determine paths
    script_dir = Path(__file__).parent
    project_root = script_dir.parent
    
    if args.template:
        template_path = Path(args.template)
    else:
        template_path = project_root / "assets" / "template" / "template.pptx"
    
    if not template_path.exists():
        print(f"Error: Template not found at {template_path}")
        return 1
    
    print("=" * 70)
    print("ADD ALT-TEXT TO TEMPLATE PLACEHOLDERS")
    print("=" * 70)
    print(f"\nTemplate: {template_path}")
    print(f"Mode: {'DRY RUN (no changes)' if args.dry_run else 'LIVE (will modify template)'}")
    print(f"Scope: {'MVP layouts only' if args.mvp_only else 'All layouts'}")
    print()
    
    # Create backup (unless dry run or explicitly skipped)
    if not args.dry_run and not args.no_backup:
        backup_path = create_backup(template_path)
        print(f"Backup created: {backup_path}")
        print()
    
    # Load template
    prs = Presentation(str(template_path))
    
    # Process layouts
    report = {
        "template_path": str(template_path),
        "timestamp": datetime.now().isoformat(),
        "dry_run": args.dry_run,
        "mvp_only": args.mvp_only,
        "layouts_processed": [],
        "summary": {
            "total_layouts": 0,
            "layouts_with_placeholders": 0,
            "total_placeholders": 0,
            "successful_assignments": 0,
        }
    }
    
    print("-" * 70)
    print("PROCESSING LAYOUTS")
    print("-" * 70)
    
    for master_idx, master in enumerate(prs.slide_masters):
        for layout_idx, layout in enumerate(master.slide_layouts):
            # Filter by MVP if requested
            if args.mvp_only and not is_mvp_layout(layout.name):
                continue
            
            layout_report = process_layout(layout, master_idx, layout_idx, args.dry_run)
            report["layouts_processed"].append(layout_report)
            
            report["summary"]["total_layouts"] += 1
            
            if layout_report["placeholders"]:
                report["summary"]["layouts_with_placeholders"] += 1
                report["summary"]["total_placeholders"] += len(layout_report["placeholders"])
                
                # Count successful assignments
                for ph in layout_report["placeholders"]:
                    if ph.get("success") is True:
                        report["summary"]["successful_assignments"] += 1
                
                # Print progress
                print(f"\n[{master_idx}.{layout_idx}] {layout.name}")
                for ph in layout_report["placeholders"]:
                    status = ""
                    if not args.dry_run:
                        status = " [OK]" if ph.get("success") else " [FAIL]"
                    print(f"  {ph['shape_name']} → {ph['field_key']}{status}")
    
    # Save template if not dry run
    if not args.dry_run:
        prs.save(str(template_path))
        print(f"\n\nTemplate saved: {template_path}")
    
    # Save report
    report_path = script_dir / "alt_text_report.json"
    with open(report_path, 'w') as f:
        json.dump(report, f, indent=2)
    print(f"Report saved: {report_path}")
    
    # Print summary
    print("\n" + "=" * 70)
    print("SUMMARY")
    print("=" * 70)
    print(f"Layouts processed: {report['summary']['total_layouts']}")
    print(f"Layouts with placeholders: {report['summary']['layouts_with_placeholders']}")
    print(f"Total placeholders: {report['summary']['total_placeholders']}")
    if not args.dry_run:
        print(f"Successful assignments: {report['summary']['successful_assignments']}")
    
    return 0


if __name__ == "__main__":
    exit(main())
