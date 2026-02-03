#!/usr/bin/env python3
"""
Inspect a PowerPoint template to understand its structure.
Reports on: slide layouts, placeholders, alt-text, masters, etc.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import sys
from pathlib import Path


def get_placeholder_info(shape):
    """Extract placeholder information from a shape."""
    info = {
        "shape_id": shape.shape_id,
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "is_placeholder": shape.is_placeholder,
    }
    
    # Get alt text if available
    # Alt-text is stored in p:cNvPr (presentation namespace), not a:cNvPr (drawing namespace)
    try:
        if hasattr(shape, 'element'):
            element = shape.element
            # Try shape placeholder structure (p:nvSpPr/p:cNvPr)
            nvSpPr = element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr')
            if nvSpPr is not None:
                cNvPr = nvSpPr.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
                if cNvPr is not None:
                    descr = cNvPr.get('descr', '')
                    if descr:
                        info["alt_text"] = descr
            
            # Try picture placeholder structure (p:nvPicPr/p:cNvPr)
            if 'alt_text' not in info:
                nvPicPr = element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}nvPicPr')
                if nvPicPr is not None:
                    cNvPr = nvPicPr.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
                    if cNvPr is not None:
                        descr = cNvPr.get('descr', '')
                        if descr:
                            info["alt_text"] = descr
    except Exception as e:
        pass
        
    if shape.is_placeholder:
        try:
            placeholder_format = shape.placeholder_format
            info["placeholder_type"] = str(placeholder_format.type)
            info["placeholder_idx"] = placeholder_format.idx
        except Exception as e:
            info["placeholder_error"] = str(e)
    
    # Get position and size
    try:
        info["left_inches"] = round(shape.left.inches, 2) if shape.left else None
        info["top_inches"] = round(shape.top.inches, 2) if shape.top else None
        info["width_inches"] = round(shape.width.inches, 2) if shape.width else None
        info["height_inches"] = round(shape.height.inches, 2) if shape.height else None
    except Exception:
        pass
    
    return info


def inspect_template(pptx_path: str) -> dict:
    """Inspect a PowerPoint template and return a structured report."""
    prs = Presentation(pptx_path)
    
    report = {
        "file_path": str(pptx_path),
        "slide_width_inches": round(prs.slide_width.inches, 2),
        "slide_height_inches": round(prs.slide_height.inches, 2),
        "slide_masters": [],
        "slide_layouts": [],
        "existing_slides": [],
    }
    
    # Inspect slide masters
    for master_idx, master in enumerate(prs.slide_masters):
        master_info = {
            "master_index": master_idx,
            "name": master.name if hasattr(master, 'name') else f"Master {master_idx}",
            "layout_count": len(master.slide_layouts),
            "shapes": []
        }
        
        for shape in master.shapes:
            shape_info = get_placeholder_info(shape)
            master_info["shapes"].append(shape_info)
        
        report["slide_masters"].append(master_info)
    
    # Inspect slide layouts
    for master_idx, master in enumerate(prs.slide_masters):
        for layout_idx, layout in enumerate(master.slide_layouts):
            layout_info = {
                "master_index": master_idx,
                "layout_index": layout_idx,
                "name": layout.name,
                "placeholders": [],
                "other_shapes": []
            }
            
            for shape in layout.shapes:
                shape_info = get_placeholder_info(shape)
                if shape.is_placeholder:
                    layout_info["placeholders"].append(shape_info)
                else:
                    layout_info["other_shapes"].append(shape_info)
            
            # Also check placeholders collection
            layout_info["placeholder_count"] = len(layout.placeholders)
            
            report["slide_layouts"].append(layout_info)
    
    # Inspect existing slides (if any)
    for slide_idx, slide in enumerate(prs.slides):
        slide_info = {
            "slide_index": slide_idx,
            "layout_name": slide.slide_layout.name,
            "shapes": []
        }
        
        for shape in slide.shapes:
            shape_info = get_placeholder_info(shape)
            
            # Try to get text content
            if shape.has_text_frame:
                try:
                    text = shape.text_frame.text[:100]
                    if text:
                        shape_info["text_preview"] = text
                except:
                    pass
            
            slide_info["shapes"].append(shape_info)
        
        report["existing_slides"].append(slide_info)
    
    return report


def print_summary(report: dict):
    """Print a human-readable summary of the template."""
    print("=" * 80)
    print(f"TEMPLATE ANALYSIS: {Path(report['file_path']).name}")
    print("=" * 80)
    print(f"\nSlide dimensions: {report['slide_width_inches']}\" x {report['slide_height_inches']}\"")
    print(f"Number of slide masters: {len(report['slide_masters'])}")
    print(f"Total slide layouts: {len(report['slide_layouts'])}")
    print(f"Existing slides in file: {len(report['existing_slides'])}")
    
    print("\n" + "-" * 80)
    print("SLIDE LAYOUTS")
    print("-" * 80)
    
    for layout in report["slide_layouts"]:
        print(f"\n[{layout['layout_index']:2d}] {layout['name']}")
        print(f"     Placeholders: {len(layout['placeholders'])}")
        
        for ph in layout['placeholders']:
            alt_text = ph.get('alt_text', 'NO ALT TEXT')
            ph_type = ph.get('placeholder_type', 'Unknown')
            ph_idx = ph.get('placeholder_idx', '?')
            print(f"       - {ph['name']} | type: {ph_type} | idx: {ph_idx} | alt: {alt_text}")
    
    # Check for alt text coverage
    print("\n" + "-" * 80)
    print("ALT TEXT COVERAGE ANALYSIS")
    print("-" * 80)
    
    layouts_with_alt = 0
    layouts_without_alt = 0
    placeholders_with_alt = 0
    placeholders_without_alt = 0
    
    for layout in report["slide_layouts"]:
        has_any_alt = False
        for ph in layout['placeholders']:
            if ph.get('alt_text'):
                placeholders_with_alt += 1
                has_any_alt = True
            else:
                placeholders_without_alt += 1
        
        if has_any_alt:
            layouts_with_alt += 1
        else:
            layouts_without_alt += 1
    
    print(f"\nLayouts with at least one alt-text: {layouts_with_alt}")
    print(f"Layouts with NO alt-text: {layouts_without_alt}")
    print(f"\nPlaceholders with alt-text: {placeholders_with_alt}")
    print(f"Placeholders without alt-text: {placeholders_without_alt}")
    
    if placeholders_without_alt > 0:
        print("\n⚠️  WARNING: Some placeholders lack alt-text for field_key binding!")
    
    # Show existing slides summary
    if report["existing_slides"]:
        print("\n" + "-" * 80)
        print("EXISTING SLIDES IN TEMPLATE")
        print("-" * 80)
        for slide in report["existing_slides"]:
            print(f"\nSlide {slide['slide_index'] + 1}: Layout = {slide['layout_name']}")
            print(f"  Shapes: {len(slide['shapes'])}")


def main():
    if len(sys.argv) < 2:
        # Default paths to check
        default_paths = [
            "Assets/Ascendion_Presentation-Template_Apr25.pptx",
            "Assets/Corp Deck 2025 - Nov.pptx",
        ]
        
        base_path = Path(__file__).parent.parent
        
        for rel_path in default_paths:
            full_path = base_path / rel_path
            if full_path.exists():
                print(f"\nAnalyzing: {full_path}")
                report = inspect_template(str(full_path))
                print_summary(report)
                
                # Save JSON report
                report_path = base_path / "scripts" / f"template_report_{Path(rel_path).stem}.json"
                with open(report_path, 'w') as f:
                    json.dump(report, f, indent=2, default=str)
                print(f"\n📄 Full JSON report saved to: {report_path}")
    else:
        for path in sys.argv[1:]:
            if Path(path).exists():
                report = inspect_template(path)
                print_summary(report)
                
                # Save JSON report
                report_path = Path(path).parent / f"template_report_{Path(path).stem}.json"
                with open(report_path, 'w') as f:
                    json.dump(report, f, indent=2, default=str)
                print(f"\n📄 Full JSON report saved to: {report_path}")
            else:
                print(f"File not found: {path}")


if __name__ == "__main__":
    main()
