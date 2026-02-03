#!/usr/bin/env python3
"""
Generate layout_catalog.json from the PowerPoint template.

This script:
1. Reads the template and extracts layout/placeholder information
2. Identifies MVP-priority layouts
3. Generates stable layout_ids from layout names
4. Computes initial fit constraints from placeholder dimensions
5. Outputs assets/layout/layout_catalog.json
6. Validates the generated catalog against the template (drift detection)

Usage:
    python generate_layout_catalog.py                    # Generate catalog
    python generate_layout_catalog.py --validate-only    # Only validate existing catalog
    python generate_layout_catalog.py --all-layouts      # Include all layouts (not just MVP)
"""

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.oxml.ns import qn


# MVP Layout definitions: (layout_name_pattern, layout_id, master_index, layout_index)
# These are the 12 core layouts for MVP
MVP_LAYOUTS = [
    ("Title with Image 2", "title_image_light", 0, 0),
    ("Section Break - Light 1", "section_break_light", 0, 26),
    ("Header Only - Light", "header_only_light", 0, 8),
    ("One Content - Light", "one_content_light", 0, 10),
    ("Two Content - Light", "two_content_light", 0, 12),
    ("Three content - Light", "three_content_light", 0, 21),
    ("Four content - Light", "four_content_light", 0, 24),
    ("One Content With Image - Light", "content_image_light", 0, 15),
    ("Two Content with Image - Light", "two_content_image_light", 0, 19),
    ("Statement - Light", "statement_light", 0, 30),
    ("Agenda - Light", "agenda_light", 0, 6),
    ("BoilerPlate - Light", "boilerplate_light", 0, 41),
]

# Map MVP layout names to their layout_ids for quick lookup
MVP_LAYOUT_MAP = {name: (layout_id, mi, li) for name, layout_id, mi, li in MVP_LAYOUTS}

# Placeholder type mapping
PH_TYPE_MAP = {
    "TITLE (1)": "title",
    "CENTER_TITLE (3)": "title",
    "SUBTITLE (4)": "subtitle",
    "BODY (2)": "body",
    "OBJECT (7)": "content",
    "PICTURE (18)": "image",
    "DATE (16)": "date",
    "FOOTER (15)": "footer",
    "SLIDE_NUMBER (13)": "slide_number",
}


def get_field_key(shape) -> Optional[str]:
    """Read field_key from placeholder's alt-text (descr attribute)."""
    try:
        elem = shape.element
        # For shape placeholders: p:nvSpPr/p:cNvPr
        nvSpPr = elem.find(qn('p:nvSpPr'))
        if nvSpPr is not None:
            cNvPr = nvSpPr.find(qn('p:cNvPr'))
            if cNvPr is not None:
                descr = cNvPr.get('descr', '')
                if descr:
                    return descr
        # For picture placeholders: p:nvPicPr/p:cNvPr
        nvPicPr = elem.find(qn('p:nvPicPr'))
        if nvPicPr is not None:
            cNvPr = nvPicPr.find(qn('p:cNvPr'))
            if cNvPr is not None:
                descr = cNvPr.get('descr', '')
                if descr:
                    return descr
    except Exception:
        pass
    return None


def get_placeholder_type(shape) -> str:
    """Get the placeholder type as a string key."""
    if not shape.is_placeholder:
        return "unknown"
    try:
        ph_type_str = str(shape.placeholder_format.type)
        return PH_TYPE_MAP.get(ph_type_str, "unknown")
    except Exception:
        return "unknown"


def normalize_layout_id(layout_name: str) -> str:
    """Convert layout name to a stable snake_case layout_id."""
    # Remove special characters, convert to lowercase, replace spaces with underscores
    name = layout_name.lower()
    name = re.sub(r'[^a-z0-9\s]', '', name)
    name = re.sub(r'\s+', '_', name.strip())
    # Remove duplicate underscores
    name = re.sub(r'_+', '_', name)
    return name


def compute_constraints(placeholders: list, layout_name: str) -> dict:
    """
    Compute initial fit constraints based on placeholder dimensions.
    
    Heuristics:
    - avg_chars_per_line ≈ width_inches × 7 (for body text)
    - body_line_budget ≈ height_inches × 2.5
    - max_title_chars = 50-80 based on width
    - max_bullets = 4-8 based on height and column count
    """
    constraints = {
        "max_title_chars": 60,
        "max_bullets": 6,
        "max_words_per_bullet": 15,
        "max_total_body_chars": 600,
        "body_line_budget": 12,
        "avg_chars_per_line": 50,
    }
    
    # Find title and body placeholders
    title_ph = None
    body_phs = []
    
    for ph in placeholders:
        if ph.get("type") == "title":
            title_ph = ph
        elif ph.get("type") in ("body", "content"):
            body_phs.append(ph)
    
    # Adjust title constraints based on width
    if title_ph and title_ph.get("width_inches"):
        width = title_ph["width_inches"]
        # Roughly 8-10 chars per inch for titles
        constraints["max_title_chars"] = min(100, max(40, int(width * 9)))
    
    # Adjust body constraints based on dimensions and column count
    if body_phs:
        # Get the first (or only) body placeholder for sizing
        primary_body = body_phs[0]
        width = primary_body.get("width_inches", 6)
        height = primary_body.get("height_inches", 4)
        
        # Calculate chars per line (roughly 7 chars per inch for body text)
        constraints["avg_chars_per_line"] = max(30, int(width * 7))
        
        # Calculate line budget (roughly 2.5 lines per inch)
        constraints["body_line_budget"] = max(4, int(height * 2.5))
        
        # Adjust for multi-column layouts
        num_columns = len(body_phs)
        if num_columns >= 3:
            constraints["max_bullets"] = 4
            constraints["max_words_per_bullet"] = 10
            constraints["max_total_body_chars"] = 300
        elif num_columns == 2:
            constraints["max_bullets"] = 5
            constraints["max_words_per_bullet"] = 12
            constraints["max_total_body_chars"] = 400
        else:
            # Single column - more generous
            constraints["max_bullets"] = 7
            constraints["max_words_per_bullet"] = 18
            constraints["max_total_body_chars"] = 700
        
        # Update total chars based on line budget
        constraints["max_total_body_chars"] = min(
            constraints["max_total_body_chars"],
            constraints["body_line_budget"] * constraints["avg_chars_per_line"]
        )
    
    # Special cases based on layout type
    layout_lower = layout_name.lower()
    if "statement" in layout_lower:
        # Statement layouts: fewer, bigger text
        constraints["max_bullets"] = 1
        constraints["max_words_per_bullet"] = 30
        constraints["max_total_body_chars"] = 200
        constraints["body_line_budget"] = 4
    elif "agenda" in layout_lower:
        # Agenda layouts: more items, shorter text
        constraints["max_bullets"] = 10
        constraints["max_words_per_bullet"] = 8
    elif "boilerplate" in layout_lower:
        # Boilerplate: longer text block
        constraints["max_bullets"] = 1
        constraints["max_words_per_bullet"] = 100
        constraints["max_total_body_chars"] = 800
    elif "section" in layout_lower or "header only" in layout_lower:
        # Section breaks: title-focused
        constraints["max_bullets"] = 0
        constraints["max_total_body_chars"] = 0
    
    return constraints


def extract_layout_info(layout, master_idx: int, layout_idx: int) -> dict:
    """Extract all relevant information from a slide layout."""
    placeholders = []
    
    for shape in layout.shapes:
        if not shape.is_placeholder:
            continue
        
        field_key = get_field_key(shape)
        ph_type = get_placeholder_type(shape)
        
        # Skip placeholders without field_key (should not happen after add_alt_text.py)
        if not field_key:
            continue
        
        # Skip metadata placeholders for constraint computation
        if ph_type in ("date", "footer", "slide_number"):
            continue
        
        ph_info = {
            "field_key": field_key,
            "type": ph_type,
            "required": ph_type in ("title", "body", "content"),
        }
        
        # Get dimensions
        try:
            if shape.width:
                ph_info["width_inches"] = round(shape.width.inches, 2)
            if shape.height:
                ph_info["height_inches"] = round(shape.height.inches, 2)
            if shape.left:
                ph_info["left_inches"] = round(shape.left.inches, 2)
            if shape.top:
                ph_info["top_inches"] = round(shape.top.inches, 2)
        except Exception:
            pass
        
        placeholders.append(ph_info)
    
    # Sort placeholders by position (top-to-bottom, left-to-right)
    placeholders.sort(key=lambda p: (p.get("top_inches", 0), p.get("left_inches", 0)))
    
    # Build fields list (without dimension info for the catalog)
    fields = []
    for ph in placeholders:
        field_entry = {
            "field_key": ph["field_key"],
            "type": ph["type"],
            "required": ph["required"],
        }
        fields.append(field_entry)
    
    # Check if this is an MVP layout
    is_mvp = layout.name in MVP_LAYOUT_MAP
    layout_id = MVP_LAYOUT_MAP.get(layout.name, (None, None, None))[0]
    
    # If not MVP, generate layout_id from name
    if not layout_id:
        layout_id = normalize_layout_id(layout.name)
    
    # Compute constraints
    constraints = compute_constraints(placeholders, layout.name)
    
    return {
        "layout_id": layout_id,
        "template_layout_name": layout.name,
        "master_index": master_idx,
        "layout_index": layout_idx,
        "mvp": is_mvp,
        "fields": fields,
        "constraints": constraints,
    }


def generate_catalog(template_path: Path, include_all: bool = False) -> dict:
    """Generate the layout catalog from the template."""
    prs = Presentation(str(template_path))
    
    layouts = []
    seen_layout_ids = set()  # Track seen layout_ids to avoid duplicates
    
    for master_idx, master in enumerate(prs.slide_masters):
        for layout_idx, layout in enumerate(master.slide_layouts):
            layout_info = extract_layout_info(layout, master_idx, layout_idx)
            
            # Include only MVP layouts unless --all-layouts is specified
            if include_all or layout_info["mvp"]:
                # Skip duplicates (same layout_id from different masters)
                if layout_info["layout_id"] in seen_layout_ids:
                    continue
                seen_layout_ids.add(layout_info["layout_id"])
                layouts.append(layout_info)
    
    # Sort: MVP layouts first, then by master/layout index
    layouts.sort(key=lambda x: (not x["mvp"], x["master_index"], x["layout_index"]))
    
    catalog = {
        "version": "1.0",
        "template_path": "assets/template/template.pptx",
        "generated_from": template_path.name,
        "layouts": layouts,
    }
    
    return catalog


def validate_catalog(catalog: dict, template_path: Path) -> tuple[bool, list[str]]:
    """
    Validate catalog against template (template drift detection).
    
    Returns (is_valid, list of error messages).
    """
    errors = []
    prs = Presentation(str(template_path))
    
    # Build a map of template layouts
    template_layouts = {}
    for master_idx, master in enumerate(prs.slide_masters):
        for layout_idx, layout in enumerate(master.slide_layouts):
            key = (master_idx, layout_idx)
            # Extract field_keys
            field_keys = set()
            for shape in layout.shapes:
                if shape.is_placeholder:
                    fk = get_field_key(shape)
                    if fk:
                        field_keys.add(fk)
            template_layouts[key] = {
                "name": layout.name,
                "field_keys": field_keys,
            }
    
    # Validate each catalog entry
    for entry in catalog.get("layouts", []):
        layout_id = entry.get("layout_id")
        master_idx = entry.get("master_index")
        layout_idx = entry.get("layout_index")
        expected_name = entry.get("template_layout_name")
        
        key = (master_idx, layout_idx)
        
        # Check layout exists
        if key not in template_layouts:
            errors.append(
                f"[{layout_id}] Layout not found at master={master_idx}, layout={layout_idx}"
            )
            continue
        
        template_layout = template_layouts[key]
        
        # Check layout name matches
        if template_layout["name"] != expected_name:
            errors.append(
                f"[{layout_id}] Name mismatch: expected '{expected_name}', "
                f"found '{template_layout['name']}'"
            )
        
        # Check required field_keys exist
        for field in entry.get("fields", []):
            if field.get("required"):
                field_key = field.get("field_key")
                if field_key not in template_layout["field_keys"]:
                    errors.append(
                        f"[{layout_id}] Required field_key '{field_key}' not found in template"
                    )
    
    return len(errors) == 0, errors


def main():
    parser = argparse.ArgumentParser(
        description="Generate and validate layout_catalog.json from template"
    )
    parser.add_argument(
        "--validate-only",
        action="store_true",
        help="Only validate existing catalog without regenerating"
    )
    parser.add_argument(
        "--all-layouts",
        action="store_true",
        help="Include all layouts, not just MVP priority layouts"
    )
    parser.add_argument(
        "--template",
        type=str,
        default=None,
        help="Path to template file (default: assets/template/template.pptx)"
    )
    parser.add_argument(
        "--output",
        type=str,
        default=None,
        help="Output path for catalog (default: assets/layout/layout_catalog.json)"
    )
    args = parser.parse_args()
    
    # Determine paths
    script_dir = Path(__file__).parent
    project_root = script_dir.parent
    
    template_path = Path(args.template) if args.template else project_root / "assets" / "template" / "template.pptx"
    catalog_path = Path(args.output) if args.output else project_root / "assets" / "layout" / "layout_catalog.json"
    
    if not template_path.exists():
        print(f"Error: Template not found at {template_path}")
        return 1
    
    print("=" * 70)
    print("LAYOUT CATALOG GENERATOR")
    print("=" * 70)
    print(f"\nTemplate: {template_path}")
    print(f"Output: {catalog_path}")
    
    if args.validate_only:
        # Validate existing catalog
        if not catalog_path.exists():
            print(f"\nError: Catalog not found at {catalog_path}")
            return 1
        
        print("\nMode: VALIDATE ONLY")
        with open(catalog_path) as f:
            catalog = json.load(f)
        
        is_valid, errors = validate_catalog(catalog, template_path)
        
        if is_valid:
            print("\n[PASS] Catalog validation successful!")
            print(f"  Validated {len(catalog.get('layouts', []))} layouts")
            return 0
        else:
            print("\n[FAIL] Catalog validation failed!")
            for error in errors:
                print(f"  - {error}")
            return 1
    
    # Generate catalog
    print(f"\nMode: {'ALL LAYOUTS' if args.all_layouts else 'MVP LAYOUTS ONLY'}")
    
    catalog = generate_catalog(template_path, include_all=args.all_layouts)
    
    # Validate before saving
    print("\n" + "-" * 70)
    print("VALIDATING CATALOG")
    print("-" * 70)
    
    is_valid, errors = validate_catalog(catalog, template_path)
    
    if not is_valid:
        print("\n[WARN] Validation issues found:")
        for error in errors:
            print(f"  - {error}")
        print("\nCatalog will still be saved, but issues should be investigated.")
    else:
        print("\n[OK] Catalog validation passed")
    
    # Ensure output directory exists
    catalog_path.parent.mkdir(parents=True, exist_ok=True)
    
    # Save catalog
    with open(catalog_path, 'w') as f:
        json.dump(catalog, f, indent=2)
    
    print("\n" + "-" * 70)
    print("CATALOG GENERATED")
    print("-" * 70)
    print(f"\nSaved to: {catalog_path}")
    print(f"Total layouts: {len(catalog['layouts'])}")
    
    mvp_count = sum(1 for l in catalog['layouts'] if l.get('mvp'))
    print(f"MVP layouts: {mvp_count}")
    
    # Print layout summary
    print("\nLayouts included:")
    for layout in catalog['layouts']:
        mvp_tag = "[MVP]" if layout.get('mvp') else "     "
        field_count = len(layout.get('fields', []))
        print(f"  {mvp_tag} {layout['layout_id']}: {field_count} fields")
    
    return 0


if __name__ == "__main__":
    sys.exit(main())
